#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AIRUN PPTX Generator 전체 테스트
모든 레이아웃의 표, 이미지, 본문 기능 검증
"""

import os
import sys
import importlib
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# API 기본 설정 (환경변수로 오버라이드 가능)
DEFAULT_API_URL = "https://api.hamonize.com/api/v1/chat/sync"
DEFAULT_API_DOCS = [doc.strip() for doc in os.getenv("HAMONIZE_API_DOCUMENTS", "test_inputsheet.xlsx").split(",") if doc.strip()]
DEFAULT_API_USERNAME = os.getenv("HAMONIZE_API_USERNAME", "ryan")
DEFAULT_API_KEY = (
    os.getenv("HAMONIZE_API_KEY")
    or os.getenv("hamonize_api_key")
    or "airun_3_690a7671d2e0b3fa81211f7fae3a8356"
)


def _extract_text_from_api_response(data: Any) -> Optional[str]:
    """API 응답(JSON)에서 텍스트 본문을 안전하게 추출."""

    if data is None:
        return None

    if isinstance(data, str):
        return data.strip()

    if isinstance(data, dict):
        # 흔히 쓰일 수 있는 키 우선순위
        candidates = [
            data.get("answer"),
            data.get("content"),
            data.get("message"),
            data.get("text"),
            data.get("data"),
        ]
        for candidate in candidates:
            if isinstance(candidate, str) and candidate.strip():
                return candidate.strip()
        # dict 전체가 중요할 수 있으므로 문자열로 직렬화
        return str(data)

    # 리스트 등 기타 타입은 문자열 변환
    return str(data)


def call_slide_api(prompt: str,
                   *,
                   mentioned_documents: Optional[List[str]] = None,
                   rag: bool = False,
                   web: bool = False,
                   username: str = DEFAULT_API_USERNAME,
                   timeout: int = 60,
                   retries: int = 2,
                   backoff: int = 2) -> Optional[str]:
    """슬라이드별 프롬프트로 API를 호출하여 텍스트를 반환.

    HAMONIZE_API_KEY가 없거나 requests 미설치 시 None을 반환하고 로그만 남긴다.
    """

    api_key = DEFAULT_API_KEY
    if not api_key:
        print("⚠️ HAMONIZE_API_KEY 환경변수가 없어 API 호출을 건너뜁니다.")
        return None
# airun_3_690a7671d2e0b3fa81211f7fae3a8356
    try:
        import requests  # type: ignore
    except ImportError:
        print("⚠️ requests 라이브러리가 없어 API 호출을 건너뜁니다.")
        return None

    api_url = os.getenv("HAMONIZE_API_URL", DEFAULT_API_URL)
    docs = mentioned_documents if mentioned_documents is not None else DEFAULT_API_DOCS

    payload: Dict[str, Any] = {
        "prompt": prompt,
        "options": {
            "mentionedDocuments": docs,
            "rag": rag,
            "web": web,
            "username": username,
        },
    }

    headers = {
        "accept": "application/json",
        "X-API-Key": api_key,
        "Content-Type": "application/json",
    }

    for attempt in range(1, retries + 2):
        try:
            resp = requests.post(api_url, headers=headers, json=payload, timeout=timeout)
            resp.raise_for_status()
            data = resp.json()
            text = _extract_text_from_api_response(data)
            if text:
                return text
            print(f"⚠️ API 응답에서 텍스트를 찾지 못했습니다. 응답: {data}")
            return None
        except Exception as exc:  # pylint: disable=broad-except
            if attempt > retries:
                print(f"⚠️ API 호출 실패 ({prompt[:40]}...): {exc}")
                return None
            sleep_for = backoff * attempt
            print(f"⚠️ API 호출 실패, {sleep_for}s 후 재시도 {attempt}/{retries} ({prompt[:40]}...): {exc}")
            try:
                import time
                time.sleep(sleep_for)
            except Exception:
                pass


def get_api_text_or_default(prompt: str,
                            default: str,
                            *,
                            documents: Optional[List[str]] = None,
                            rag: bool = False,
                            web: bool = False,
                            username: str = DEFAULT_API_USERNAME) -> str:
    """API 호출 결과를 반환하되 실패 시 기본값을 사용."""

    text = call_slide_api(prompt, mentioned_documents=documents, rag=rag, web=web, username=username)
    if text:
        return text
    return default

# 현재 디렉토리를 sys.path 최우선으로 설정 (pv_solar 로컬 solar_pptx 사용)
current_dir = Path(__file__).resolve().parent
sys.path.insert(0, str(current_dir))
# 프로젝트 루트가 sys.path에 이미 있으면 제거해 로컬 모듈 충돌을 방지
project_root = str(current_dir.parent)
sys.path = [p for p in sys.path if p != project_root]
# 이전에 로드된 모듈이 있으면 제거 후 재로딩
sys.modules.pop('airun_pptx', None)
sys.modules.pop('solar_pptx', None)

from solar_pptx import PPTXGenerator, SlideContent, create_pptx_document


def render_pdf_cover(pdf_path: Path, output_dir: Path, dpi: int = 220) -> Path | None:
    """Render the first page of a PDF to PNG for use as a slide cover.

    Returns the path to the generated image, or None if rendering fails.
    """

    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        print(f"⚠️ PDF 파일을 찾을 수 없습니다: {pdf_path}")
        return None

    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"{pdf_path.stem}_page1.png"

    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("⚠️ PyMuPDF(fitz)가 설치되지 않아 PDF 표지를 사용할 수 없습니다. 기존 표지를 유지합니다.")
        return None

    try:
        doc = fitz.open(pdf_path)
        page = doc.load_page(0)
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        pix.save(output_path)
        print(f"✅ PDF 표지 이미지를 생성했습니다: {output_path}")
        return output_path
    except Exception as exc:  # pylint: disable=broad-except
        print(f"⚠️ PDF 표지 변환 중 오류가 발생했습니다: {exc}")
        return None


def apply_cover_image(pptx_path: Path, cover_image: Path) -> None:
    """Replace slide 1 with the provided cover image."""

    pptx_path = Path(pptx_path)
    cover_image = Path(cover_image)
    if not cover_image.exists():
        print(f"⚠️ 표지 이미지가 없어 교체를 건너뜁니다: {cover_image}")
        return

    prs = Presentation(str(pptx_path))
    if not prs.slides:
        print("⚠️ 슬라이드가 없어 표지를 교체하지 않습니다.")
        return

    slide = prs.slides[0]
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)  # type: ignore[attr-defined]

    zero = Emu(0)
    slide.shapes.add_picture(str(cover_image), zero, zero, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(pptx_path))
    print(f"✅ 1번 슬라이드를 PDF 표지로 교체했습니다: {cover_image}")


def prune_slides(pptx_path: Path, keep: int) -> None:
    """생성 후 슬라이드 개수를 기대치에 맞춰 정리한다."""

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if len(prs.slides) <= keep:
        return

    sld_id_lst = prs.slides._sldIdLst
    while len(sld_id_lst) > keep:
        sld_id_lst.remove(sld_id_lst[-1])

    prs.save(str(pptx_path))
    print(f"⚙️  불필요 슬라이드 제거: 총 {len(prs.slides)}장으로 정리")


def apply_session_slide(pptx_path: Path, slide_idx: int, number: str, title_text: str) -> None:
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if slide_idx >= len(prs.slides):
        return

    slide = prs.slides[slide_idx]
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)

    # 슬라이드 비우기
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)

    # 배경 적용 (bleed 2%)
    bg_path = Path(__file__).parent / "session_title_bg2.png"
    if not bg_path.exists():
        print(f"⚠️ 세션 배경 이미지를 찾을 수 없습니다: {bg_path}")
        return
    bleed_ratio = 1.02
    adj_w_val = int(slide_width * bleed_ratio)
    adj_h_val = int(slide_height * bleed_ratio)
    adj_w = Emu(adj_w_val)
    adj_h = Emu(adj_h_val)
    offset_x = Emu(int(-(adj_w_val - slide_width) / 2))
    offset_y = Emu(int(-(adj_h_val - slide_height) / 2))
    slide.shapes.add_picture(str(bg_path), offset_x, offset_y, width=adj_w, height=adj_h)

    # 텍스트 박스
    slide_w_inch = slide_width / 914400
    slide_h_inch = slide_height / 914400
    box_left = Inches(1)
    box_width = Inches(slide_w_inch * 0.45)
    box_height = Inches(2.0)
    box_top = Inches(slide_h_inch * 0.45)

    textbox = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{number}\n{title_text}"
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.1
    run = p.runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    prs.save(str(pptx_path))



def build_content_list_from_api(documents: Optional[List[str]] = None) -> List[SlideContent]:
    """슬라이드마다 API를 호출해 내용을 채운 뒤 SlideContent 리스트를 생성."""

    docs = documents if documents else DEFAULT_API_DOCS

    cover_subtitle = get_api_text_or_default(
        "레그에 등록된 test_inputsheet.xlsx 데이터를 기반으로 이 문서를 한 줄로 요약해 표지 부제로만 반환하세요.",
        "Project overview subtitle",
        documents=docs,
        rag=True,
        username=DEFAULT_API_USERNAME,
    )

    exec_summary_text = get_api_text_or_default(
        "레그에 등록된 test_inputsheet.xlsx 데이터를 기반으로, 샘플 PDF(/home/no/ryan/11.airun/solar/workspaces_pv_solar/pv_solar_sample_docs/[Sample] CAPEX Approval PPT.pdf) 2페이지 톤과 구조처럼 Executive Summary를 'BOLT #2 - Phase 1'과 'BOLT #2 - Phase 2'로 구분해 작성하세요. 각 구간마다 불릿(•) 3~4개를 한 줄 요약으로 작성하고, 가능하면 아래 항목을 Phase별로 포함하세요: 용량/프로젝트 수/지역, RtB·EBL·DAP 등 인허가 상태, COD 목표 시점, 계통연계/REC 판매(단가) 전략, CAPEX(통화/환율), 수익성 지표(IRR 등), 자본이익 추정, 실행 주체. 데이터가 부족하면 '데이터 확인 필요'라고 명시하세요. 포맷 예시:\nBOLT #2 - Phase 1\n• 포인트1\n• 포인트2\nBOLT #2 - Phase 2\n• 포인트1\n• 포인트2",
        "BOLT #2 - Phase 1\n• Phase 1 주요 포인트 1\n• Phase 1 주요 포인트 2\n• Phase 1 주요 포인트 3\nBOLT #2 - Phase 2\n• Phase 2 주요 포인트 1\n• Phase 2 주요 포인트 2\n• Phase 2 주요 포인트 3",
        documents=docs,
        rag=True,
        username=DEFAULT_API_USERNAME,
    )

    session_text = get_api_text_or_default(
        "이 문서의 섹션 표지에 들어갈 한 문장 설명을 작성하세요. 레그 test_inputsheet.xlsx 내용에서 전체 프로젝트의 핵심 목적을 1문장으로 정리해 주세요.",
        "프로젝트 핵심 목적 요약",
        documents=docs,
        rag=True,
        username=DEFAULT_API_USERNAME,
    )

    content_list: List[SlideContent] = []

    # 슬라이드 1: 표지
    content_list.append(SlideContent(
        layout="title",
        title="Capex Approval – S. Korea DG BOLT#2",
        subtitle=cover_subtitle,
        date="November 15, 2023",
        order=1,
    ))

    # 슬라이드 2: Executive Summary (7개 핵심 포인트)
    content_list.append(SlideContent(
        layout="executive_summary",
        title="Executive Summary",
        content=exec_summary_text,
        order=2,
    ))

    # 슬라이드 3: 섹션 표지
    content_list.append(SlideContent(
        layout="session_title",
        title="01",
        content=session_text,
        order=3,
    ))


    # # 슬라이드 2: Executive Summary (표 2개 + 요약 텍스트)
    # content_list.append(SlideContent(
    #     title="Executive Summary",
    #     subtitle=table_caption,
    #     content=table_caption,
    #     table_headers=["#", "BOLT #2 - Phase 1"],
    #     table_data=[
    #         ["1", "240.69"],
    #         ["2", "386.24"],
    #         ["3", "210.00"],
    #         ["4", "171.00"],
    #     ],
    #     table_headers_2=["#", "BOLT #2 - Phase 2"],
    #     table_data_2=[
    #         ["1", "100"],
    #         ["2", "80"],
    #         ["3", "20"],
    #     ],
    #     layout='two_tables',
    #     order=2,
    # ))

    # # 슬라이드 3: Approval request
    # content_list.append(SlideContent(
    #     layout="approval_request",
    #     title="Approval request",
    #     content=approval_text,
    #     order=3,
    # ))

    
    return content_list


def test_airun_pptx():
    """AIRUN PPTX Generator 전체 기능 테스트"""
    print("=" * 80)
    print("AIRUN PPTX Generator 전체 테스트")
    print("=" * 80)

    # 출력 디렉토리 생성
    output_dir = Path("/home/no/문서/generated-pptx/custom")
    output_dir.mkdir(parents=True, exist_ok=True)

    # v2 방식: content_list에 모든 슬라이드 내용을 담기 (API 연동)
    content_list = build_content_list_from_api()
    
    
    
    # # ========== 슬라이드 2: 레이아웃 0 - basic (제목 + 본문) ==========
    # content_list.append(SlideContent(
    #     title="레이아웃 0: Basic (제목 + 본문)",
    #     content="• 이것은 기본 레이아웃입니다\n• 제목과 본문 텍스트로 구성\n• 가장 단순한 형태의 슬라이드\n\n"
    #            "주요 특징:\n"
    #            "- 깔끔한 제목 영역\n"
    #            "- 넓은 본문 영역\n"
    #            "- 일반적인 텍스트 콘텐츠에 적합",
    #     layout='content_title_text'
    # ))

    # # ========== 슬라이드 3: 레이아웃 1 - two_column (제목 + 2단 텍스트) ==========
    # content_list.append(SlideContent(
    #     title="레이아웃 1: Two Column (제목 + 2단 텍스트)",
    #     content="• 왼쪽 단락: 이 레이아웃은 화면을 좌우로 나누어\n  두 가지 내용을 병렬로 보여줍니다\n\n"
    #            "• 텍스트가 자동으로 두 단으로 나뉩니다\n"
    #            "  각 단은 균형있는 크기를 가집니다\n\n"
    #            "• 내용이 많은 경우 유용합니다",
    #     layout='two_column'
    # ))

    # # ========== 슬라이드 4: 레이아웃 2 - three_column (제목 + 3단 텍스트) ==========
    # content_list.append(SlideContent(
    #     title="레이아웃 2: Three Column (제목 + 3단 텍스트)",
    #     content="• 첫 번째: 3단 레이아웃입니다\n"
    #            "• 두 번째: 화면을 3등분하여\n"
    #            "• 세 번째: 세 가지 내용을 비교합니다",
    #     layout='three_column'
    # ))

    # # ========== 슬라이드 5: 레이아웃 3 - large_table (제목 + 대형 테이블) ==========
    # content_list.append(SlideContent(
    #     title="레이아웃 3: Large Table (제목 + 대형 테이블)",
    #     table_headers=["항목", "설명", "값", "비고"],
    #     table_data=[
    #         ["항목 1", "첫 번째 데이터", "100", "메모"],
    #         ["항목 2", "두 번째 데이터", "200", "메모"],
    #         ["항목 3", "세 번째 데이터", "300", "메모"],
    #         ["항목 4", "네 번째 데이터", "400", "메모"],
    #         ["항목 5", "다섯 번째 데이터", "500", "메모"],
    #         ["합계", "총계", "1500", "완료"],
    #     ],
    #     layout='large_table'
    # ))


    # # ========== 슬라이드 7: 레이아웃 5 - gantt_chart (제목 + 간트 차트 + 하단 텍스트) ==========
    # content_list.append(SlideContent(
    #     title="레이아웃 5: Gantt Chart (제목 + 간트 차트 + 하단 텍스트)",
    #     content="• 1월: 기획 및 설계\n"
    #            "• 2월: 개발 시작\n"
    #            "• 3월: 중간 점검\n"
    #            "• 4월: 테스트 및 수정\n"
    #            "• 5월: 최종 배포\n\n"
    #            "하단 설명: 프로젝트는 5개월 동안 진행되며, 각 단계별로 마일스톤이 있습니다.",
    #     layout='gantt_chart'
    # ))

    # # ========== 슬라이드 8: 레이아웃 6 - text_small_table (제목 + 본문 + 작은 테이블) ==========
    # content_list.append(SlideContent(
    #     title="레이아웃 6: Text + Small Table (제목 + 본문 + 작은 테이블)",
    #     content="• 이 레이아웃은 상단에 본문 텍스트가 있고\n"
    #            "• 하단에 작은 표가 위치합니다\n"
    #            "• 텍스트로 내용을 설명하고 표로 데이터를 보충할 때 유용합니다\n\n"
    #            "주요 활용처:\n"
    #            "- 요약과 상세 데이터\n"
    #            "- 설명과 수치 비교",
    #     table_headers=["구분", "Q1", "Q2", "Q3", "Q4"],
    #     table_data=[
    #         ["매출", "100", "120", "140", "160"],
    #         ["비용", "80", "90", "95", "100"],
    #         ["이익", "20", "30", "45", "60"],
    #     ],
    #     layout='text_small_table'
    # ))

    # # ========== 슬라이드 9: 레이아웃 7 - text_large_table (제목 + 상단 텍스트 + 대형 테이블) ==========
    # content_list.append(SlideContent(
    #     title="레이아웃 7: Text + Large Table (제목 + 상단 텍스트 + 대형 테이블)",
    #     content="• 상단에 짧은 설명 텍스트가 있습니다\n• 하단에 대형 테이블이 위치합니다",
    #     table_headers=["프로젝트", "위치", "용량(kWp)", "REC 가격", "상태"],
    #     table_data=[
    #         ["프로젝트 A", "부산", "240.69", "KRW 189/kWh", "Ready"],
    #         ["프로젝트 B", "충북", "386.24", "KRW 189/kWh", "Ready"],
    #         ["프로젝트 C", "경기", "210.00", "KRW 189/kWh", "Ready"],
    #         ["프로젝트 D", "인천", "171.00", "KRW 189/kWh", "Ready"],
    #         ["프로젝트 E", "경남", "500.00", "KRW 189/kWh", "Ready"],
    #         ["합계", "", "1,507.93", "", "완료"],
    #     ],
    #     layout='text_large_table'
    # ))

    # # ========== 슬라이드 10: 추가 테스트 - large_table 더 많은 데이터 ==========
    # content_list.append(SlideContent(
    #     title="추가 테스트: 대형 테이블 (많은 데이터)",
    #     table_headers=["No", "프로젝트명", "지역", "용량(kWp)", "모듈", "인버터", "구조", "예산(M KRW)"],
    #     table_data=[
    #         ["1", "프로젝트 A", "부산", "240.69", "710W", "100kW", "강구조", "265"],
    #         ["2", "프로젝트 B", "충북", "386.24", "710W", "100kW", "강구조", "425"],
    #         ["3", "프로젝트 C", "경기", "210.00", "710W", "100kW", "강구조", "231"],
    #         ["4", "프로젝트 D", "인천", "171.00", "710W", "100kW", "강구조", "188"],
    #         ["5", "프로젝트 E", "경남", "500.00", "710W", "100kW", "강구조", "550"],
    #         ["6", "프로젝트 F", "경기", "122.00", "710W", "100kW", "강구조", "134"],
    #         ["7", "프로젝트 G", "부산", "234.00", "710W", "100kW", "강구조", "257"],
    #         ["8", "프로젝트 H", "경기", "438.00", "710W", "100kW", "강구조", "482"],
    #     ],
    #     layout='large_table'
    # ))

    # # ========== 슬라이드 11: 레이아웃 9 - image_left (제목 + 본문 + 이미지) ==========
    # # 실제 이미지 사용: test_smart_tune.png (작은 이미지, 67KB)
    # sample_image1 = "/home/hamonikr/work/airun/assets/images/sample/test_smart_tune.png"
    # content_list.append(SlideContent(
    #     title="레이아웃 9: Image Left (제목 + 본문 + 이미지)",
    #     content="• 이 레이아웃은 왼쪽에 본문 텍스트가 있고\n• 오른쪽에 이미지가 배치됩니다\n• 본문과 이미지를 함께 보여줄 때 유용합니다\n\n"
    #            "테스트 이미지: test_smart_tune.png (67KB)\n"
    #            "원본 비율 유지하며 크기 조정",
    #     image_path=sample_image1,
    #     layout='image_left'
    # ))

    # # ========== 슬라이드 12: 레이아웃 10 - image_right (제목 + 본문 왼쪽 + 이미지 오른쪽) ==========
    # # 실제 이미지 사용: test_rag_page.png (391KB)
    # sample_image2 = "/home/hamonikr/work/airun/assets/images/sample/test_rag_page.png"
    # content_list.append(SlideContent(
    #     title="레이아웃 10: Image Right (제목 + 본문 왼쪽 + 이미지 오른쪽)",
    #     content="• 이 레이아웃은 왼쪽에 본문 텍스트가 있고\n• 오른쪽에 이미지가 배치됩니다\n• image_left와 반대 배치입니다\n\n"
    #            "테스트 이미지: test_rag_page.png (391KB)\n"
    #            "원본 비율 유지하며 크기 조정",
    #     image_path=sample_image2,
    #     layout='image_right'
    # ))

    # # ========== 슬라이드 13: 레이아웃 11 - large_image (제목 + 큰 이미지) ==========
    # # 실제 이미지 사용: test_infographic_large.png (6.4MB - 큰 인포그래픽)
    # sample_image3 = "/home/hamonikr/work/airun/assets/images/sample/test_infographic_large.png"
    # content_list.append(SlideContent(
    #     title="레이아웃 11: Large Image (제목 + 큰 이미지 영역)",
    #     content="큰 이미지 영역\n\n"
    #            "테스트 이미지: test_infographic_large.png (6.4MB)\n"
    #            "원본 비율 유지하며 전체 영역에 채움",
    #     image_path=sample_image3,
    #     layout='large_image'
    # ))

    # # ========== 슬라이드 14: 레이아웃 12 - two_images (제목 + 이미지 2개) ==========
    # # 실제 이미지 사용: test_admin_rag.png & test_background_blue.jpg
    # sample_image4 = "/home/hamonikr/work/airun/assets/images/sample/test_admin_rag.png"
    # sample_image5 = "/home/hamonikr/work/airun/assets/images/sample/test_background_blue.jpg"
    # content_list.append(SlideContent(
    #     title="레이아웃 12: Two Images (제목 + 이미지 2개 좌우 배치)",
    #     content="좌우 이미지 비교\n\n"
    #            "왼쪽: test_admin_rag.png (182KB)\n"
    #            "오른쪽: test_background_blue.jpg (15KB)\n\n"
    #            "각 이미지는 원본 비율 유지",
    #     image_path=sample_image4,
    #     image_path_2=sample_image5,
    #     layout='two_images'
    # ))

    # # ========== 슬라이드 15: 레이아웃 8 - two_tables_vertical (제목 + 표 2개 상하) ==========
    # content_list.append(SlideContent(
    #     title="레이아웃 8: Two Tables Vertical (제목 + 표 2개 상하 분할)",
    #     content="• 상단 표와 하단 표가 상하로 배치됩니다\n• 각 표는 독립적인 데이터를 보여줍니다",
    #     table_headers=["카테고리", "항목", "값"],
    #     table_data=[
    #         ["매출", "Q1", "100"],
    #         ["매출", "Q2", "120"],
    #         ["비용", "Q1", "80"],
    #         ["비용", "Q2", "90"],
    #     ],
    #     table_headers_2=["지역", "프로젝트", "용량"],
    #     table_data_2=[
    #         ["부산", "A", "100"],
    #         ["서울", "B", "200"],
    #         ["경기", "C", "150"],
    #     ],
    #     layout='two_tables_vertical'
    # ))

    # ========== 문서 생성 ==========
    output_path = output_dir / "test_airun_pptx_16_9.pptx"

    # PPTXGenerator 생성
    gen = PPTXGenerator(aspect_ratio='16_9', template_type='custom')
    gen.generate_with_template(content_list, str(output_path))

    # 세션 타이틀 슬라이드 재구성 (템플릿 인덱스 불일치 대비)
    # for idx, c in enumerate(content_list):
    #     if c.layout == "session_title":
    #         apply_session_slide(output_path, idx, c.title, c.content or "")

    # 템플릿 기본 슬라이드가 남지 않도록 생성된 슬라이드를 content_list 길이에 맞춰 정리
    prune_slides(output_path, keep=len(content_list))

    print("=" * 80)
    print(f"\n✅ 문서 생성 완료: {output_path}")
    print(f"   총 슬라이드: {len(content_list)}")

    # 로고 확인
    from pptx import Presentation
    prs = Presentation(str(output_path))
    logo_count = 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                left_inch = float(shape.left) / 914400
                if left_inch > 10:  # 오른쪽 상단 로고
                    logo_count += 1

    print(f"   로고 포함: {logo_count}개 슬라이드")

    # 레이아웃 요약
    print("\n" + "=" * 80)
    print("테스트된 레이아웃 요약")
    print("=" * 80)
    print("슬라이드 1: 표지 (title)")
    print("슬라이드 2: Executive Summary (two_tables)")
    print("슬라이드 3: Approval request")
    print("슬라이드 4: 세션 표지 (session_title)")
    print("슬라이드 5: Executive Summary 상세 (executive_summary)")
    print("=" * 80)

    return output_path


if __name__ == "__main__":
    test_airun_pptx()
