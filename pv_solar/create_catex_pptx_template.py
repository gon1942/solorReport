#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Custom PPTX 템플릿 생성 - CAPEX/DEVEX Approval 문서용
- 통합 템플릿 방식 (제목 + 내용이 하나의 파일에)
- 대형 테이블, 2단 테이블, 3단 레이아웃 등 지원
"""

import sys
import os
import subprocess
from pathlib import Path

# 현재 디렉토리 우선 검색 (pv_solar 내 로컬 solar_pptx 사용)
sys.path.insert(0, str(Path(__file__).parent))
# 프로젝트 루트 경로 설정 (백업 경로)
sys.path.insert(1, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from utils import getVarVal
import cairosvg
from io import BytesIO
import os
from PIL import Image

# 프로젝트 루트 경로 (설정 파일에서 가져옴)
AIRUN_PATH = Path(getVarVal('AIRUN_PATH', str(Path(__file__).parent.parent)))
# 출력 디렉토리 (custom 폴더)
OUTPUT_DIR = AIRUN_PATH / "templates" / "pptx" / "custom"

# 페이지 크기
SLIDE_WIDTH_4_3 = 10  # 인치
SLIDE_WIDTH_16_9 = 13.333  # 인치
SLIDE_HEIGHT = 7.5  # 인치

# 페이지 여백
MARGIN = 0.5
CONTENT_TOP = 1.3  # 제목 아래부터 시작

# 로고 경로
LOGO_SVG_PATH = AIRUN_PATH / "assets" / "images" / "sample" / "test_logo.svg"

# 표지 배경 이미지 (직접 제공된 PNG)
COVER_IMAGE_PATH = Path(__file__).parent / "session_title_bg.png"

# 표지 PDF 경로(대체용) 및 캐시
PDF_COVER_PATH = Path("/work/proj/airun/002_poc/airun/workspaces_pv_solar/sample_docs/[Sample] CAPEX Approval PPT.pdf")
PDF_COVER_DPI = 220
PDF_COVER_OUTPUT = OUTPUT_DIR / "_cache"

# 전체 샘플 PDF(슬라이드 전체 배경 반영용)
_SAMPLE_CAPEX_PDF_CANDIDATES = [
    AIRUN_PATH / "pv_solar_sample_docs" / "[Sample] CAPEX Approval PPT.pdf",
    Path(__file__).resolve().parents[1] / "pv_solar_sample_docs" / "[Sample] CAPEX Approval PPT.pdf",
]
SAMPLE_CAPEX_PDF_PATH = next((p for p in _SAMPLE_CAPEX_PDF_CANDIDATES if p.exists()), _SAMPLE_CAPEX_PDF_CANDIDATES[0])
SAMPLE_PDF_DPI = 170
SAMPLE_PDF_OUTPUT = OUTPUT_DIR / "_cache" / "sample_pages"

# 세션 타이틀 배경 이미지
SESSION_BG_PATH = Path(__file__).parent / "session_title_bg2.png"


def set_slide_full_background(slide, image_path: Path, width=None, height=None) -> bool:
    """슬라이드 전체를 배경 이미지로 채움 (비율 유지 없이 스트레치)."""

    image_path = Path(image_path)
    if not image_path.exists():
        print(f"  ⚠️  표지 이미지 없음: {image_path}")
        return False

    if width is None:
        width = getattr(slide, "slide_width", None) or getattr(getattr(slide, "_parent", None), "slide_width", None)
    if height is None:
        height = getattr(slide, "slide_height", None) or getattr(getattr(slide, "_parent", None), "slide_height", None)

    if not width or not height:
        print("  ⚠️  슬라이드 크기를 가져올 수 없어 배경 적용을 건너뜁니다.")
        return False

    # 약간의 블리드(여백 메우기) 적용
    bleed_ratio = 1.02
    adj_width = int(width * bleed_ratio)
    adj_height = int(height * bleed_ratio)
    offset_x = int(-(adj_width - width) / 2)
    offset_y = int(-(adj_height - height) / 2)

    zero = Emu(0)
    try:
        slide.shapes.add_picture(
            str(image_path),
            Emu(offset_x),
            Emu(offset_y),
            width=Emu(adj_width),
            height=Emu(adj_height)
        )
        return True
    except Exception as exc:  # pylint: disable=broad-except
        print(f"  ⚠️  표지 이미지 적용 실패: {exc}")
        return False


def render_pdf_cover(pdf_path: Path, output_dir: Path, dpi: int = 220) -> Path | None:
    """PDF 1페이지를 PNG로 변환해 반환 (실패 시 None)."""

    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        print(f"  ⚠️  표지 PDF를 찾을 수 없습니다: {pdf_path}")
        return None

    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"{pdf_path.stem}_page1.png"

    if output_path.exists():
        return output_path

    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("  ⚠️  PyMuPDF(fitz)가 설치되지 않아 PDF 표지를 사용할 수 없습니다.")
        return None

    try:
        doc = fitz.open(pdf_path)
        page = doc.load_page(0)
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        pix.save(output_path)
        print(f"  ✅ PDF 표지를 이미지로 변환했습니다: {output_path}")
        return output_path
    except Exception as exc:  # pylint: disable=broad-except
        print(f"  ⚠️  PDF 표지 변환 실패: {exc}")
        return None


def render_pdf_pages(pdf_path: Path, output_dir: Path, dpi: int = 170) -> list[Path]:
    """pdftoppm으로 PDF 전체 페이지를 PNG로 렌더링한다."""
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        print(f"  ⚠️  샘플 PDF를 찾을 수 없습니다: {pdf_path}")
        return []

    output_dir.mkdir(parents=True, exist_ok=True)
    prefix = output_dir / pdf_path.stem

    # 기존 파일 정리
    prefix_name = f"{pdf_path.stem}-"

    for old in output_dir.iterdir():
        if not old.is_file() or old.suffix.lower() != ".png":
            continue
        if not old.name.startswith(prefix_name):
            continue
        try:
            old.unlink()
        except OSError:
            pass

    cmd = [
        "pdftoppm",
        "-png",
        "-r", str(dpi),
        str(pdf_path),
        str(prefix),
    ]

    try:
        subprocess.run(cmd, check=True, capture_output=True, text=True)
    except Exception as exc:  # pylint: disable=broad-except
        print(f"  ⚠️  PDF 페이지 렌더링 실패: {exc}")
        return []

    pages = sorted(
        p for p in output_dir.iterdir()
        if p.is_file() and p.suffix.lower() == ".png" and p.name.startswith(prefix_name)
    )
    if not pages:
        print("  ⚠️  렌더링된 페이지가 없습니다.")
        return []

    print(f"  ✅ PDF 전체 페이지 렌더링 완료: {len(pages)}장")
    return pages


def _send_shape_to_back(slide, shape) -> None:
    """shape를 슬라이드 뒤쪽으로 이동한다."""
    try:
        sp_tree = slide.shapes._spTree
        element = shape._element
        sp_tree.remove(element)
        sp_tree.insert(2, element)
    except Exception:
        pass


def apply_pdf_backgrounds(prs, page_images: list[Path]) -> None:
    """템플릿 각 슬라이드에 PDF 페이지 이미지를 배경으로 적용한다."""
    if not page_images:
        return

    while len(prs.slides) < len(page_images):
        prs.slides.add_slide(prs.slide_layouts[6])

    for idx, slide in enumerate(prs.slides):
        if idx >= len(page_images):
            break
        img_path = page_images[idx]
        try:
            bg = slide.shapes.add_picture(
                str(img_path),
                Emu(0),
                Emu(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )
            _send_shape_to_back(slide, bg)
        except Exception as exc:  # pylint: disable=broad-except
            print(f"  ⚠️  배경 적용 실패(slide {idx + 1}): {exc}")


def add_logo_to_slide(slide, left_inches, top_inches, width_inches, height_inches):
    """SVG 로고를 슬라이드에 추가 (cairosvg로 PNG 변환 후 추가)"""
    if not os.path.exists(LOGO_SVG_PATH):
        print(f"  ⚠️  로고 파일 없음: {LOGO_SVG_PATH}")
        return None

    try:
        # SVG를 PNG로 변환 (지정된 크기로)
        scale_factor = 2  # 고해상도를 위해 2배
        png_data = cairosvg.svg2png(
            url=str(LOGO_SVG_PATH),
            output_width=int(width_inches * 96 * scale_factor),
            output_height=int(height_inches * 96 * scale_factor)
        )

        # BytesIO로 래핑
        from io import BytesIO
        img_stream = BytesIO(png_data)

        # 슬라이드에 추가
        logo_pic = slide.shapes.add_picture(
            img_stream,
            Inches(left_inches),
            Inches(top_inches),
            width=Inches(width_inches),
            height=Inches(height_inches)
        )

        print(f"  ✅ 로고 추가: 왼쪽 상단 ({left_inches}\", {top_inches}\")")
        return logo_pic

    except Exception as e:
        print(f"  ⚠️  로고 추가 실패: {e}")
        return None


def create_title_template_1(prs, cover_image_path: Path | None = None):
    """제목 템플릿 1: 중앙 정렬 큰 제목 + 부제목 + 날짜"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 슬라이드 너비 계산
    slide_width = prs.slide_width / 914400  # EMU -> 인치

    has_cover = False
    if cover_image_path:
        has_cover = set_slide_full_background(slide, cover_image_path, width=prs.slide_width, height=prs.slide_height)
        print(f"  ✅ 표지 이미지 적용: {cover_image_path}")
        print(f"  ✅ 표지 이미지 적용 has_cover: {has_cover}")

    if not has_cover:
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width,
            prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(235, 241, 246)
        background.line.fill.background()

    # 로고 추가 (왼쪽 상단)
    add_logo_to_slide(slide, 0.3, 0.3, 1.2, 0.8)

    # 제목 (중앙 상단) - 전체 너비 사용
    title_left = Inches(0.5)
    title_top = Inches(2.0)
    title_width = Inches(slide_width - 1)  # 좌우 0.5인치 여백
    title_height = Inches(1.2)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "제목을 입력하세요"
    p.alignment = PP_ALIGN.CENTER

    # Run 수준에서 폰트 설정
    run = p.runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255) if has_cover else RGBColor(68, 114, 196)

    # 부제목 (중앙)
    subtitle_left = Inches(0.5)
    subtitle_top = Inches(3.4)
    subtitle_width = Inches(slide_width - 1)
    subtitle_height = Inches(0.8)

    subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True

    p = subtitle_frame.paragraphs[0]
    p.text = "부제목을 입력하세요"
    p.alignment = PP_ALIGN.CENTER

    # Run 수준에서 폰트 설정
    run = p.runs[0]
    run.font.size = Pt(26)
    run.font.color.rgb = RGBColor(235, 235, 235) if has_cover else RGBColor(68, 114, 196)

    # 날짜/추가 정보 (하단)
    date_left = Inches(0.5)
    date_top = Inches(4.4)
    date_width = Inches(slide_width - 1)
    date_height = Inches(0.6)

    date_box = slide.shapes.add_textbox(date_left, date_top, date_width, date_height)
    date_frame = date_box.text_frame
    date_frame.word_wrap = True

    p = date_frame.paragraphs[0]
    p.text = "날짜 또는 추가 정보"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(230, 230, 230) if has_cover else RGBColor(89, 89, 89)

    return slide


def _set_fill(shape, rgb: tuple[int, int, int], line_none: bool = True):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    if line_none:
        shape.line.fill.background()


def _add_rich_text(tf, parts, font_size=14, color=(0, 0, 0), line_spacing=1.2):
    """
    parts: [(text, bold_bool), ...]
    """
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing

    # python-pptx는 paragraph.text 할당 후 run 조작이 까다로워서 run 단위로 구성
    for i, (txt, is_bold) in enumerate(parts):
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = bool(is_bold)
        run.font.color.rgb = RGBColor(*color)

    return p




def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _fill_rect(shape, rgb, line=False, line_rgb=(255, 255, 255), line_width_pt=1):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(*rgb)
    if line:
        shape.line.color.rgb = _rgb(*line_rgb)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()


def _add_textbox(slide, left, top, width, height, text="", font_size=12, bold=False,
                 color=(0, 0, 0), align=PP_ALIGN.LEFT, line_spacing=1.15, word_wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    if p.runs:
        r = p.runs[0]
    else:
        r = p.add_run()
    r.font.size = Pt(font_size)
    r.font.bold = bool(bold)
    r.font.color.rgb = _rgb(*color)
    return tb


def _add_rich_parts(tf, parts, font_size=10.5, color=(40, 40, 40), line_spacing=1.15):
    """
    parts: [(text, bold_bool), ...]
    """
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing

    for txt, is_bold in parts:
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = bool(is_bold)
        run.font.color.rgb = _rgb(*color)


def _rgb(t):
    return RGBColor(t[0], t[1], t[2])


def _set_cell_fill(cell, rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(rgb)


def _set_cell_text(cell, text, font_size=10.5, bold=False, color=(0, 0, 0),
                   align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.12):
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing

    # run 스타일
    if p.runs:
        r = p.runs[0]
    else:
        r = p.add_run()
    r.font.size = Pt(font_size)
    r.font.bold = bool(bold)
    r.font.color.rgb = _rgb(color)


def _add_rich_runs(cell, parts, font_size=10.2, color=(50, 50, 50),
                   align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.12):
    """
    parts: [(text, bold_bool), ...]
    """
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing

    for txt, is_bold in parts:
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = bool(is_bold)
        run.font.color.rgb = _rgb(color)


def create_executive_summary_table_template(
    prs,
    page_no: int = 2,
    left_header: str = "BOLT #2 – Phase 1",
    right_header: str = "BOLT #2 – Phase 2",
):
    """
    Executive Summary 템플릿(표 기반)
    - 표: 8행 x 4열 (0행 헤더 + 1~7행 항목)
      [번호|내용|번호|내용]
    """

    # 색상
    BLUE_TITLE = (0x1A, 0x56, 0xC8)
    BLUE_BAR = (0x0E, 0x63, 0xB5)
    TEXT = (0x33, 0x33, 0x33)
    HILITE = (0xFF, 0xF2, 0x00)
    PAGE_GREEN = (0x00, 0xB0, 0x50)

    # 슬라이드 생성
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w = prs.slide_width / 914400
    slide_h = prs.slide_height / 914400

    # 타이틀
    title_box = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.35), Inches(slide_w - 1.10), Inches(0.55)
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.size = Pt(26)
    r.font.bold = False
    r.font.color.rgb = _rgb(BLUE_TITLE)

    # 테이블 영역
    left_margin = 0.55
    right_margin = 0.55
    top = 0.95
    bottom_margin = 0.55

    table_left = left_margin
    table_top = top
    table_w = slide_w - left_margin - right_margin
    table_h = slide_h - table_top - bottom_margin

    rows = 8   # header + 7 items
    cols = 4   # numL | textL | numR | textR

    tbl_shape = slide.shapes.add_table(
        rows, cols,
        Inches(table_left), Inches(table_top),
        Inches(table_w), Inches(table_h)
    )
    table = tbl_shape.table

    # 열 너비 (번호열은 좁게)
    num_w = 0.32
    text_w = (table_w - 2 * num_w) / 2
    table.columns[0].width = Inches(num_w)
    table.columns[1].width = Inches(text_w)
    table.columns[2].width = Inches(num_w)
    table.columns[3].width = Inches(text_w)

    # 0행 헤더 병합
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 2).merge(table.cell(0, 3))

    # 헤더 스타일
    for c in range(cols):
        _set_cell_fill(table.cell(0, c), BLUE_BAR)

    _set_cell_text(
        table.cell(0, 0),
        left_header,
        font_size=13,
        bold=True,
        color=(255, 255, 255),
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0
    )
    _set_cell_text(
        table.cell(0, 2),
        right_header,
        font_size=13,
        bold=True,
        color=(255, 255, 255),
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0
    )

    # 샘플 데이터(템플릿용) — 필요 시 실제 데이터로 교체
    left_items = []
    right_items = []

    return slide


def create_approval_request_template(prs, page_no: int = 3):
    """
    첨부 이미지(Approval request) 스타일 템플릿 슬라이드 생성
    - Title: "Approval request"
    - 2 rows:
      [Navy square label] + [Light gray explanation box]
    - Footer page number (green) at bottom-left
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # -------------------------
    # Base / constants (16:9 기준)
    # -------------------------
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    # Colors (샘플 톤)
    BLUE_TITLE = (0x12, 0x4E, 0xB7)   # 타이틀 블루
    NAVY = (0x06, 0x1D, 0x3A)         # 좌측 라벨 박스
    LIGHT_GRAY = (0xF2, 0xF2, 0xF2)   # 우측 설명 박스
    TEXT_GRAY = (0x33, 0x33, 0x33)
    GREEN_PAGE = (0x00, 0xB0, 0x50)   # 하단 페이지 번호 (샘플처럼 초록)

    # Layout (인치)
    margin_l = 0.85
    margin_r = 0.85
    title_top = 0.75

    row_top_1 = 1.85
    row_h = 1.15
    row_gap = 0.28

    label_w = 1.25
    gap = 0.25

    content_w = (prs.slide_width / 914400) - margin_l - margin_r - label_w - gap
    # content_w는 인치값, 아래에서 Inches로 변환

    # -------------------------
    # Background (white)
    # -------------------------
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    _set_fill(bg, (255, 255, 255), line_none=True)

    # -------------------------
    # Title
    # -------------------------
    title_box = slide.shapes.add_textbox(
        Inches(margin_l), Inches(title_top), Inches(8.0), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Approval request"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = False
    run.font.color.rgb = RGBColor(*BLUE_TITLE)

    # -------------------------
    # Helper to draw one row
    # -------------------------
    def add_row(y_top_in, label_text, rich_parts):
        # Left label square
        label = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(margin_l), Inches(y_top_in),
            Inches(label_w), Inches(row_h)
        )
        _set_fill(label, NAVY, line_none=True)

        lt = slide.shapes.add_textbox(
            Inches(margin_l), Inches(y_top_in),
            Inches(label_w), Inches(row_h)
        )
        ltf = lt.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lp.text = label_text
        lp.alignment = PP_ALIGN.CENTER
        lp.line_spacing = 1.1
        # runs가 여러 개로 쪼개지는 경우(개행 등) 모두 동일 스타일을 입혀 색상 차이를 방지
        for lrun in lp.runs:
            lrun.font.size = Pt(13)
            lrun.font.bold = True
            lrun.font.color.rgb = RGBColor(255, 255, 255)

        # Right explanation box
        right_left = margin_l + label_w + gap
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(right_left), Inches(y_top_in),
            Inches(content_w), Inches(row_h)
        )
        _set_fill(box, LIGHT_GRAY, line_none=True)

        # Text inside right box (with padding 느낌을 위해 약간 inset)
        pad_x = 0.25
        pad_y = 0.18
        text = slide.shapes.add_textbox(
            Inches(right_left + pad_x), Inches(y_top_in + pad_y),
            Inches(content_w - 2 * pad_x), Inches(row_h - 2 * pad_y)
        )
        ttf = text.text_frame
        ttf.word_wrap = True
        _add_rich_text(
            ttf,
            rich_parts,
            font_size=13,
            color=TEXT_GRAY,
            line_spacing=1.2
        )

    # -------------------------
    # Row 1 (Expected return)
    # -------------------------
    add_row(
        row_top_1,
        "Expected\nreturn",
        [
            ("Approve the expected return based on the latest revenues and capex:\n", False),
            ("- Approval criterion: ", False),
            ("Portfolio blended Equity IRR of 9.64% (≈9.65%)", True),
            ("\n- Average IRR (auxiliary, unweighted average of project IRRs): ", False),
            ("8.95%", True),
        ],
    )

    # -------------------------
    # Row 2 (CAPEX)
    # -------------------------
    add_row(
        row_top_1 + row_h + row_gap,
        "CAPEX",
        [
            ("Approval of total capex of ", False),
            ("KRW 5,086M", True),
            (" approximately ", False),
            ("EUR 3,459,727", True),
            (" for the 4.203MWp portfolio", False),
        ],
    )

    # -------------------------
    # Footer page number (bottom-left)
    # -------------------------
    footer = slide.shapes.add_textbox(
        Inches(0.9), Inches((prs.slide_height / 914400) - 0.45),
        Inches(1.0), Inches(0.3)
    )
    ftf = footer.text_frame
    ftf.clear()
    fp = ftf.paragraphs[0]
    fp.text = str(page_no)
    fp.alignment = PP_ALIGN.LEFT
    frun = fp.runs[0]
    frun.font.size = Pt(11)
    frun.font.color.rgb = RGBColor(*GREEN_PAGE)

    return slide


def create_session_title_template(prs, bg_path: Path = SESSION_BG_PATH):
    """세션 구분용 타이틀 페이지 (배경 이미지 + 좌측 정렬 텍스트)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w_inch = prs.slide_width / 914400
    slide_h_inch = prs.slide_height / 914400

    has_bg = set_slide_full_background(slide, bg_path, width=prs.slide_width, height=prs.slide_height) if bg_path else False

    if not has_bg:
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width,
            prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(10, 40, 80)
        background.line.fill.background()

    # 텍스트 박스 (좌측, 수직 중앙 부근)
    box_left = Inches(1)
    box_width = Inches(slide_w_inch * 0.45)
    box_height = Inches(2.0)
    box_top = Inches(slide_h_inch * 0.45)

    textbox = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "01\n세션 제목"
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.1

    run = p.runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255) if has_bg else RGBColor(68, 114, 196)

    return slide


def create_title_template_2(prs):
    """제목 템플릿 2: 왼쪽 정렬 제목 + 오른쪽 데코레이션"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 왼쪽 데코레이션 바
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        Inches(0.3),
        prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(68, 114, 196)
    bar.line.fill.background()

    # 제목 (왼쪽 상단)
    title_left = Inches(1)
    title_top = Inches(2.5)
    title_width = Inches(8)
    title_height = Inches(1)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "제목을 입력하세요"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 부제목 (왼쪽 하단)
    subtitle_left = Inches(1)
    subtitle_top = Inches(3.8)
    subtitle_width = Inches(8)
    subtitle_height = Inches(0.6)

    subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True

    p = subtitle_frame.paragraphs[0]
    p.text = "부제목을 입력하세요"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 날짜 (오른쪽 하단)
    slide_width = prs.slide_width / 914400
    date_left = Inches(slide_width - 3.5)
    date_top = Inches(6)
    date_width = Inches(3)
    date_height = Inches(0.5)

    date_box = slide.shapes.add_textbox(date_left, date_top, date_width, date_height)
    date_frame = date_box.text_frame
    date_frame.word_wrap = True

    p = date_frame.paragraphs[0]
    p.text = "YYYY년 MM월"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(89, 89, 89)

    return slide


def create_content_basic(prs):
    """내용 템플릿 1: 제목 + 본문 텍스트 (기본)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 본문 영역
    content_left = Inches(MARGIN)
    content_top = Inches(CONTENT_TOP)
    body_width = Inches(content_width)
    body_height = Inches(slide_height - CONTENT_TOP - MARGIN)

    content_box = slide.shapes.add_textbox(content_left, content_top, body_width, body_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    p = content_frame.paragraphs[0]
    p.text = "• 내용을 입력하세요\n• 여러 줄 작성 가능\n• 줄간격 자동 조정"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.5

    # 텍스트 박스 내부 여백 (python-pptx에서는 직접 설정 불가, 들여쓰기로 대체)
    # 첫 줄에 공백 추가 대신 텍스트 들여쓰기 사용
    # 텍스트 프레임의 margin 설정 (지원되지 않을 수 있음)
    try:
        content_frame.margin_left = Inches(0.1)
    except:
        pass

    return slide


def create_content_two_column_text(prs):
    """내용 템플릿 2: 제목 + 2단 내용 (왼쪽/오른쪽 텍스트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 두 칼럼 크기 계산
    column_gap = 0.5
    left_width = (content_width - column_gap) / 2
    right_width = left_width

    # 왼쪽 칼럼
    left_left = Inches(MARGIN)
    left_top = Inches(CONTENT_TOP)
    left_height = Inches(slide_height - CONTENT_TOP - MARGIN)

    left_box = slide.shapes.add_textbox(left_left, left_top, Inches(left_width), left_height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "왼쪽 내용"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.4
    try:
        left_frame.margin_left = Inches(0.1)
    except:
        pass

    # 오른쪽 칼럼
    right_left = MARGIN + left_width + column_gap
    right_top = Inches(CONTENT_TOP)
    right_height = Inches(slide_height - CONTENT_TOP - MARGIN)

    right_box = slide.shapes.add_textbox(Inches(right_left), right_top, Inches(right_width), right_height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    p.text = "오른쪽 내용"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.4
    try:
        right_frame.margin_left = Inches(0.1)
    except:
        pass

    return slide


def create_content_three_column_text(prs):
    """내용 템플릿 3: 제목 + 3단 내용"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 세 칼럼 크기 계산
    column_gap = 0.3
    col_width = (content_width - 2 * column_gap) / 3

    # 칼럼 1
    col1_left = Inches(MARGIN)
    col1_top = Inches(CONTENT_TOP)
    col1_height = Inches(slide_height - CONTENT_TOP - MARGIN)

    col1_box = slide.shapes.add_textbox(col1_left, col1_top, Inches(col_width), col1_height)
    col1_frame = col1_box.text_frame
    col1_frame.word_wrap = True

    p = col1_frame.paragraphs[0]
    p.text = "첫 번째\n내용"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.3

    # 칼럼 2
    col2_left = MARGIN + col_width + column_gap
    col2_top = Inches(CONTENT_TOP)
    col2_height = Inches(slide_height - CONTENT_TOP - MARGIN)

    col2_box = slide.shapes.add_textbox(Inches(col2_left), col2_top, Inches(col_width), col2_height)
    col2_frame = col2_box.text_frame
    col2_frame.word_wrap = True

    p = col2_frame.paragraphs[0]
    p.text = "두 번째\n내용"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.3

    # 칼럼 3
    col3_left = MARGIN + 2 * (col_width + column_gap)
    col3_top = Inches(CONTENT_TOP)
    col3_height = Inches(slide_height - CONTENT_TOP - MARGIN)

    col3_box = slide.shapes.add_textbox(Inches(col3_left), col3_top, Inches(col_width), col3_height)
    col3_frame = col3_box.text_frame
    col3_frame.word_wrap = True

    p = col3_frame.paragraphs[0]
    p.text = "세 번째\n내용"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.3

    return slide


def create_content_large_table(prs):
    """내용 템플릿 4: 제목 + 대형 테이블 (페이지 꽉 찬 테이블)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.7)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 표 영역
    table_width = content_width
    table_height = slide_height - CONTENT_TOP - MARGIN

    table_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN),
        Inches(CONTENT_TOP),
        Inches(table_width),
        Inches(table_height)
    )
    table_placeholder.fill.solid()
    table_placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
    table_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    table_placeholder.line.width = Pt(2)

    # 표 플레이스홀더 텍스트
    table_frame = table_placeholder.text_frame
    table_frame.word_wrap = True

    p = table_frame.paragraphs[0]
    p.text = "표영역"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(150, 150, 150)

    return slide


def create_content_two_tables(prs):
    """내용 템플릿 5: 제목 + 테이블 2개 (나란히)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.7)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 두 테이블 영역
    table_gap = 0.4
    table_width = (content_width - table_gap) / 2
    table_height = slide_height - CONTENT_TOP - MARGIN

    # 왼쪽 테이블
    table1_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN),
        Inches(CONTENT_TOP),
        Inches(table_width),
        Inches(table_height)
    )
    table1_placeholder.fill.solid()
    table1_placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
    table1_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    table1_placeholder.line.width = Pt(2)

    table1_frame = table1_placeholder.text_frame
    table1_frame.word_wrap = True

    p = table1_frame.paragraphs[0]
    p.text = "표영역1"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(150, 150, 150)

    # 오른쪽 테이블
    table2_left = MARGIN + table_width + table_gap
    table2_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(table2_left),
        Inches(CONTENT_TOP),
        Inches(table_width),
        Inches(table_height)
    )
    table2_placeholder.fill.solid()
    table2_placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
    table2_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    table2_placeholder.line.width = Pt(2)

    table2_frame = table2_placeholder.text_frame
    table2_frame.word_wrap = True

    p = table2_frame.paragraphs[0]
    p.text = "표영역2"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(150, 150, 150)

    return slide


def create_content_gantt_chart(prs):
    """내용 템플릿 6: 제목 + 간트 차트 영역 + 하단 텍스트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.7)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 간트 차트 영역 (상단)
    gantt_height = 4.5
    gantt_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN),
        Inches(CONTENT_TOP),
        Inches(content_width),
        Inches(gantt_height)
    )
    gantt_placeholder.fill.solid()
    gantt_placeholder.fill.fore_color.rgb = RGBColor(248, 250, 252)
    gantt_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    gantt_placeholder.line.width = Pt(2)

    gantt_frame = gantt_placeholder.text_frame
    gantt_frame.word_wrap = True

    p = gantt_frame.paragraphs[0]
    p.text = "간트 차트 / 타임라인 영역"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(150, 150, 150)

    # 하단 텍스트 영역
    text_top = CONTENT_TOP + gantt_height + 0.2
    text_height = slide_height - text_top - MARGIN

    text_box = slide.shapes.add_textbox(
        Inches(MARGIN),
        Inches(text_top),
        Inches(content_width),
        Inches(text_height)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "• 하단 설명 텍스트\n• 간트 차트에 대한 추가 정보"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.3

    return slide


def create_content_text_small_table(prs):
    """내용 템플릿 7: 제목 + 본문 + 작은 테이블 (오른쪽)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 왼쪽 본문 영역 (60%)
    text_width = content_width * 0.55
    text_height = slide_height - CONTENT_TOP - MARGIN

    text_box = slide.shapes.add_textbox(
        Inches(MARGIN),
        Inches(CONTENT_TOP),
        Inches(text_width),
        Inches(text_height)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "• 왼쪽 본문 내용\n• 텍스트 영역\n• 여러 줄 작성 가능"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.4

    # 오른쪽 작은 테이블 영역 (40%)
    table_left = MARGIN + text_width + 0.3
    table_width = content_width - text_width - 0.3

    table_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(table_left),
        Inches(CONTENT_TOP),
        Inches(table_width),
        Inches(text_height)
    )
    table_placeholder.fill.solid()
    table_placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
    table_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    table_placeholder.line.width = Pt(2)

    table_frame = table_placeholder.text_frame
    table_frame.word_wrap = True

    p = table_frame.paragraphs[0]
    p.text = "표영역"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(150, 150, 150)

    return slide


def create_content_text_large_table(prs):
    """내용 템플릿 8: 제목 + 상단 텍스트 + 대형 하단 테이블"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.7)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 상단 짧은 텍스트 영역
    text_top = CONTENT_TOP
    text_height = 1.2

    text_box = slide.shapes.add_textbox(
        Inches(MARGIN),
        Inches(text_top),
        Inches(content_width),
        Inches(text_height)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    # p.text = "• 상단 설명 텍스트 (간단)"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(64, 64, 64)
    p.line_spacing = 1.3

    # 하단 대형 테이블 영역
    table_top = text_top + text_height + 0.2
    table_height = slide_height - table_top - MARGIN

    table_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN),
        Inches(table_top),
        Inches(content_width),
        Inches(table_height)
    )
    table_placeholder.fill.solid()
    table_placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
    table_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    table_placeholder.line.width = Pt(2)

    table_frame = table_placeholder.text_frame
    table_frame.word_wrap = True

    p = table_frame.paragraphs[0]
    p.text = "표영역"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(150, 150, 150)

    return slide


def create_content_two_tables_vertical(prs):
    """내용 템플릿 9: 제목 + 2개 표 (상하 세로 분할)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.7)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정 (python-pptx에서는 run.level에서 설정해야 함)
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 두 테이블 영역 (상하 분할)
    table_gap = 0.3
    available_height = slide_height - CONTENT_TOP - MARGIN
    table_height = (available_height - table_gap) / 2

    # 상단 테이블
    table1_top = CONTENT_TOP
    table1_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN),
        Inches(table1_top),
        Inches(content_width),
        Inches(table_height)
    )
    table1_placeholder.fill.solid()
    table1_placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
    table1_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    table1_placeholder.line.width = Pt(2)

    table1_frame = table1_placeholder.text_frame
    table1_frame.word_wrap = True

    p = table1_frame.paragraphs[0]
    p.text = "표영역1"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(150, 150, 150)

    # 하단 테이블
    table2_top = table1_top + table_height + table_gap
    table2_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN),
        Inches(table2_top),
        Inches(content_width),
        Inches(table_height)
    )
    table2_placeholder.fill.solid()
    table2_placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
    table2_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    table2_placeholder.line.width = Pt(2)

    table2_frame = table2_placeholder.text_frame
    table2_frame.word_wrap = True

    p = table2_frame.paragraphs[0]
    p.text = "표영역2"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(150, 150, 150)

    return slide


def create_content_title_subtitle_table(prs):
    """내용 템플릿: 제목 + 부제목 + 대형 테이블"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(MARGIN), Inches(content_width), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    subtitle_top = MARGIN + 0.62
    subtitle_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(subtitle_top), Inches(content_width), Inches(0.45)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    sp = subtitle_frame.paragraphs[0]
    sp.text = "부제목을 입력하세요"
    sp.alignment = PP_ALIGN.LEFT
    srun = sp.runs[0]
    srun.font.size = Pt(16)
    srun.font.bold = False
    srun.font.color.rgb = RGBColor(90, 90, 90)

    table_top = subtitle_top + 0.58
    table_height = slide_height - table_top - MARGIN
    table_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN),
        Inches(table_top),
        Inches(content_width),
        Inches(table_height)
    )
    table_placeholder.fill.solid()
    table_placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
    table_placeholder.line.color.rgb = RGBColor(200, 200, 200)
    table_placeholder.line.width = Pt(2)

    table_frame = table_placeholder.text_frame
    table_frame.word_wrap = True
    tp = table_frame.paragraphs[0]
    tp.text = "표영역"
    tp.alignment = PP_ALIGN.CENTER
    tp.font.size = Pt(24)
    tp.font.color.rgb = RGBColor(150, 150, 150)

    return slide


def create_content_image_left(prs):
    """내용 템플릿 10: 제목 + 본문 (왼쪽) + 이미지 (오른쪽)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 콘텐츠 영역 계산 (제목 아래부터)
    CONTENT_TOP_LOCAL = CONTENT_TOP
    CONTENT_HEIGHT = slide_height - CONTENT_TOP_LOCAL - MARGIN

    # 본문 영역 (왼쪽 50%)
    text_width = content_width * 0.5
    content_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(CONTENT_TOP_LOCAL),
        Inches(text_width), Inches(CONTENT_HEIGHT)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    p = content_frame.paragraphs[0]
    p.text = "• 본문 내용\n• 여기에 입력\n• 이미지는 오른쪽"
    p.font.size = Pt(18)
    p.line_spacing = 1.5

    # 이미지 영역 (오른쪽 50% - 0.3인치 간격)
    image_width = content_width * 0.5 - 0.3
    image_box = slide.shapes.add_textbox(
        Inches(MARGIN + text_width + 0.3), Inches(CONTENT_TOP_LOCAL),
        Inches(image_width), Inches(CONTENT_HEIGHT)
    )
    image_frame = image_box.text_frame
    image_frame.word_wrap = True  # 텍스트 랩핑 활성화하여 박스 크기 유지
    p = image_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)

    # 수직 중앙 정렬을 위한 여백 및 박스 크기 유지를 위한 더미 텍스트
    p.text = "이미지영역\n\n[이미지가 들어갈 자리]\n\n(이 텍스트는 이미지로 대체됩니다)"

    return slide


def create_content_image_right(prs):
    """내용 템플릿 11: 제목 + 이미지 (왼쪽) + 본문 (오른쪽)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 콘텐츠 영역 계산 (제목 아래부터)
    CONTENT_TOP_LOCAL = CONTENT_TOP
    CONTENT_HEIGHT = slide_height - CONTENT_TOP_LOCAL - MARGIN

    # 이미지 영역 (왼쪽 50% - 0.3인치 간격)
    image_width = content_width * 0.5 - 0.3
    image_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(CONTENT_TOP_LOCAL),
        Inches(image_width), Inches(CONTENT_HEIGHT)
    )
    image_frame = image_box.text_frame
    image_frame.word_wrap = True  # 텍스트 랩핑 활성화하여 박스 크기 유지
    p = image_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.text = "이미지영역\n\n[이미지가 들어갈 자리]\n\n(이 텍스트는 이미지로 대체됩니다)"

    # 본문 영역 (오른쪽 50%)
    text_width = content_width * 0.5
    content_box = slide.shapes.add_textbox(
        Inches(MARGIN + image_width + 0.3), Inches(CONTENT_TOP_LOCAL),
        Inches(text_width), Inches(CONTENT_HEIGHT)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    p = content_frame.paragraphs[0]
    p.text = "• 본문 내용\n• 여기에 입력\n• 이미지는 왼쪽"
    p.font.size = Pt(18)
    p.line_spacing = 1.5

    return slide


def create_content_large_image(prs):
    """내용 템플릿 12: 제목 + 큰 이미지 영역 (전체)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 큰 이미지 영역 (중앙)
    CONTENT_TOP_LOCAL = CONTENT_TOP
    CONTENT_HEIGHT = slide_height - CONTENT_TOP_LOCAL - MARGIN

    image_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(CONTENT_TOP_LOCAL),
        Inches(content_width), Inches(CONTENT_HEIGHT)
    )
    image_frame = image_box.text_frame
    image_frame.word_wrap = True  # 텍스트 랩핑 활성화하여 박스 크기 유지
    p = image_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)

    # 수직/수평 중앙 정렬 및 박스 크기 유지를 위한 더미 텍스트
    p.text = "이미지영역\n\n\n[큰 이미지가 들어갈 자리]\n\n\n(이 텍스트는 이미지로 대체됩니다)"

    return slide


def create_content_two_images(prs):
    """내용 템플릿 13: 제목 + 이미지 2개 (좌우 배치)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    # 배경색 (흰색)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # 제목 (상단)
    title_left = Inches(MARGIN)
    title_top = Inches(MARGIN)
    title_width = Inches(content_width)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT

    # Run 수준에서 폰트 설정
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    # 이미지 영역 계산
    CONTENT_TOP_LOCAL = CONTENT_TOP
    CONTENT_HEIGHT = slide_height - CONTENT_TOP_LOCAL - MARGIN
    image_gap = 0.5
    image_width = (content_width - image_gap) / 2  # 이미지 2개로 나누기

    # 왼쪽 이미지 영역
    left_image_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(CONTENT_TOP_LOCAL),
        Inches(image_width), Inches(CONTENT_HEIGHT)
    )
    left_frame = left_image_box.text_frame
    left_frame.word_wrap = True  # 텍스트 랩핑 활성화하여 박스 크기 유지
    p = left_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.text = "이미지영역1\n\n[왼쪽 이미지]\n\n(이 텍스트는 이미지로 대체됩니다)"

    # 오른쪽 이미지 영역
    right_image_box = slide.shapes.add_textbox(
        Inches(MARGIN + image_width + image_gap), Inches(CONTENT_TOP_LOCAL),
        Inches(image_width), Inches(CONTENT_HEIGHT)
    )
    right_frame = right_image_box.text_frame
    right_frame.word_wrap = True  # 텍스트 랩핑 활성화하여 박스 크기 유지
    p = right_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.text = "이미지영역2\n\n[오른쪽 이미지]\n\n(이 텍스트는 이미지로 대체됩니다)"

    return slide


def create_main_agreements_template(prs):
    """Main Agreements 전용 템플릿 (내용은 이후 렌더러가 대체)"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.3), Inches(SLIDE_WIDTH_16_9 - 2 * MARGIN), Inches(0.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Main Agreements"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x56, 0xC8)
    return slide


def create_bos_negotiation_template(prs, base_image_path: Path | None = None):
    """BOS negotiation 전용 템플릿 (표 데이터는 공란 유지)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / 914400
    slide_height = prs.slide_height / 914400
    content_width = slide_width - 2 * MARGIN

    if base_image_path and Path(base_image_path).exists():
        set_slide_full_background(
            slide,
            Path(base_image_path),
            width=prs.slide_width,
            height=prs.slide_height,
        )
    else:
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            prs.slide_height,
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background.line.fill.background()

    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(MARGIN - 0.06),
        Inches(MARGIN - 0.06),
        Inches(content_width + 0.12),
        Inches(1.05),
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(MARGIN), Inches(content_width), Inches(0.65)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "슬라이드 제목"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(68, 114, 196)

    subtitle_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(MARGIN + 0.64), Inches(content_width), Inches(0.35)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    sp = subtitle_frame.paragraphs[0]
    sp.text = "부제목을 입력하세요"
    sp.alignment = PP_ALIGN.LEFT
    srun = sp.runs[0]
    srun.font.size = Pt(14)
    srun.font.color.rgb = RGBColor(90, 90, 90)

    table_top = CONTENT_TOP + 0.05
    table_height = slide_height - table_top - MARGIN
    table_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(MARGIN),
        Inches(table_top),
        Inches(content_width),
        Inches(table_height),
    )
    table_placeholder.fill.solid()
    table_placeholder.fill.fore_color.rgb = RGBColor(245, 247, 250)
    table_placeholder.line.color.rgb = RGBColor(200, 208, 220)
    table_placeholder.line.width = Pt(1.5)

    tf = table_placeholder.text_frame
    tf.word_wrap = True
    tp = tf.paragraphs[0]
    tp.text = "표영역"
    tp.alignment = PP_ALIGN.CENTER
    tp.font.size = Pt(22)
    tp.font.color.rgb = RGBColor(150, 150, 150)

    return slide


def create_templates():
    """모든 템플릿 생성 - 하나의 통합 파일로 (제목 + 내용)"""
    print("=" * 80)
    print("Custom PPTX 템플릿 생성 (4:3, 16:9) - 통합 템플릿 방식")
    print("=" * 80)

    # 출력 디렉토리 생성
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    result = {}
    cover_image_path = COVER_IMAGE_PATH if COVER_IMAGE_PATH.exists() else render_pdf_cover(PDF_COVER_PATH, PDF_COVER_OUTPUT)

    # ============ 16:9 비율 템플릿 ============
    print("\n" + "="*80)
    print("[16:9 비율 (13.333\" × 7.5\")]")
    print("="*80)

    # 통합 템플릿 파일 (16:9)
    prs_169 = Presentation()
    prs_169.slide_width = Inches(SLIDE_WIDTH_16_9)
    prs_169.slide_height = Inches(SLIDE_HEIGHT)

    print("\n[제목 템플릿]")
    print("  슬라이드 0: 제목 템플릿 1 (중앙 정렬 + 날짜)")
    create_title_template_1(prs_169, cover_image_path)
    print("  슬라이드 1: 제목 템플릿 2 (왼쪽 정렬 + 데코레이션)")
    create_title_template_2(prs_169)

    print("\n[내용 템플릿]")
    print("  슬라이드 3: 내용 템플릿 1 (제목 + 본문)")
    create_content_basic(prs_169)
    print("  슬라이드 4: 내용 템플릿 2 (제목 + 2단 텍스트)")
    create_content_two_column_text(prs_169)
    print("  슬라이드 5: 내용 템플릿 3 (제목 + 3단 텍스트)")
    create_content_three_column_text(prs_169)
    print("  슬라이드 6: 내용 템플릿 4 (제목 + 대형 테이블)")
    create_content_large_table(prs_169)
    print("  슬라이드 7: 내용 템플릿 5 (제목 + 표 2개 나란히)")
    create_content_two_tables(prs_169)
    print("  슬라이드 8: 내용 템플릿 6 (제목 + 간트 차트 + 하단 텍스트)")
    create_content_gantt_chart(prs_169)
    print("  슬라이드 9: 내용 템플릿 7 (제목 + 본문 + 작은 테이블)")
    create_content_text_small_table(prs_169)
    print("  슬라이드 10: 내용 템플릿 8 (제목 + 상단 텍스트 + 대형 테이블)")
    create_content_text_large_table(prs_169)
    print("  슬라이드 11: 내용 템플릿 9 (제목 + 표 2개 상하)")
    create_content_two_tables_vertical(prs_169)
    print("  슬라이드 12: 내용 템플릿 10 (제목 + 본문 왼쪽 + 이미지 오른쪽)")
    create_content_image_left(prs_169)
    print("  슬라이드 13: 내용 템플릿 11 (제목 + 이미지 왼쪽 + 본문 오른쪽)")
    create_content_image_right(prs_169)
    print("  슬라이드 14: 내용 템플릿 12 (제목 + 큰 이미지)")
    create_content_large_image(prs_169)
    print("  슬라이드 15: 내용 템플릿 13 (제목 + 이미지 2개)")
    create_content_two_images(prs_169)
    
    print("  슬라이드 16: 세션 타이틀 (배경이미지)")
    create_session_title_template(prs_169, SESSION_BG_PATH)
    
    print("  슬라이드 17: Approval request 템플릿")
    create_approval_request_template(prs_169, page_no=3)
    
    
    
    print("  슬라이드 XX: Executive Summary (table)")
    create_executive_summary_table_template(prs_169, page_no=2)

    print("  슬라이드 19: Main Agreements 템플릿")
    create_main_agreements_template(prs_169)

    print("  슬라이드 20: 제목 + 부제목 + 테이블 템플릿")
    create_content_title_subtitle_table(prs_169)

    # 샘플 PDF 전체 슬라이드(페이지)를 템플릿 배경으로 반영
    sample_pages = render_pdf_pages(SAMPLE_CAPEX_PDF_PATH, SAMPLE_PDF_OUTPUT, dpi=SAMPLE_PDF_DPI)
    if sample_pages:
        apply_pdf_backgrounds(prs_169, sample_pages)

    print("  슬라이드 30: BOS negotiation 템플릿 (표 데이터 공란)")
    bos_base = sample_pages[15] if len(sample_pages) >= 16 else None
    create_bos_negotiation_template(prs_169, bos_base)



    # ====================================================================== #
    template_path_169 = os.path.join(OUTPUT_DIR, "templates_16_9.pptx")
    prs_169.save(template_path_169)
    print(f"\n  ✅ 저장: {template_path_169}")
    print(
        f"  총 {len(prs_169.slides)}개 슬라이드 (기본 레이아웃 + 샘플 PDF 배경 전체 페이지 반영)"
    )

    result['16_9'] = template_path_169

    # 요약 출력
    print(f"\n{'='*80}")
    print("Custom 통합 템플릿 생성 완료!")
    print(f"{'='*80}")
    print(f"""
📁 생성된 파일:


【16:9 비율 (13.333" × 7.5")】
   - {template_path_169}
   - 총 {len(prs_169.slides)}개 슬라이드 (기본 레이아웃 + 샘플 PDF 배경 전체 페이지 반영)

🎨 레이아웃 매핑 (airun_pptx.py):
   제목 슬라이드:
   - 0: 제목 템플릿 1 (중앙 정렬 + 날짜)
   - 1: 제목 템플릿 2 (왼쪽 정렬 + 데코레이션)

   내용 슬라이드 (content_list에서 layout으로 지정):
   - 'title' 또는 'basic' 또는 'content_title_text': 2 (제목 + 본문)
   - 'two_column': 3 (제목 + 2단 텍스트)
   - 'three_column': 4 (제목 + 3단 텍스트)
   - 'large_table': 5 (제목 + 대형 테이블)
   - 'two_tables': 6 (제목 + 표 2개 나란히)
   - 'gantt_chart': 7 (제목 + 간트 차트 + 하단 텍스트)
   - 'text_small_table': 8 (제목 + 본문 + 작은 테이블)
   - 'text_large_table': 9 (제목 + 상단 텍스트 + 대형 테이블)
   - 'two_tables_vertical': 10 (제목 + 표 2개 상하)
   - 'image_left': 11 (제목 + 본문 왼쪽 + 이미지 오른쪽)
   - 'image_right': 12 (제목 + 이미지 왼쪽 + 본문 오른쪽)
   - 'large_image': 13 (제목 + 큰 이미지)
   - 'two_images': 14 (제목 + 이미지 2개)
   - 'session_title': 15 (세션 표지)
   - 'approval_request': 16 (승인 요청)
   - 'executive_summary': 17 (Executive Summary)
   - 'main_agreements': 18 (Main Agreements)
   - 'title_subtitle_table': 19 (제목 + 부제목 + 테이블)
   - 'bos_negotiation': 29 (BOS negotiation, 표 데이터 공란)
""")

    return result


def main():
    """메인 실행"""
    print("\n" + "=" * 80)
    print("Custom PPTX 템플릿 생성 - CAPEX/DEVEX Approval 문서용")
    print("=" * 80)

    result = create_templates()

    return result


if __name__ == "__main__":
    main()
