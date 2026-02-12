#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Generate CAPEX Executive Summary (Phase 1 / Phase 2) from INPUTSHEET.xlsx

Outputs:
  - CAPEX_Executive_Summary.pptx
  - CAPEX_Executive_Summary.pdf

Usage:
  python make_exec_summary.py \
    --xlsx "/work/proj/airun/002_poc/airun/pv_solar_sample_docs/지상형태양광_Inputsheet.xlsx" \
    --outdir "/home/no/test" \
    --portfolio "BOLT #2"

Notes:
- Assumes projects are in row 2 starting from column E (5), e.g. "Project 1", "Project 2", ...
- Assumes field labels are in column B.
"""

from __future__ import annotations

import argparse
import os
import re
from dataclasses import dataclass
from datetime import date, datetime
from typing import Any, Dict, List, Optional, Tuple

import openpyxl

# PPTX
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# PDF
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch


# -----------------------------
# Helpers
# -----------------------------
def iso_date(v: Any) -> Optional[str]:
    """Convert Excel cell value to ISO date string if possible."""
    if v is None:
        return None
    if isinstance(v, datetime):
        return v.date().isoformat()
    if isinstance(v, date):
        return v.isoformat()
    # sometimes a string like "2026-01-31" or "31/01/2026"
    if isinstance(v, str):
        s = v.strip()
        # try YYYY-MM-DD
        m = re.match(r"^(\d{4})-(\d{2})-(\d{2})$", s)
        if m:
            return s
        # try DD/MM/YYYY
        m = re.match(r"^(\d{2})/(\d{2})/(\d{4})$", s)
        if m:
            dd, mm, yyyy = m.group(1), m.group(2), m.group(3)
            return f"{yyyy}-{mm}-{dd}"
    return None


def to_float(v: Any) -> Optional[float]:
    if v is None:
        return None
    if isinstance(v, (int, float)):
        return float(v)
    if isinstance(v, str):
        s = v.strip().replace(",", "")
        try:
            return float(s)
        except ValueError:
            return None
    return None


def safe_str(v: Any) -> str:
    if v is None:
        return ""
    return str(v).strip()


def fmt_mwp(x: float) -> str:
    return f"{x:.2f}"


# -----------------------------
# Data model
# -----------------------------
@dataclass
class Project:
    id: str
    country: str = ""
    region: str = ""
    technology: str = ""
    status: str = ""
    transaction_structure: str = ""
    currency: str = ""
    rtb_date: Optional[str] = None
    soc_date: Optional[str] = None
    cod_date: Optional[str] = None
    capacity_mwp: float = 0.0


@dataclass
class PhaseAgg:
    phase: int
    soc_date: Optional[str]
    cod_max: Optional[str]
    project_count: int
    capacity_mwp: float
    regions: List[str]
    projects: List[Project]


# -----------------------------
# Extract from INPUTSHEET
# -----------------------------
FIELD_MAP = {
    "Country": "country",
    "Region": "region",
    "Technology": "technology",
    "Project Status": "status",
    "Transaction structure": "transaction_structure",
    "Currency": "currency",
    "RTB date": "rtb_date",
    "SOC date": "soc_date",
    "COD date": "cod_date",
    "Total Capacity DC": "capacity_mwp",  # MWp
}


def load_inputsheet(xlsx_path: str, sheet_name: Optional[str] = None) -> List[Project]:
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)

    ws = None
    if sheet_name and sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
    else:
        # common name
        if "INPUTSHEET" in wb.sheetnames:
            ws = wb["INPUTSHEET"]
        else:
            # fallback: first sheet
            ws = wb[wb.sheetnames[0]]

    # project columns: row 2, from col E(5) to right until empty streak
    project_cols: List[Tuple[str, int]] = []
    empty_run = 0
    for c in range(5, ws.max_column + 1):
        name = ws.cell(2, c).value
        if name:
            project_cols.append((safe_str(name), c))
            empty_run = 0
        else:
            empty_run += 1
            if empty_run >= 5:  # stop after 5 empty cols
                break

    # field rows: column B (2)
    field_rows: Dict[str, int] = {}
    for r in range(1, ws.max_row + 1):
        label = ws.cell(r, 2).value
        if label in FIELD_MAP and label not in field_rows:
            field_rows[label] = r

    projects: List[Project] = []
    for proj_name, c in project_cols:
        p = Project(id=proj_name)

        for label, key in FIELD_MAP.items():
            r = field_rows.get(label)
            if not r:
                continue
            val = ws.cell(r, c).value

            if key in ("rtb_date", "soc_date", "cod_date"):
                setattr(p, key, iso_date(val))
            elif key == "capacity_mwp":
                p.capacity_mwp = to_float(val) or 0.0
            else:
                setattr(p, key, safe_str(val))

        projects.append(p)

    return projects


def split_phases_by_soc(projects: List[Project]) -> Tuple[List[Project], List[Project], Optional[str], Optional[str]]:
    """Phase 1: earliest SOC date group, Phase 2: rest"""
    socs = sorted({p.soc_date for p in projects if p.soc_date})
    if not socs:
        return projects, [], None, None
    phase1_soc = socs[0]
    phase2_soc = socs[1] if len(socs) > 1 else None

    p1 = [p for p in projects if p.soc_date == phase1_soc]
    p2 = [p for p in projects if p.soc_date != phase1_soc]

    return p1, p2, phase1_soc, phase2_soc


def aggregate_phase(phase: int, soc_date: Optional[str], projects: List[Project]) -> PhaseAgg:
    capacity = sum(p.capacity_mwp for p in projects)
    regions = sorted({p.region for p in projects if p.region})
    cod_dates = sorted([p.cod_date for p in projects if p.cod_date])
    cod_max = cod_dates[-1] if cod_dates else None

    return PhaseAgg(
        phase=phase,
        soc_date=soc_date,
        cod_max=cod_max,
        project_count=len(projects),
        capacity_mwp=capacity,
        regions=regions,
        projects=projects,
    )


# -----------------------------
# Text generation (Executive Summary bullets)
# -----------------------------
def make_exec_bullets(portfolio: str, agg: PhaseAgg) -> List[str]:
    """
    샘플 슬라이드 스타일에 맞추되, INPUTSHEET에 없는 값(REC 가격, IRR 등)은 넣지 않습니다.
    """
    regions = ", ".join(agg.regions) if agg.regions else "N/A"
    soc = agg.soc_date or "N/A"
    cod = agg.cod_max or "N/A"

    return [
        f"In {soc[:7].replace('-', ' ')}, {agg.project_count} rooftop solar PV projects "
        f"with a total capacity of {fmt_mwp(agg.capacity_mwp)} MWp are progressing.",
        f"Projects are located in: {regions}.",
        f"All projects are expected to reach Commercial Operation Date (COD) by {cod}.",
    ]


# -----------------------------
# PPTX generation
# -----------------------------
def build_pptx(portfolio: str, phase1: PhaseAgg, phase2: PhaseAgg, out_pptx: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Page constants (16:9 default)
    W = prs.slide_width
    H = prs.slide_height

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), W - Inches(1.4), Inches(0.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Executive Summary"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x12, 0x4E, 0xB7)  # blue
    p.alignment = PP_ALIGN.LEFT

    # Two columns frame
    margin_x = Inches(0.7)
    top_y = Inches(1.2)
    col_gap = Inches(0.4)

    col_w = (W - 2 * margin_x - col_gap) / 2
    col_h = H - top_y - Inches(0.6)

    def add_phase_col(x, header, bullets: List[str]):
        # Header bar
        header_h = Inches(0.45)
        hdr = slide.shapes.add_shape(
            1, x, top_y, col_w, header_h  # 1 == rect
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = RGBColor(0x12, 0x4E, 0xB7)
        hdr.line.color.rgb = RGBColor(0x12, 0x4E, 0xB7)

        ht = slide.shapes.add_textbox(x, top_y, col_w, header_h)
        htf = ht.text_frame
        htf.clear()
        pp = htf.paragraphs[0]
        rr = pp.add_run()
        rr.text = header
        rr.font.size = Pt(16)
        rr.font.bold = True
        rr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        pp.alignment = PP_ALIGN.CENTER

        # Body bullets
        body = slide.shapes.add_textbox(x, top_y + header_h + Inches(0.1), col_w, col_h - header_h)
        btf = body.text_frame
        btf.word_wrap = True
        btf.clear()

        for i, line in enumerate(bullets, start=1):
            para = btf.paragraphs[0] if i == 1 else btf.add_paragraph()
            para.text = f"{i}. {line}"
            para.level = 0
            para.font.size = Pt(12)
            para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            para.space_after = Pt(6)

    b1 = make_exec_bullets(portfolio, phase1)
    b2 = make_exec_bullets(portfolio, phase2)

    add_phase_col(margin_x, f"{portfolio} – Phase 1", b1)
    add_phase_col(margin_x + col_w + col_gap, f"{portfolio} – Phase 2", b2)

    prs.save(out_pptx)


# -----------------------------
# PDF generation (same content)
# -----------------------------
def build_pdf(portfolio: str, phase1: PhaseAgg, phase2: PhaseAgg, out_pdf: str) -> None:
    c = canvas.Canvas(out_pdf, pagesize=landscape(letter))
    page_w, page_h = landscape(letter)

    # Title
    c.setFont("Helvetica-Bold", 26)
    c.setFillColorRGB(0.07, 0.31, 0.72)
    c.drawString(0.8 * inch, page_h - 0.9 * inch, "Executive Summary")

    # Columns
    margin_x = 0.8 * inch
    top_y = page_h - 1.4 * inch
    col_gap = 0.5 * inch
    col_w = (page_w - 2 * margin_x - col_gap)
    col_w /= 2
    col_h = page_h - 2.0 * inch

    def draw_col(x, header, bullets: List[str]):
        header_h = 0.45 * inch

        # header bar
        c.setFillColorRGB(0.07, 0.31, 0.72)
        c.rect(x, top_y - header_h, col_w, header_h, fill=1, stroke=0)

        c.setFillColorRGB(1, 1, 1)
        c.setFont("Helvetica-Bold", 14)
        c.drawCentredString(x + col_w / 2, top_y - header_h + 0.15 * inch, header)

        # body box
        body_top = top_y - header_h - 0.2 * inch
        c.setFillColorRGB(0, 0, 0)
        c.setFont("Helvetica", 11)

        y = body_top
        for i, line in enumerate(bullets, start=1):
            txt = f"{i}. {line}"
            # simple wrap
            wrapped = wrap_text(txt, max_chars=92)
            for wline in wrapped:
                c.drawString(x, y, wline)
                y -= 0.22 * inch
            y -= 0.06 * inch

    def wrap_text(text: str, max_chars: int) -> List[str]:
        words = text.split()
        lines, cur = [], []
        for w in words:
            if sum(len(t) for t in cur) + len(cur) + len(w) > max_chars:
                lines.append(" ".join(cur))
                cur = [w]
            else:
                cur.append(w)
        if cur:
            lines.append(" ".join(cur))
        return lines

    b1 = make_exec_bullets(portfolio, phase1)
    b2 = make_exec_bullets(portfolio, phase2)

    draw_col(margin_x, f"{portfolio} – Phase 1", b1)
    draw_col(margin_x + col_w + col_gap, f"{portfolio} – Phase 2", b2)

    c.showPage()
    c.save()


# -----------------------------
# Main
# -----------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--xlsx", required=True, help="Path to INPUTSHEET.xlsx")
    ap.add_argument("--sheet", default=None, help="Sheet name (optional)")
    ap.add_argument("--outdir", required=True, help="Output directory")
    ap.add_argument("--portfolio", default="BOLT #2", help="Portfolio name used in headers")
    args = ap.parse_args()

    os.makedirs(args.outdir, exist_ok=True)

    projects = load_inputsheet(args.xlsx, args.sheet)

    p1, p2, soc1, soc2 = split_phases_by_soc(projects)
    agg1 = aggregate_phase(1, soc1, p1)
    agg2 = aggregate_phase(2, soc2, p2)

    out_pptx = os.path.join(args.outdir, "CAPEX_Executive_Summary.pptx")
    out_pdf = os.path.join(args.outdir, "CAPEX_Executive_Summary.pdf")

    build_pptx(args.portfolio, agg1, agg2, out_pptx)
    build_pdf(args.portfolio, agg1, agg2, out_pdf)

    print("[OK] Generated:")
    print(" -", out_pptx)
    print(" -", out_pdf)
    print()
    print("Phase 1:", agg1.project_count, "projects,", fmt_mwp(agg1.capacity_mwp), "MWp, SOC:", agg1.soc_date, "COD:", agg1.cod_max)
    print("Phase 2:", agg2.project_count, "projects,", fmt_mwp(agg2.capacity_mwp), "MWp, SOC:", agg2.soc_date, "COD:", agg2.cod_max)


if __name__ == "__main__":
    main()

