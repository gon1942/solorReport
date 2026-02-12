#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AIRUN PPTX Generator - 깔끔한 템플릿 기반 PPTX 문서 생성 클래스

services/report/report_generator.py에서 사용할 수 있는 PPTX 생성 클래스입니다.

주요 기능:
1. 템플릿 자동 관리 (생성, 로드, 갱신)
2. AI가 자동으로 레이아웃 선택 가능
3. 파라미터로 직접 레이아웃 지정 가능
4. 제목, 본문, 표, 이미지 지원
5. 4:3, 16:9 두 가지 화면 비율 지원
"""

import os
import sys
import shutil
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Union, Any
from dataclasses import dataclass
from enum import Enum
from copy import deepcopy
from datetime import datetime

# 현재 디렉토리를 sys.path에 추가하고, 한 단계 위(프로젝트 루트)도 추가하여 utils 등 공용 모듈 참조
current_dir = Path(__file__).parent
sys.path.insert(0, str(current_dir))
sys.path.insert(1, str(current_dir.parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
from utils import getVarVal


class SlideType(Enum):
    """슬라이드 타입"""
    TITLE = "title"
    CONTENT = "content"
    SECTION = "section"
    IMAGE = "image"
    TABLE = "table"
    TWO_COLUMN = "two_column"
    BLANK = "blank"


class AspectRatio(Enum):
    """화면 비율"""
    RATIO_4_3 = "4_3"  # 10" × 7.5"
    RATIO_16_9 = "16_9"  # 13.333" × 7.5"


@dataclass
class SlideContent:
    """슬라이드 콘텐츠 데이터 클래스"""
    title: str
    content: Optional[str] = None
    date: Optional[str] = None
    image_path: Optional[str] = None
    image_path_2: Optional[str] = None  # 두 번째 이미지 (two_images 레이아웃용)
    table_data: Optional[List[List[str]]] = None
    table_headers: Optional[List[str]] = None
    table_data_2: Optional[List[List[str]]] = None  # 두 번째 테이블 데이터
    table_headers_2: Optional[List[str]] = None  # 두 번째 테이블 헤더
    subtitle: Optional[str] = None
    layout: Optional[str] = None  # AI가 선택하거나 사용자가 지정
    order: Optional[int] = None  # 사용자 정의 슬라이드 순서


class PPTXGenerator:
    """
    AIRUN PPTX 문서 생성 클래스

    깔끔한 템플릿을 사용하여 PPTX 문서를 생성합니다.
    AI가 자동으로 레이아웃을 선택하거나, 파라미터로 직접 지정할 수 있습니다.
    """

    # 레이아웃 매핑 (슬라이드 타입 -> 템플릿 인덱스)
    # 통합 템플릿: 0(title1), 1(title2), 2(session), 3~15(content)
    LAYOUT_MAPPING = {
        SlideType.TITLE: [0, 1],  # 제목 템플릿 2개
        SlideType.CONTENT: [3],  # 기본 본문
        SlideType.SECTION: [2],  # 세션 타이틀
        SlideType.IMAGE: [12],  # 제목 + 본문 + 이미지 (image_left)
        SlideType.TABLE: [6],  # 제목 + 표 (large_table)
        SlideType.TWO_COLUMN: [4],  # 제목 + 두 칼럼
        SlideType.BLANK: [3],  # 기본 본문
    }

    def __init__(self, aspect_ratio: Union[str, AspectRatio] = AspectRatio.RATIO_4_3,
                 template_type: str = 'default'):
        """
        PPTX 생성기 초기화

        Args:
            aspect_ratio: 화면 비율 ('4_3', '16_9' 또는 AspectRatio enum)
            template_type: 템플릿 타입 ('default', 'custom', 'auto-dev')
        """
        # 프로젝트 루트 경로 (설정 파일에서 가져옴)
        workspace_root = Path(__file__).resolve().parents[1]
        airun_path_str = getVarVal('AIRUN_PATH', str(workspace_root))
        airun_path = Path(airun_path_str)

        # 현재 워크스페이스 템플릿이 있으면 우선 사용
        preferred_template_dir = workspace_root / "templates" / "pptx" / template_type
        preferred_template_path = preferred_template_dir / "templates_16_9.pptx"
        if preferred_template_path.exists():
            airun_path = workspace_root

        self.airun_path = airun_path
        # 템플릿 타입
        self.template_type = template_type
        # 템플릿 디렉토리
        if template_type == "auto-dev":
            self.template_dir = airun_path / "templates" / "pptx" / "custom"
            self.template_path_override = self.template_dir / "templates_16_9.pptx"
        else:
            self.template_dir = airun_path / "templates" / "pptx" / template_type
            self.template_path_override = None
        # 화면 비율 설정
        if isinstance(aspect_ratio, str):
            self.aspect_ratio = AspectRatio(aspect_ratio)
        else:
            self.aspect_ratio = aspect_ratio

        # 페이지 크기 설정
        if self.aspect_ratio == AspectRatio.RATIO_16_9:
            self.slide_width = 13.333  # 인치
        else:
            self.slide_width = 10  # 인치
        self.slide_height = 7.5  # 인치

        # 템플릿 로드
        self._load_templates()

        # 슬라이드 카운터
        self.slide_count = 0

    def _load_templates(self):
        """템플릿 파일 로드"""
        ratio_str = self.aspect_ratio.value
        # 모든 템플릿 타입이 동일한 파일명 사용
        template_path = self.template_path_override or (
            self.template_dir / f"templates_{ratio_str}.pptx"
        )

        if not template_path.exists():
            # 템플릿이 없으면 생성
            self._create_templates()

        # 작업용 임시 파일 생성 (하나의 템플릿만 사용)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            self.template_temp_path = tmp.name
            shutil.copy(str(template_path), self.template_temp_path)

        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            self.working_temp_path = tmp.name
            shutil.copy(str(template_path), self.working_temp_path)

        # Presentation 객체 생성
        self.template_source = Presentation(self.template_temp_path)  # 복사용 (원본 유지)
        self.presentation = Presentation(self.working_temp_path)  # 작업용

        # 작업용 프레젠테이션의 모든 슬라이드 삭제 (빈 문서로 시작)
        self._clear_presentation(self.presentation)

    def _create_templates(self):
        """템플릿 파일 생성"""
        from scripts.create_pptx_templates import create_templates

        print(f"[INFO] 템플릿을 생성합니다...")
        create_templates()

    def _clear_presentation(self, prs):
        """프레젠테이션의 모든 슬라이드 삭제"""
        # 역순으로 삭제하여 인덱스 문제 방지
        for sldId in list(prs.slides._sldIdLst):
            rId = sldId.rId  # relationship ID 문자열 가져오기
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(sldId)

    def _cleanup_temp_files(self):
        """임시 파일 삭제"""
        try:
            if hasattr(self, 'template_temp_path') and os.path.exists(self.template_temp_path):
                os.unlink(self.template_temp_path)
            if hasattr(self, 'working_temp_path') and os.path.exists(self.working_temp_path):
                os.unlink(self.working_temp_path)
        except Exception as e:
            print(f"Warning: Could not cleanup temp files: {e}")

    def __del__(self):
        """소멸자: 임시 파일 정리"""
        self._cleanup_temp_files()

    def _save_frame_format(self, text_frame):
        """text_frame의 첫 번째 단락 서식 저장 (단락 수준 defRPr 포함)"""
        if not text_frame.paragraphs:
            return None

        paragraph = text_frame.paragraphs[0]

        # 기본값 초기화
        format_info = {
            'name': None,
            'size': None,
            'bold': None,
            'italic': None,
            'underline': None,
            'color': None
        }

        # 1. 단락 수준의 defRPr (default run properties) 확인
        try:
            pPr_element = paragraph._element.get_or_add_pPr()
            def_rpr = pPr_element.defRPr

            if def_rpr is not None:
                # sz: font size (in 100ths of a point, e.g., 3600 = 36pt)
                if def_rpr.sz is not None:
                    size_value = int(def_rpr.sz)
                    if size_value < 100:
                        size_value = 100  # 최소 1pt
                    elif size_value > 400000:
                        size_value = 400000  # 최대 4000pt
                    format_info['size'] = size_value

                # b: bold (1 or 0)
                if def_rpr.b is not None:
                    format_info['bold'] = def_rpr.b == "1" or def_rpr.b is True

                # i: italic
                if def_rpr.i is not None:
                    format_info['italic'] = def_rpr.i == "1" or def_rpr.i is True

                # u: underline
                if def_rpr.u is not None:
                    format_info['underline'] = def_rpr.u == "1" or def_rpr.u is True

                # latin: font name
                if def_rpr.latin is not None and def_rpr.latin.typeface is not None:
                    format_info['name'] = def_rpr.latin.typeface

                # 색상 정보 (solidFill > srgbClr)
                if def_rpr.solidFill is not None:
                    srgb_clr = def_rpr.solidFill.srgbClr
                    if srgb_clr is not None and srgb_clr.val is not None:
                        hex_color = srgb_clr.val
                        if len(hex_color) == 6:
                            r = int(hex_color[0:2], 16)
                            g = int(hex_color[2:4], 16)
                            b = int(hex_color[4:6], 16)
                            format_info['color'] = (r, g, b)
        except Exception:
            pass

        # 2. 단락 수준에 서식이 없으면 첫 번째 run에서 확인
        if format_info['size'] is None and paragraph.runs:
            first_run = paragraph.runs[0]
            font = first_run.font

            if font.name:
                format_info['name'] = font.name
            if font.size:
                size_value = font.size
                if hasattr(size_value, 'pt'):
                    size_value = int(size_value.pt * 100)
                elif isinstance(size_value, (int, float)):
                    size_value = int(size_value)

                if size_value < 100:
                    size_value = 100
                elif size_value > 400000:
                    size_value = 400000
                format_info['size'] = size_value

            if font.bold is not None:
                format_info['bold'] = font.bold
            if font.italic is not None:
                format_info['italic'] = font.italic
            if font.underline is not None:
                format_info['underline'] = font.underline

            if font.color:
                try:
                    format_info['color'] = font.color.rgb
                except AttributeError:
                    pass

        return format_info

    def _apply_frame_format(self, text_frame, format_info):
        """저장된 서식을 text_frame에 적용"""
        if format_info is None:
            return

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                if format_info['name']:
                    font.name = format_info['name']
                if format_info['size']:
                    size = format_info['size']
                    size_in_hundredths = None

                    if isinstance(size, str):
                        try:
                            size_in_hundredths = int(size)
                        except (ValueError, TypeError):
                            pass
                    elif hasattr(size, 'pt'):
                        size_in_hundredths = int(size.pt * 100)
                    elif isinstance(size, (int, float)):
                        size_in_hundredths = int(size * 100)

                    if size_in_hundredths is not None:
                        if size_in_hundredths < 100:
                            size_in_hundredths = 100
                        elif size_in_hundredths > 400000:
                            size_in_hundredths = 400000
                        font.size = size_in_hundredths

                if format_info['bold'] is not None:
                    font.bold = format_info['bold']
                if format_info['italic'] is not None:
                    font.italic = format_info['italic']
                if format_info['underline'] is not None:
                    font.underline = format_info['underline']
                if format_info['color']:
                    from pptx.dml.color import RGBColor
                    r, g, b = format_info['color']
                    font.color.rgb = RGBColor(r, g, b)

    def _set_frame_text(self, text_frame, text, preserve_format=True):
        """text_frame의 텍스트 설정 - 서식 유지"""
        if not text_frame.paragraphs:
            text_frame.add_paragraph()
            if text:
                text_frame.paragraphs[0].text = text
            return

        # 기존 서식 저장
        format_info = self._save_frame_format(text_frame) if preserve_format else None

        # 첫 번째 단락 유지하면서 텍스트만 교체
        first_paragraph = text_frame.paragraphs[0]

        # XML 레벨에서 직접 텍스트 설정
        p_element = first_paragraph._element
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        # 첫 번째 run 찾기
        r_elements = p_element.findall('a:r', namespaces=ns)
        if r_elements:
            first_r = r_elements[0]
            t_elements = first_r.findall('a:t', namespaces=ns)
            if t_elements:
                t_elements[0].text = text
            else:
                from lxml import etree
                t = etree.SubElement(first_r, '{%s}t' % ns['a'])
                t.text = text

            # 나머지 run 삭제
            for r in r_elements[1:]:
                p_element.remove(r)
        else:
            first_paragraph.text = text

        # 서식 유지 모드인 경우 run 수준의 rPr 제거
        if preserve_format:
            for rpr in p_element.findall('a:r/a:rPr', namespaces=ns):
                r_element = rpr.getparent()
                r_element.remove(rpr)

        # 서식 복원
        if preserve_format and format_info:
            self._apply_frame_format(text_frame, format_info)

        # 나머지 단락들 삭제
        txBody = text_frame._txBody
        paragraphs_to_remove = list(text_frame.paragraphs[1:])
        for paragraph in paragraphs_to_remove:
            p_element = paragraph._element
            txBody.remove(p_element)

    def replace_text_in_shape(self, shape, text_map: Dict[str, str]):
        """shape의 텍스트를 text_map에 따라 교체 - 부분 일치 지원"""
        if not shape.has_text_frame:
            return

        text_frame = shape.text_frame
        full_text = text_frame.text

        # 더 긴 key를 우선하기 위해 정렬
        sorted_keys = sorted(text_map.keys(), key=len, reverse=True)

        # 1. 정확한 일치 우선
        for old_text in sorted_keys:
            new_val = text_map[old_text]

            if full_text == old_text:
                self._set_frame_text(text_frame, new_val)
                return

        # 2. 부분 일치
        for old_text in sorted_keys:
            new_val = text_map[old_text]

            if full_text.startswith(old_text) and new_val == '':
                self._set_frame_text(text_frame, full_text[len(old_text):])
                return
            elif old_text in full_text:
                new_full_text = full_text.replace(old_text, new_val)
                self._set_frame_text(text_frame, new_full_text)
                return

    def style_table(self, table, headers: List[str], data: List[List[str]]):
        """표 스타일 적용 (헤더는 Executive Summary 스타일로 적용)"""
        from lxml import etree

        # XML 네임스페이스
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        # 테두리 색상과 스타일
        border_color = '000000'  # 검은색 테두리
        border_width = 12700  # 테두리 두께 (1pt)
        # Executive Summary의 헤더 톤에 맞춰 블루 배경과 화이트 텍스트 적용
        header_bg_color = '0E63B5'  # BLUE_BAR (0x0E, 0x63, 0xB5)

        def set_cell_border(cell_element, border_color=border_color, border_width=border_width):
            """셀에 테두리 설정"""
            # tcPr 요소 찾기 또는 생성
            tc_pr = cell_element.find('./a:tcPr', namespaces=ns)
            if tc_pr is None:
                tc_pr = etree.SubElement(cell_element, '{%s}tcPr' % ns['a'])

            # tcBorders 요소 찾기 또는 생성
            tc_borders = tc_pr.find('./a:tcBorders', namespaces=ns)
            if tc_borders is None:
                tc_borders = etree.SubElement(tc_pr, '{%s}tcBorders' % ns['a'])

            # 모든 테두리 설정 (상, 하, 좌, 우)
            for side in ['top', 'left', 'bottom', 'right']:
                border = tc_borders.find('./a:%s' % side, namespaces=ns)
                if border is None:
                    border = etree.SubElement(tc_borders, '{%s}%s' % (ns['a'], side))
                border.set('val', 'single')
                border.set('w', str(border_width))
                border.set('color', border_color)

        def set_cell_fill(cell_element, fill_color):
            """셀 배경색 설정"""
            # tcPr 요소 찾기 또는 생성
            tc_pr = cell_element.find('./a:tcPr', namespaces=ns)
            if tc_pr is None:
                tc_pr = etree.SubElement(cell_element, '{%s}tcPr' % ns['a'])

            # solidFill 요소 생성
            solid_fill = etree.SubElement(tc_pr, '{%s}solidFill' % ns['a'])
            srgb_clr = etree.SubElement(solid_fill, '{%s}srgbClr' % ns['a'])
            srgb_clr.set('val', fill_color)

        is_route_table = any("REC Tariff" in h for h in headers) and any("PJT Name" in h for h in headers)
        is_permits_table = any("Hanjeon PPA" in h for h in headers) and any("EBL" in h for h in headers)
        normalized_headers = [str(h or "").replace("\n", " ").strip().lower() for h in headers]
        is_main_agreements_table = normalized_headers == ["agreement", "completion", "current status", "next steps"]
        is_saved_gantt_table = normalized_headers == ["task", "start w", "end w", "duration", "timeline"]
        is_construction_plan_table = normalized_headers == [
            "spv",
            "phase #",
            "no.",
            "pjt name",
            "capacity (kwp)",
            "rtb date",
            "expected soc",
            "expected cod",
        ]

        if is_route_table and len(headers) == 10:
            total_width = sum(col.width for col in table.columns)
            # Status 컬럼의 서술형 텍스트 공간 확보를 위해,
            # SPV~Offtaker 컬럼은 최소 가독 폭으로 조정하고 Status에 폭을 크게 배정
            ratios = [0.05, 0.05, 0.04, 0.16, 0.09, 0.08, 0.10, 0.08, 0.08, 0.27]
            for idx, ratio in enumerate(ratios):
                table.columns[idx].width = int(total_width * ratio)

            # 한 페이지 내 가시성을 위해 행 높이를 표 높이 안에서 재분배
            total_height = sum(int(row.height) for row in table.rows)
            if total_height > 0 and len(table.rows) > 1:
                header_height = min(int(0.42 * 914400), int(total_height * 0.14))
                body_height = int((total_height - header_height) / (len(table.rows) - 1))
                if body_height > 0:
                    table.rows[0].height = header_height
                    for ridx in range(1, len(table.rows)):
                        table.rows[ridx].height = body_height

        if is_permits_table and len(headers) == 12:
            total_width = sum(col.width for col in table.columns)
            ratios = [0.06, 0.05, 0.06, 0.15, 0.08, 0.09, 0.09, 0.09, 0.09, 0.09, 0.09, 0.16]
            ratio_sum = sum(ratios)
            if ratio_sum <= 0:
                ratio_sum = 1.0
            normalized = [ratio / ratio_sum for ratio in ratios]
            for idx, ratio in enumerate(normalized):
                table.columns[idx].width = int(total_width * ratio)

            total_height = sum(int(row.height) for row in table.rows)
            if total_height > 0 and len(table.rows) > 1:
                header_height = min(int(0.5 * 914400), int(total_height * 0.16))
                body_height = int((total_height - header_height) / (len(table.rows) - 1))
                if body_height > 0:
                    table.rows[0].height = header_height
                    for ridx in range(1, len(table.rows)):
                        table.rows[ridx].height = body_height

        if is_main_agreements_table and len(headers) == 4:
            total_width = sum(col.width for col in table.columns)
            ratios = [0.16, 0.16, 0.42, 0.26]
            ratio_sum = sum(ratios)
            if ratio_sum <= 0:
                ratio_sum = 1.0
            normalized = [ratio / ratio_sum for ratio in ratios]
            for idx, ratio in enumerate(normalized):
                table.columns[idx].width = int(total_width * ratio)

            total_height = sum(int(row.height) for row in table.rows)
            if total_height > 0 and len(table.rows) == 4:
                header_height = int(total_height * 0.15)
                epc_height = int(total_height * 0.45)
                other_height = int((total_height - header_height - epc_height) / 2)
                table.rows[0].height = header_height
                table.rows[1].height = epc_height
                table.rows[2].height = other_height
                table.rows[3].height = total_height - header_height - epc_height - other_height

        if is_saved_gantt_table and len(headers) == 5:
            total_width = sum(col.width for col in table.columns)
            ratios = [0.26, 0.10, 0.10, 0.10, 0.44]
            ratio_sum = sum(ratios)
            if ratio_sum <= 0:
                ratio_sum = 1.0
            normalized = [ratio / ratio_sum for ratio in ratios]
            for idx, ratio in enumerate(normalized):
                table.columns[idx].width = int(total_width * ratio)

            total_height = sum(int(row.height) for row in table.rows)
            if total_height > 0 and len(table.rows) > 1:
                header_height = min(int(0.42 * 914400), int(total_height * 0.15))
                body_height = int((total_height - header_height) / (len(table.rows) - 1))
                if body_height > 0:
                    table.rows[0].height = header_height
                    for ridx in range(1, len(table.rows)):
                        table.rows[ridx].height = body_height

        if is_construction_plan_table and len(headers) == 8:
            total_width = sum(col.width for col in table.columns)
            ratios = [0.08, 0.07, 0.05, 0.25, 0.11, 0.14, 0.14, 0.16]
            ratio_sum = sum(ratios)
            if ratio_sum <= 0:
                ratio_sum = 1.0
            normalized = [r / ratio_sum for r in ratios]
            for idx, ratio in enumerate(normalized):
                table.columns[idx].width = int(total_width * ratio)

        # 헤더 행 스타일
        header_font_size = Pt(9 if (is_route_table or is_permits_table or is_main_agreements_table or is_saved_gantt_table or is_construction_plan_table) else 13)
        for i, header in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = header

            # 헤더 폰트 설정
            text_frame = cell.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            if text_frame.paragraphs:
                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.line_spacing = 1.0
                    if paragraph.runs:
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.size = header_font_size
                            run.font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        run = paragraph.add_run()
                        run.text = header
                        run.font.bold = True
                        run.font.size = header_font_size
                        run.font.color.rgb = RGBColor(255, 255, 255)

            # 헤더 셀에 배경색과 테두리 설정
            # python-pptx에서 셀 엘리먼트 접근
            cell_element = cell._tc  # 또는 cell.element
            set_cell_fill(cell_element, header_bg_color)
            set_cell_border(cell_element)

        # 데이터 행 스타일
        for row_idx, row_data in enumerate(data):
            is_total_row = any(
                isinstance(value, str) and value.strip().upper() == "TOTAL" for value in row_data
            )
            is_subtotal_row = any(
                isinstance(value, str) and value.strip().lower() == "sub total" for value in row_data
            )
            is_permits_project_row = (
                len(row_data) > 2 and str(row_data[2]).strip().isdigit()
            )
            for col_idx, value in enumerate(row_data):
                cell = table.rows[row_idx + 1].cells[col_idx]
                cell.text = str(value)

                # 데이터 폰트 설정
                text_frame = cell.text_frame
                if text_frame.paragraphs:
                    paragraph = text_frame.paragraphs[0]
                    if paragraph.runs:
                        run = paragraph.runs[0]
                        run.font.size = Pt(9 if (is_route_table or is_permits_table or is_main_agreements_table or is_saved_gantt_table or is_construction_plan_table) else 10)
                    else:
                        # run이 없는 경우, paragraph.add_run() 사용
                        run = paragraph.add_run()
                        run.text = str(value)
                        run.font.size = Pt(9 if (is_route_table or is_permits_table or is_main_agreements_table or is_saved_gantt_table or is_construction_plan_table) else 10)

                # 데이터 셀에 테두리 설정
                cell_element = cell._tc  # 또는 cell.element
                set_cell_border(cell_element)

                if is_route_table:
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if text_frame.paragraphs:
                        paragraph = text_frame.paragraphs[0]
                        if col_idx in {0, 1, 2, 5, 6, 7, 8}:
                            paragraph.alignment = PP_ALIGN.CENTER
                        elif col_idx in {4, 9}:
                            paragraph.alignment = PP_ALIGN.LEFT
                        else:
                            paragraph.alignment = PP_ALIGN.LEFT

                    if col_idx == 0:
                        set_cell_fill(cell_element, '2EDC6F')

                    if is_total_row and col_idx != 0:
                        set_cell_fill(cell_element, 'FFF200')
                        if text_frame.paragraphs:
                            paragraph = text_frame.paragraphs[0]
                            if paragraph.runs:
                                paragraph.runs[0].font.bold = True

                if is_permits_table:
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if text_frame.paragraphs:
                        paragraph = text_frame.paragraphs[0]
                        if col_idx in {0, 1, 2, 4, 5, 6, 7, 8, 9, 10, 11}:
                            paragraph.alignment = PP_ALIGN.CENTER
                        else:
                            paragraph.alignment = PP_ALIGN.LEFT

                    if col_idx == 0:
                        set_cell_fill(cell_element, '3DDC84')
                    elif is_permits_project_row and 2 <= col_idx <= 11:
                        set_cell_fill(cell_element, 'FFFFFF')
                    elif col_idx == 2:
                        set_cell_fill(cell_element, 'FFFFFF')
                    else:
                        set_cell_fill(cell_element, 'D9D9D9')

                    if is_total_row:
                        if col_idx == 0:
                            set_cell_fill(cell_element, 'D9D9D9')
                        else:
                            set_cell_fill(cell_element, 'FFF200')
                        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                            text_frame.paragraphs[0].runs[0].font.bold = True

                    if is_subtotal_row:
                        set_cell_fill(cell_element, 'D9D9D9')
                        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                            text_frame.paragraphs[0].runs[0].font.bold = True

                    if is_subtotal_row:
                        set_cell_fill(cell_element, 'D9D9D9')
                        if text_frame.paragraphs:
                            paragraph = text_frame.paragraphs[0]
                            if paragraph.runs:
                                paragraph.runs[0].font.bold = True

                if is_main_agreements_table:
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    text_frame.word_wrap = True
                    if text_frame.paragraphs:
                        paragraph = text_frame.paragraphs[0]
                        paragraph.line_spacing = 1.0
                        if col_idx in {0, 1}:
                            paragraph.alignment = PP_ALIGN.CENTER
                        else:
                            paragraph.alignment = PP_ALIGN.LEFT

                if is_saved_gantt_table:
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    text_frame.word_wrap = False
                    if text_frame.paragraphs:
                        paragraph = text_frame.paragraphs[0]
                        paragraph.line_spacing = 1.0
                        paragraph.alignment = PP_ALIGN.LEFT if col_idx in {0, 4} else PP_ALIGN.CENTER
                        if paragraph.runs and col_idx == 4:
                            paragraph.runs[0].font.name = 'Courier New'

                if is_construction_plan_table:
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    text_frame.word_wrap = True
                    if text_frame.paragraphs:
                        paragraph = text_frame.paragraphs[0]
                        paragraph.line_spacing = 1.0
                        paragraph.alignment = PP_ALIGN.LEFT if col_idx in {0, 3} else PP_ALIGN.CENTER

                    label_text = (cell.text or "").strip().lower().replace(" ", "")
                    if label_text in {"subtotal", "total"}:
                        set_cell_fill(cell_element, 'FFF200' if label_text == 'total' else 'D9D9D9')

        if is_route_table and len(table.rows) > 2:
            subtotal_row_indices: list[int] = []
            for ridx in range(1, len(table.rows)):
                try:
                    label = (table.cell(ridx, 3).text or "").strip().lower()
                    normalized_label = label.replace(" ", "")
                    if normalized_label == "subtotal":
                        subtotal_row_indices.append(ridx)
                except Exception:
                    continue

            merge_end_row = len(table.rows) - 1
            try:
                last_label = (table.cell(merge_end_row, 3).text or "").strip().upper()
                if last_label == "TOTAL" and merge_end_row > 1:
                    merge_end_row -= 1
            except Exception:
                pass

            spv_cell = table.cell(1, 0)
            if merge_end_row >= 1:
                spv_cell.merge(table.cell(merge_end_row, 0))
                spv_cell.text = spv_cell.text or "BOLT#2"
                set_cell_fill(spv_cell._tc, '2EDC6F')
                spv_frame = spv_cell.text_frame
                spv_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if spv_frame.paragraphs:
                    spv_paragraph = spv_frame.paragraphs[0]
                    spv_paragraph.alignment = PP_ALIGN.CENTER
                    if spv_paragraph.runs:
                        spv_paragraph.runs[0].font.bold = True
                        spv_paragraph.runs[0].font.size = Pt(9)

            # Phase # 세로 병합 (1~7 + SubTotal, 8~N + SubTotal)
            phase_col = 1
            if subtotal_row_indices:
                p1_end = subtotal_row_indices[0]
                if p1_end >= 1:
                    p1_cell = table.cell(1, phase_col)
                    p1_cell.merge(table.cell(p1_end, phase_col))
                    p1_cell.text = "1"
                    p1_frame = p1_cell.text_frame
                    p1_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if p1_frame.paragraphs:
                        p1_paragraph = p1_frame.paragraphs[0]
                        p1_paragraph.alignment = PP_ALIGN.CENTER
                        if p1_paragraph.runs:
                            p1_paragraph.runs[0].font.bold = True
                            p1_paragraph.runs[0].font.size = Pt(9)

                if len(subtotal_row_indices) >= 2:
                    p2_start = subtotal_row_indices[0] + 1
                    p2_end = subtotal_row_indices[1]
                else:
                    p2_start = subtotal_row_indices[0] + 1
                    p2_end = merge_end_row

                if p2_end >= p2_start:
                    p2_cell = table.cell(p2_start, phase_col)
                    p2_cell.merge(table.cell(p2_end, phase_col))
                    p2_cell.text = "2"
                    p2_frame = p2_cell.text_frame
                    p2_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if p2_frame.paragraphs:
                        p2_paragraph = p2_frame.paragraphs[0]
                        p2_paragraph.alignment = PP_ALIGN.CENTER
                        if p2_paragraph.runs:
                            p2_paragraph.runs[0].font.bold = True
                            p2_paragraph.runs[0].font.size = Pt(9)

            # REC Tariff 컬럼 세로 병합
            rec_tariff_idx = None
            offtaker_idx = None
            status_idx = None
            for idx, header in enumerate(headers):
                normalized = header.replace("\n", " ")
                if "REC Tariff" in normalized:
                    rec_tariff_idx = idx
                if "Offtaker" in normalized:
                    offtaker_idx = idx
                if "Status" in normalized:
                    status_idx = idx

            for merge_col in [rec_tariff_idx]:
                if merge_col is None:
                    continue
                if merge_end_row <= 1:
                    continue
                top_cell = table.cell(1, merge_col)
                top_cell.merge(table.cell(merge_end_row, merge_col))
                merged_lines = [ln.strip() for ln in (top_cell.text or "").splitlines() if ln.strip()]
                if merged_lines:
                    top_cell.text = merged_lines[0]
                top_frame = top_cell.text_frame
                top_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if top_frame.paragraphs:
                    top_paragraph = top_frame.paragraphs[0]
                    top_paragraph.alignment = PP_ALIGN.CENTER
                    if top_paragraph.runs:
                        top_paragraph.runs[0].font.bold = False
                        top_paragraph.runs[0].font.size = Pt(9)

            # Offtaker 컬럼 세로 병합 + 고정 문구 입력
            if offtaker_idx is not None and merge_end_row > 1:
                offtaker_cell = table.cell(1, offtaker_idx)
                offtaker_cell.merge(table.cell(merge_end_row, offtaker_idx))
                offtaker_cell.text = "SBP\n(Samcheok\nBluepower)"
                offtaker_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                offtaker_frame = offtaker_cell.text_frame
                offtaker_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                offtaker_frame.word_wrap = True
                if offtaker_frame.paragraphs:
                    for para in offtaker_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        para.space_before = Pt(0)
                        para.space_after = Pt(0)
                        if para.runs:
                            for run in para.runs:
                                run.font.bold = True
                                run.font.size = Pt(9)

            # Status 컬럼 세로 병합
            if status_idx is not None and merge_end_row > 1:
                status_lines: list[str] = []
                for ridx in range(1, merge_end_row + 1):
                    try:
                        txt = (table.cell(ridx, status_idx).text or "").strip()
                        if txt and txt.lower() != "subtotal":
                            status_lines.append(txt)
                    except Exception:
                        continue

                merged_status = "\n".join(dict.fromkeys(status_lines))
                status_cell = table.cell(1, status_idx)
                status_cell.merge(table.cell(merge_end_row, status_idx))
                status_cell.text = merged_status
                status_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                status_frame = status_cell.text_frame
                status_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                status_frame.word_wrap = True
                if status_frame.paragraphs:
                    for para in status_frame.paragraphs:
                        para.alignment = PP_ALIGN.LEFT
                        para.space_before = Pt(0)
                        para.space_after = Pt(0)
                        if para.runs:
                            for run in para.runs:
                                run.font.size = Pt(9)

            for subtotal_row_idx in subtotal_row_indices:
                subtotal_cell = table.cell(subtotal_row_idx, 3)
                subtotal_cell.merge(table.cell(subtotal_row_idx, 4))
                subtotal_cell.text = "SubTotal"
                subtotal_frame = subtotal_cell.text_frame
                subtotal_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if subtotal_frame.paragraphs:
                    subtotal_paragraph = subtotal_frame.paragraphs[0]
                    subtotal_paragraph.alignment = PP_ALIGN.CENTER
                    if subtotal_paragraph.runs:
                        subtotal_paragraph.runs[0].font.bold = True
                        subtotal_paragraph.runs[0].font.size = Pt(9)

        if is_permits_table and len(table.rows) > 2:
            subtotal_row_indices: list[int] = []
            for ridx in range(1, len(table.rows)):
                label = (table.cell(ridx, 3).text or "").strip().lower().replace(" ", "")
                if label == "subtotal":
                    subtotal_row_indices.append(ridx)

            total_row_idx = len(table.rows) - 1
            if (table.cell(total_row_idx, 3).text or "").strip().upper() != "TOTAL":
                total_row_idx = -1

            spv_end = total_row_idx if total_row_idx >= 1 else len(table.rows) - 1
            spv_cell = table.cell(1, 0)
            spv_cell.merge(table.cell(spv_end, 0))
            spv_cell.text = spv_cell.text or "BOLT#2"
            set_cell_fill(spv_cell._tc, '3DDC84')
            spv_frame = spv_cell.text_frame
            spv_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            if spv_frame.paragraphs:
                p = spv_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    p.runs[0].font.bold = True
                    p.runs[0].font.size = Pt(9)

            if subtotal_row_indices:
                p1_end = subtotal_row_indices[0]
                p1_cell = table.cell(1, 1)
                p1_cell.merge(table.cell(p1_end, 1))
                p1_cell.text = "1"
                p1_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if p1_cell.text_frame.paragraphs:
                    p = p1_cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    if p.runs:
                        p.runs[0].font.bold = True
                        p.runs[0].font.size = Pt(9)

                p2_start = p1_end + 1
                p2_end = subtotal_row_indices[1] if len(subtotal_row_indices) >= 2 else (total_row_idx - 1 if total_row_idx >= 0 else len(table.rows) - 1)
                if p2_end >= p2_start:
                    p2_cell = table.cell(p2_start, 1)
                    p2_cell.merge(table.cell(p2_end, 1))
                    p2_cell.text = "2"
                    p2_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if p2_cell.text_frame.paragraphs:
                        p = p2_cell.text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        if p.runs:
                            p.runs[0].font.bold = True
                            p.runs[0].font.size = Pt(9)

            for subtotal_row_idx in subtotal_row_indices:
                subtotal_cell = table.cell(subtotal_row_idx, 2)
                subtotal_cell.merge(table.cell(subtotal_row_idx, 4))
                subtotal_cell.text = "Sub Total"
                subtotal_frame = subtotal_cell.text_frame
                subtotal_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if subtotal_frame.paragraphs:
                    p = subtotal_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    if p.runs:
                        p.runs[0].font.bold = True
                        p.runs[0].font.size = Pt(9)

            if total_row_idx >= 1:
                total_cell = table.cell(total_row_idx, 2)
                total_cell.merge(table.cell(total_row_idx, 4))
                total_cell.text = "TOTAL"
                total_frame = total_cell.text_frame
                total_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if total_frame.paragraphs:
                    p = total_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    if p.runs:
                        p.runs[0].font.bold = True
                        p.runs[0].font.size = Pt(9)

            # EBL/EIA/DAP/Hanjeon 영역 병합 (이미지 스타일)
            data_rows = [
                ridx for ridx in range(1, len(table.rows))
                if (table.cell(ridx, 2).text or "").strip().isdigit()
            ]

            def _normalize_merge_text(text: str) -> str:
                return " ".join((text or "").replace("\n", " ").split()).strip().lower()

            def _apply_permits_cell_text_style(cell):
                tf = cell.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.word_wrap = True
                normalized_text = _normalize_merge_text(cell.text or "")
                for paragraph in tf.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.line_spacing = 1.0
                    for run in paragraph.runs:
                        run.font.size = Pt(9)
                        run.font.bold = "issued" in normalized_text

            def _merge_and_keep_top_text(start_row: int, end_row: int, col_idx: int):
                top_cell = table.cell(start_row, col_idx)
                top_text = top_cell.text or ""
                if end_row > start_row:
                    top_cell.merge(table.cell(end_row, col_idx))
                    top_cell = table.cell(start_row, col_idx)
                top_cell.text = top_text
                _apply_permits_cell_text_style(top_cell)

            def _merge_contiguous_same(col_idx: int):
                if not data_rows:
                    return
                start = None
                prev = None
                for ridx in data_rows:
                    cell_obj = table.cell(ridx, col_idx)
                    if getattr(cell_obj, "is_spanned", False):
                        continue
                    txt_raw = (cell_obj.text or "")
                    txt = _normalize_merge_text(txt_raw)
                    if not txt:
                        if start is not None and prev is not None and prev > start:
                            _merge_and_keep_top_text(start, prev, col_idx)
                        start = None
                        prev = None
                        continue

                    if start is None:
                        start = ridx
                        prev = ridx
                    elif txt == _normalize_merge_text(table.cell(prev, col_idx).text or "") and ridx == prev + 1:
                        prev = ridx
                    else:
                        if prev is not None and prev > start:
                            _merge_and_keep_top_text(start, prev, col_idx)
                        start = ridx
                        prev = ridx

                if start is not None and prev is not None and prev > start:
                    _merge_and_keep_top_text(start, prev, col_idx)

            def _merge_contiguous_by_predicate(col_idx: int, predicate):
                if not data_rows:
                    return
                start = None
                prev = None
                for ridx in data_rows:
                    cell_obj = table.cell(ridx, col_idx)
                    if getattr(cell_obj, "is_spanned", False):
                        continue
                    txt = _normalize_merge_text(cell_obj.text or "")
                    matched = bool(txt) and predicate(txt)
                    if not matched:
                        if start is not None and prev is not None and prev > start:
                            _merge_and_keep_top_text(start, prev, col_idx)
                        start = None
                        prev = None
                        continue

                    if start is None:
                        start = ridx
                        prev = ridx
                    elif ridx == prev + 1:
                        prev = ridx
                    else:
                        if prev is not None and prev > start:
                            _merge_and_keep_top_text(start, prev, col_idx)
                        start = ridx
                        prev = ridx

                if start is not None and prev is not None and prev > start:
                    _merge_and_keep_top_text(start, prev, col_idx)

            def _merge_by_project_no_range(col_idx: int, start_no: int, end_no: int, predicate) -> bool:
                matching_rows: list[int] = []
                for ridx in data_rows:
                    try:
                        no_text = (table.cell(ridx, 2).text or "").strip()
                        project_no = int(no_text)
                    except (TypeError, ValueError):
                        continue
                    if project_no < start_no or project_no > end_no:
                        continue
                    cell_obj = table.cell(ridx, col_idx)
                    if getattr(cell_obj, "is_spanned", False):
                        continue
                    normalized = _normalize_merge_text(cell_obj.text or "")
                    if not normalized or not predicate(normalized):
                        continue
                    matching_rows.append(ridx)

                if not matching_rows:
                    return False

                block_start = matching_rows[0]
                block_prev = matching_rows[0]
                merged_any = False
                for ridx in matching_rows[1:]:
                    if ridx == block_prev + 1:
                        block_prev = ridx
                    else:
                        if block_prev > block_start:
                            _merge_and_keep_top_text(block_start, block_prev, col_idx)
                            merged_any = True
                        block_start = ridx
                        block_prev = ridx
                if block_prev > block_start:
                    _merge_and_keep_top_text(block_start, block_prev, col_idx)
                    merged_any = True
                return merged_any

            # EBL Completion: Issued 병합
            _merge_contiguous_by_predicate(5, lambda t: t == "issued")

            # DAP Completion: 요청 스타일 고정 병합(1~2 Not subject, 3~4 Issued, 5~7 Issued)
            dap_custom_applied = False
            dap_custom_applied = _merge_by_project_no_range(9, 1, 2, lambda t: "not subject to dap" in t) or dap_custom_applied
            dap_custom_applied = _merge_by_project_no_range(9, 3, 4, lambda t: t == "issued") or dap_custom_applied
            dap_custom_applied = _merge_by_project_no_range(9, 5, 7, lambda t: t == "issued") or dap_custom_applied
            if not dap_custom_applied:
                _merge_contiguous_by_predicate(9, lambda t: "not subject to dap" in t)
                _merge_contiguous_by_predicate(9, lambda t: t == "issued")

            # DAP Validity, Hanjeon Approval은 동일값 연속 병합
            for merge_col in [10, 11]:
                _merge_contiguous_same(merge_col)

            # EIA는 모든 프로젝트 행이 동일한 값일 때만 병합
            if data_rows:
                eia_vals = [(table.cell(ridx, 7).text or "").strip() for ridx in data_rows]
                eia_non_empty = [v for v in eia_vals if v]
                if (
                    eia_non_empty
                    and len(eia_non_empty) == len(data_rows)
                    and all(v == eia_non_empty[0] for v in eia_non_empty)
                ):
                    first_row = data_rows[0]
                    last_row = data_rows[-1]
                    eia_cell = table.cell(first_row, 7)
                    eia_cell.merge(table.cell(last_row, 8))
                    eia_cell.text = eia_non_empty[0]
                    eia_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if eia_cell.text_frame.paragraphs:
                        p = eia_cell.text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        if p.runs:
                            p.runs[0].font.size = Pt(9)

            # EBL/EIA/DAP/Hanjeon 영역 텍스트 정렬/강조
            for ridx in range(1, len(table.rows)):
                for cidx in [5, 6, 7, 8, 9, 10, 11]:
                    cell = table.cell(ridx, cidx)
                    tf = cell.text_frame
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if tf.paragraphs:
                        p = tf.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        if p.runs:
                            text = (cell.text or "").strip().lower()
                            for run in p.runs:
                                run.font.size = Pt(9)
                                run.font.bold = "issued" in text

    def _render_executive_summary_layout(self, slide, prs, content: SlideContent):
        """Executive Summary 전용 레이아웃을 직접 구성해 표 데이터를 채운다."""
        # 템플릿 title 위치/스타일 캡처 (하드코딩 제거)
        title_ref: Dict[str, Any] = {
            "left": Inches(0.55),
            "top": Inches(0.35),
            "width": Inches((prs.slide_width / 914400) - 1.10),
            "height": Inches(0.55),
            "size": None,
            "bold": None,
            "color": None,
            "name": None,
        }
        table_ref: Dict[str, Any] = {
            "left": Inches(0.55),
            "top": Inches(0.95),
            "width": Inches((prs.slide_width / 914400) - 1.10),
            "height": Inches((prs.slide_height / 914400) - 1.50),
            "col_widths": [],
            "row_heights": [],
        }

        def _pick_template_title_shape() -> Any | None:
            candidates = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = (shape.text or "").strip()
                if not text:
                    continue
                top = int(shape.top)
                if top > int(2.2 * 914400):
                    continue
                size_pt = 0.0
                tf = shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs and tf.paragraphs[0].runs[0].font.size:
                    size_pt = float(tf.paragraphs[0].runs[0].font.size.pt)
                candidates.append((size_pt, -top, shape))
            if not candidates:
                return None
            candidates.sort(reverse=True, key=lambda x: (x[0], x[1]))
            return candidates[0][2]

        def _pick_template_table_shape() -> Any | None:
            candidates = []
            for shape in slide.shapes:
                if not getattr(shape, "has_table", False):
                    continue
                area = int(shape.width) * int(shape.height)
                candidates.append((area, shape))
            if not candidates:
                return None
            candidates.sort(reverse=True, key=lambda x: x[0])
            return candidates[0][1]

        tpl_title = _pick_template_title_shape()
        if tpl_title is not None:
            title_ref["left"] = tpl_title.left
            title_ref["top"] = tpl_title.top
            title_ref["width"] = tpl_title.width
            title_ref["height"] = tpl_title.height
            tf = tpl_title.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                run0 = tf.paragraphs[0].runs[0]
                title_ref["size"] = run0.font.size
                title_ref["bold"] = run0.font.bold
                title_ref["name"] = run0.font.name
                try:
                    title_ref["color"] = run0.font.color.rgb if run0.font.color else None
                except Exception:
                    title_ref["color"] = None

        tpl_table = _pick_template_table_shape()
        if tpl_table is not None:
            table_ref["left"] = tpl_table.left
            table_ref["top"] = tpl_table.top
            table_ref["width"] = tpl_table.width
            table_ref["height"] = tpl_table.height
            try:
                table_ref["col_widths"] = [col.width for col in tpl_table.table.columns]
                table_ref["row_heights"] = [row.height for row in tpl_table.table.rows]
            except Exception:
                table_ref["col_widths"] = []
                table_ref["row_heights"] = []

        # 기존 도형 제거
        for shape in list(slide.shapes):
            slide.shapes._spTree.remove(shape._element)

        slide_w_in = prs.slide_width / 914400
        slide_h_in = prs.slide_height / 914400

        # 색상 팔레트
        BLUE_BAR = (0x0E, 0x63, 0xB5)
        TEXT = (0x33, 0x33, 0x33)

        # 제목
        title_box = slide.shapes.add_textbox(
            title_ref["left"], title_ref["top"], title_ref["width"], title_ref["height"]
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = content.title or "Executive Summary"
        p.alignment = PP_ALIGN.LEFT
        r = p.runs[0]
        if title_ref["size"] is not None:
            r.font.size = title_ref["size"]
        if title_ref["bold"] is not None:
            r.font.bold = title_ref["bold"]
        if title_ref["name"]:
            r.font.name = title_ref["name"]
        if title_ref["color"] is not None:
            r.font.color.rgb = title_ref["color"]

        # 데이터 준비
        left_data = content.table_data or []
        right_data = content.table_data_2 or []
        max_rows = max(len(left_data), len(right_data))
        # 템플릿 기준 최소 본문 7행 유지 (헤더 포함 총 8행)
        rows = max(7, max_rows) + 1

        cols = 4  # numL | textL | numR | textR
        tbl_shape = slide.shapes.add_table(
            rows, cols,
            table_ref["left"], table_ref["top"],
            table_ref["width"], table_ref["height"]
        )
        table = tbl_shape.table

        # 열 너비는 템플릿 표 기준 사용
        if len(table_ref["col_widths"]) >= cols:
            for i in range(cols):
                table.columns[i].width = table_ref["col_widths"][i]

        # 행 높이도 템플릿 기준 사용 (행 수가 동일하면 그대로 복제)
        ref_rows = table_ref.get("row_heights") or []
        if rows > 0 and ref_rows:
            if len(ref_rows) == rows:
                for ridx in range(rows):
                    if ref_rows[ridx] is not None:
                        table.rows[ridx].height = ref_rows[ridx]
            else:
                header_h = int(ref_rows[0]) if ref_rows[0] is not None else None
                total_h = int(table_ref["height"])
                if header_h is not None and header_h < total_h:
                    table.rows[0].height = header_h
                    if rows > 1:
                        body_h = int((total_h - header_h) / (rows - 1))
                        for ridx in range(1, rows):
                            table.rows[ridx].height = body_h

        # 헤더 병합
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(0, 2).merge(table.cell(0, 3))

        # 헤더 텍스트
        left_header = ""
        if content.table_headers:
            left_header = content.table_headers[1] if len(content.table_headers) > 1 else content.table_headers[0]
        right_header = ""
        if content.table_headers_2:
            right_header = content.table_headers_2[1] if len(content.table_headers_2) > 1 else content.table_headers_2[0]

        def _style_header(cell, text):
            tf_h = cell.text_frame
            tf_h.clear()
            tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_h = tf_h.paragraphs[0]
            p_h.text = text
            p_h.alignment = PP_ALIGN.CENTER
            p_h.line_spacing = 1.0
            run_h = p_h.runs[0] if p_h.runs else p_h.add_run()
            run_h.font.size = Pt(13)
            run_h.font.bold = True
            run_h.font.color.rgb = RGBColor(255, 255, 255)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*BLUE_BAR)

        _style_header(table.cell(0, 0), left_header)
        _style_header(table.cell(0, 2), right_header)

        # 본문 행 채우기
        def _set_cell(cell, text, bold=False, align=PP_ALIGN.LEFT):
            tf_c = cell.text_frame
            tf_c.clear()
            tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_c = tf_c.paragraphs[0]
            p_c.text = text
            p_c.alignment = align
            p_c.line_spacing = 1.05
            run_c = p_c.runs[0] if p_c.runs else p_c.add_run()
            run_c.font.size = Pt(11)
            run_c.font.bold = bold
            run_c.font.color.rgb = RGBColor(*TEXT)

        for idx in range(max_rows):
            l_num, l_val = ("", "")
            if idx < len(left_data):
                l_num = str(left_data[idx][0]) if len(left_data[idx]) > 0 else ""
                l_val = str(left_data[idx][1]) if len(left_data[idx]) > 1 else ""

            r_num, r_val = ("", "")
            if idx < len(right_data):
                r_num = str(right_data[idx][0]) if len(right_data[idx]) > 0 else ""
                r_val = str(right_data[idx][1]) if len(right_data[idx]) > 1 else ""

            row = idx + 1
            _set_cell(table.cell(row, 0), l_num, align=PP_ALIGN.CENTER)
            _set_cell(table.cell(row, 1), l_val)
            _set_cell(table.cell(row, 2), r_num, align=PP_ALIGN.CENTER)
            _set_cell(table.cell(row, 3), r_val)

    def _render_session_title_layout(self, slide, prs, content: SlideContent):
        """Session title 레이아웃을 입력값 기반 텍스트로 구성한다."""
        for shape in list(slide.shapes):
            if int(getattr(shape, "shape_type", -1)) == 13:
                continue
            slide.shapes._spTree.remove(shape._element)

        title_text = (content.title or "").strip()
        body_text = (content.content or "").strip()
        if not title_text and not body_text:
            return

        slide_w_in = prs.slide_width / 914400
        slide_h_in = prs.slide_height / 914400

        textbox = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(slide_h_in * 0.45),
            Inches(slide_w_in * 0.45),
            Inches(2.2),
        )
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True

        if title_text:
            p1 = tf.paragraphs[0]
            p1.text = title_text
            p1.alignment = PP_ALIGN.LEFT
            p1.line_spacing = 1.1
            r1 = p1.runs[0]
            r1.font.size = Pt(40)
            r1.font.bold = True
            r1.font.color.rgb = RGBColor(255, 255, 255)

        if body_text:
            p2 = tf.add_paragraph() if title_text else tf.paragraphs[0]
            p2.text = body_text
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = 1.1
            r2 = p2.runs[0]
            r2.font.size = Pt(30)
            r2.font.bold = True
            r2.font.color.rgb = RGBColor(255, 255, 255)

    def _render_approval_request_layout(self, slide, content: SlideContent):
        """Approval request 레이아웃의 sample 텍스트를 입력 데이터로 치환한다."""
        title_text = (content.title or "Approval request").strip()
        raw = (content.content or "").strip()

        expected_text = raw
        capex_text = ""
        split_token = "\ncapex :"
        if split_token in raw.lower():
            lower_raw = raw.lower()
            split_idx = lower_raw.find(split_token)
            expected_text = raw[:split_idx].strip()
            capex_text = raw[split_idx + 1 :].strip()

        if expected_text.lower().startswith("expected return:"):
            expected_text = expected_text.split(":", 1)[1].strip()
        if capex_text.lower().startswith("capex"):
            capex_text = capex_text.split(":", 1)[1].strip() if ":" in capex_text else capex_text

        sample_values = [expected_text, capex_text]
        sample_idx = 0

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text or "").strip()

            if text == "Approval request":
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title_text
                p.alignment = PP_ALIGN.LEFT
                if p.runs:
                    p.runs[0].font.size = Pt(28)
                continue

            if text == "sample" and sample_idx < len(sample_values):
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = sample_values[sample_idx]
                p.alignment = PP_ALIGN.LEFT
                if p.runs:
                    p.runs[0].font.size = Pt(13)
                sample_idx += 1

    def _render_title_subtitle_table_layout(self, slide, prs, content: SlideContent):
        """제목 + 부제목 + 테이블 레이아웃을 직접 구성한다."""
        subtitle_text = (content.subtitle or content.content or "").strip()

        for shape in list(slide.shapes):
            slide.shapes._spTree.remove(shape._element)

        slide_w = prs.slide_width
        slide_h = prs.slide_height
        margin = Inches(0.55)

        background = slide.shapes.add_shape(
            1,
            0,
            0,
            slide_w,
            slide_h,
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background.line.fill.background()

        content_w = int(slide_w - (2 * margin))
        title_box = slide.shapes.add_textbox(margin, Inches(0.55), content_w, Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.clear()
        p_title = tf_title.paragraphs[0]
        p_title.text = content.title or "슬라이드 제목"
        p_title.alignment = PP_ALIGN.LEFT
        if p_title.runs:
            r_title = p_title.runs[0]
            r_title.font.size = Pt(30)
            r_title.font.bold = True
            r_title.font.color.rgb = RGBColor(68, 114, 196)

        subtitle_box = slide.shapes.add_textbox(margin, Inches(1.17), content_w, Inches(0.45))
        tf_sub = subtitle_box.text_frame
        tf_sub.clear()
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.alignment = PP_ALIGN.LEFT
        if p_sub.runs:
            r_sub = p_sub.runs[0]
            r_sub.font.size = Pt(14)
            r_sub.font.bold = False
            r_sub.font.color.rgb = RGBColor(90, 90, 90)

        table_top = Inches(1.75)
        table_h = int(slide_h - table_top - margin)
        if content.table_headers and content.table_data:
            rows = len(content.table_data) + 1
            cols = len(content.table_headers)
            table = slide.shapes.add_table(rows, cols, margin, table_top, content_w, table_h).table
            self.style_table(table, content.table_headers, content.table_data)
        else:
            placeholder = slide.shapes.add_shape(5, margin, table_top, content_w, table_h)
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(242, 242, 242)
            placeholder.line.color.rgb = RGBColor(200, 200, 200)
            placeholder.line.width = Pt(2)
            tf = placeholder.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "표영역"
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.size = Pt(24)
                p.runs[0].font.color.rgb = RGBColor(150, 150, 150)

    def _render_gantt_plan_layout(self, slide, prs, content: SlideContent):
        """저장된 phase 간트 데이터를 실제 막대형 차트로 렌더링한다."""
        COLOR_BG = RGBColor(255, 255, 255)
        COLOR_HEADER = RGBColor(49, 53, 124)
        COLOR_GRID = RGBColor(35, 35, 35)
        COLOR_BODY = RGBColor(237, 237, 237)
        COLOR_FILL = RGBColor(255, 255, 0)

        rows = content.table_data or []
        parsed_rows: list[dict[str, Any]] = []
        for row in rows:
            if not isinstance(row, list) or len(row) < 3:
                continue
            raw_task = str(row[0] or "")
            task_text = raw_task.strip()
            if not task_text or task_text == "No saved tasks":
                continue
            start = self._to_int(row[1], default=1)
            end = self._to_int(row[2], default=start)
            if end < start:
                start, end = end, start
            is_group = not raw_task.startswith("  ")
            parsed_rows.append({
                "task": task_text,
                "start": max(1, start),
                "end": max(1, end),
                "is_group": is_group,
            })

        for shape in list(slide.shapes):
            slide.shapes._spTree.remove(shape._element)

        slide_w = prs.slide_width
        slide_h = prs.slide_height
        margin = Inches(0.08)
        title_y = Inches(0.08)
        table_top = Inches(0.64)
        table_h = int(slide_h - table_top - Inches(0.10))
        table_w = int(slide_w - (2 * margin))

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        title_box = slide.shapes.add_textbox(margin, title_y, table_w, Inches(0.45))
        tf_title = title_box.text_frame
        tf_title.clear()
        p_title = tf_title.paragraphs[0]
        p_title.text = content.title or "Project Construction Plan"
        p_title.alignment = PP_ALIGN.LEFT
        if p_title.runs:
            p_title.runs[0].font.size = Pt(17)
            p_title.runs[0].font.bold = True
            p_title.runs[0].font.color.rgb = RGBColor(31, 91, 176)

        if not parsed_rows:
            empty = slide.shapes.add_textbox(margin, table_top, table_w, Inches(0.5))
            tf_empty = empty.text_frame
            tf_empty.clear()
            p_empty = tf_empty.paragraphs[0]
            p_empty.text = "No saved gantt data"
            p_empty.alignment = PP_ALIGN.CENTER
            if p_empty.runs:
                p_empty.runs[0].font.size = Pt(10)
            return

        week_count = max(1, min(52, max((r["end"] for r in parsed_rows), default=1)))
        month_bucket = 4
        month_count = int((week_count + month_bucket - 1) / month_bucket)

        base = datetime.now()
        month_names = [
            "Jan", "Feb", "Mar", "Apr", "May", "Jun",
            "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"
        ]
        month_labels: list[str] = []
        m = base.month
        y = base.year
        for _ in range(month_count):
            month_labels.append(f"{month_names[m - 1]} {y}")
            m += 1
            if m > 12:
                m = 1
                y += 1

        col_sn = int(Inches(0.42))
        col_task = int(Inches(2.2))
        col_desc = int(Inches(4.1))
        timeline_left = margin + col_sn + col_task + col_desc
        timeline_w = int(table_w - col_sn - col_task - col_desc)

        header_month_h = int(Inches(0.26))
        header_week_h = int(Inches(0.20))
        header_total_h = header_month_h + header_week_h
        body_h = int(table_h - header_total_h)
        row_h = max(10, int(body_h / max(1, len(parsed_rows))))
        body_h = row_h * len(parsed_rows)

        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, table_top, table_w, header_total_h + body_h)
        frame.fill.solid()
        frame.fill.fore_color.rgb = COLOR_BODY
        frame.line.color.rgb = COLOR_GRID
        frame.line.width = Pt(1)

        header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, table_top, table_w, header_total_h)
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = COLOR_HEADER
        header_bg.line.fill.background()

        def _header_text(x: int, yv: int, w: int, h: int, text: str, align=PP_ALIGN.CENTER):
            tb = slide.shapes.add_textbox(x, yv, w, h)
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = align
            if p.runs:
                p.runs[0].font.size = Pt(8)
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = RGBColor(255, 255, 255)

        _header_text(margin, table_top, col_sn, header_total_h, "SN")
        _header_text(margin + col_sn, table_top, col_task + col_desc, header_total_h, "Description", PP_ALIGN.LEFT)

        week_w = timeline_w / week_count if week_count > 0 else timeline_w
        for mi in range(month_count):
            start_week = (mi * month_bucket) + 1
            end_week = min((mi + 1) * month_bucket, week_count)
            x = int(timeline_left + ((start_week - 1) * week_w))
            w = int((end_week - start_week + 1) * week_w)
            _header_text(x, table_top, w, header_month_h, month_labels[mi])
            for wk in range(start_week, end_week + 1):
                wx = int(timeline_left + ((wk - 1) * week_w))
                _header_text(wx, table_top + header_month_h, int(week_w), header_week_h, f"{((wk - 1) % 4) + 1}w")

        # Fill active ranges with yellow cells first
        body_top = table_top + header_total_h
        group_spans: list[dict[str, Any]] = []
        row_to_span_idx: dict[int, int] = {}
        cursor = 0
        sn_counter = 0
        while cursor < len(parsed_rows):
            row = parsed_rows[cursor]
            if row["is_group"]:
                start_idx = cursor
                end_idx = cursor
                while end_idx + 1 < len(parsed_rows) and not parsed_rows[end_idx + 1]["is_group"]:
                    end_idx += 1
                sn_counter += 1
                span_idx = len(group_spans)
                group_spans.append({
                    "start": start_idx,
                    "end": end_idx,
                    "sn": sn_counter,
                    "task": row["task"],
                })
                for ridx in range(start_idx, end_idx + 1):
                    row_to_span_idx[ridx] = span_idx
                cursor = end_idx + 1
            else:
                # 상위 그룹이 없는 단독 task는 1행 span으로 처리
                sn_counter += 1
                span_idx = len(group_spans)
                group_spans.append({
                    "start": cursor,
                    "end": cursor,
                    "sn": sn_counter,
                    "task": row["task"],
                })
                row_to_span_idx[cursor] = span_idx
                cursor += 1

        def _cell_text(y_row: int, x: int, w: int, h: int, txt: str, bold: bool = False, align: int = PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(x + 2, y_row + 1, max(8, w - 4), max(8, h - 2))
            tf = tb.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = txt
            p.alignment = align
            if p.runs:
                p.runs[0].font.size = Pt(7)
                p.runs[0].font.bold = bold
                p.runs[0].font.color.rgb = RGBColor(0, 0, 0)

        for ridx, row in enumerate(parsed_rows):
            y_row = int(body_top + (ridx * row_h))
            start = max(1, min(week_count, row["start"]))
            end = max(1, min(week_count, row["end"]))
            for wk in range(start, end + 1):
                x = int(timeline_left + ((wk - 1) * week_w))
                fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_row, int(week_w), row_h)
                fill.fill.solid()
                fill.fill.fore_color.rgb = COLOR_FILL
                fill.line.fill.background()

            span_idx = row_to_span_idx.get(ridx)
            span = group_spans[span_idx] if span_idx is not None else None
            if span and ridx == span["start"]:
                y_span = int(body_top + (span["start"] * row_h))
                h_span = int((span["end"] - span["start"] + 1) * row_h)
                _cell_text(y_span, margin, col_sn, h_span, str(span["sn"]), bold=True, align=PP_ALIGN.CENTER)
                _cell_text(y_span, margin + col_sn, col_task, h_span, str(span["task"]), bold=True)

            if not row["is_group"]:
                _cell_text(y_row, margin + col_sn + col_task, col_desc, row_h, row["task"], bold=False)

        # Draw grid lines on top
        x_edges = [margin, margin + col_sn, margin + col_sn + col_task, margin + col_sn + col_task + col_desc]
        for wk in range(0, week_count + 1):
            x_edges.append(int(timeline_left + (wk * week_w)))
        for x in x_edges:
            v = slide.shapes.add_connector(1, x, table_top, x, body_top + body_h)
            v.line.color.rgb = COLOR_GRID
            v.line.width = Pt(0.8)

        y_edges = [table_top, table_top + header_month_h, body_top]
        for ridx in range(1, len(parsed_rows) + 1):
            y_line = int(body_top + (ridx * row_h))
            y_edges.append(y_line)

        for y_idx, y_line in enumerate(y_edges):
            x_start = margin
            if y_idx >= 3 and y_idx < len(y_edges) - 1:
                upper_row = y_idx - 3
                lower_row = upper_row + 1
                upper_span = row_to_span_idx.get(upper_row)
                lower_span = row_to_span_idx.get(lower_row)
                if upper_span is not None and upper_span == lower_span:
                    # 같은 그룹 내부 경계는 Description 이후부터 라인 표시 (SN/Task 병합 효과)
                    x_start = margin + col_sn + col_task

            hline = slide.shapes.add_connector(1, x_start, y_line, margin + table_w, y_line)
            hline.line.color.rgb = COLOR_GRID
            hline.line.width = Pt(0.8)

    @staticmethod
    def _to_int(value: Any, default: int = 0) -> int:
        try:
            return int(float(value))
        except Exception:
            return default

    def add_title_slide(self, title: str, subtitle: str = None, template_index: Optional[int] = None):
        """제목 슬라이드 추가"""
        # 템플릿 인덱스 결정
        if template_index is None:
            template_index = 0

        # 통합 템플릿에서 슬라이드 복사
        slide = self._copy_slide_from_prs(self.template_source, template_index)

        # 텍스트 교체
        text_map = {
            '제목을 입력하세요': title,
            '부제목을 입력하세요': subtitle or '',
            '날짜 또는 추가 정보': '',
        }
        self._replace_placeholders_in_slide(slide, text_map)

        self.slide_count += 1
        return slide

    def add_content_slide(self, content: SlideContent, auto_select_layout: bool = True):
        """내용 슬라이드 추가"""
        import tempfile
        import shutil

        # 레이아웃 결정
        if auto_select_layout:
            if content.layout:
                layout_index = self._get_layout_index_from_name(content.layout)
            else:
                layout_index = self._select_best_layout(content)
        else:
            layout_index = 2  # 기본값: basic (index 2 in unified template)

        # 통합 템플릿 경로
        ratio_str = self.aspect_ratio.value
        unified_template_path = self.template_dir / f"templates_{ratio_str}.pptx"

        # generate_with_template와 동일한 방식으로 슬라이드 복사
        temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
        try:
            os.close(temp_fd)

            # 템플릿을 임시 파일로 복사
            shutil.copy(str(unified_template_path), temp_path)

            # 임시 프레젠테이션 로드
            temp_prs = Presentation(temp_path)

            # 필요한 슬라이드 하나만 남기고 모두 삭제
            for j in range(len(temp_prs.slides) - 1, -1, -1):
                if j != layout_index:
                    slide = temp_prs.slides[j]
                    for shape in list(slide.shapes):
                        sp = shape.element
                        sp.getparent().remove(sp)

            # Placeholder 교체
            slide = temp_prs.slides[layout_index]
            text_map = {
                '{slide_title}': content.title,
                '{title}': content.title,
                '{subtitle}': content.subtitle or '',
                '{content}': content.content or '',
                '슬라이드 제목': content.title,
                '제목을 입력하세요': content.title,
                '부제목을 입력하세요': content.subtitle or '',
                '01\n세션 제목': f"{content.title}\n{content.content}" if content.content else content.title,
                '세션 제목': content.content or content.title,
            }

            if content.date:
                text_map.update({
                    '{date}': content.date,
                    '날짜 또는 추가 정보': content.date,
                    'Date': content.date,
                    'DATE': content.date,
                })

            if content.content:
                text_map.update({
                    '본문 내용': content.content,
                    '내용을 입력하세요': content.content,
                    '왼쪽 내용': content.content,
                    '오른쪽 내용': content.content,
                })

            self._replace_placeholders_in_slide(slide, text_map)

            # 표 처리
            if content.table_data and content.table_headers:
                self._add_table_to_slide(slide, content.table_headers, content.table_data)

            # 두 번째 표 처리
            if content.table_data_2 and content.table_headers_2:
                self._add_table2_to_slide(slide, content.table_headers_2, content.table_data_2)

            # 이미지 처리
            if content.image_path:
                self._add_image_to_slide_v2(slide, content.image_path)

            if content.image_path_2:
                self._add_image2_to_slide(slide, content.image_path_2)

            # 임시 프레젠테이션을 출력 파일에 추가
            # 현재 presentation에 슬라이드 복사
            rId = self.presentation.part.relate_to(temp_prs.part, RT.SLIDE)
            self.presentation.part.drop_rel(rId)

            # 슬라이드 추가
            self.presentation.slides._sldIdLst.append(slide.element)

        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)

        self.slide_count += 1
        return slide


    def add_section_slide(self, title: str, subtitle: str = None):
        """섹션 슬라이드 추가"""
        return self.add_title_slide(title, subtitle, template_index=0)

    def add_image_slide(self, title: str, image_path: str, subtitle: str = None):
        """이미지 슬라이드 추가"""
        content = SlideContent(title=title, image_path=image_path, subtitle=subtitle)
        return self.add_content_slide(content)

    def add_table_slide(self, title: str, headers: List[str], data: List[List[str]]):
        """표 슬라이드 추가"""
        content = SlideContent(title=title, table_headers=headers, table_data=data)
        return self.add_content_slide(content)

    def add_blank_slide(self):
        """빈 슬라이드 추가"""
        blank_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(blank_layout)
        self.slide_count += 1
        return slide

    def save(self, output_path: str):
        """프레젠테이션 저장"""
        self.presentation.save(output_path)
        return output_path

    def _copy_slide_from_prs(self, source_prs: Presentation, source_slide_index: int):
        """다른 presentation에서 슬라이드를 복사"""
        from lxml import etree
        from copy import deepcopy

        source_slide = source_prs.slides[source_slide_index]
        source_slide_element = source_slide.element

        # 빈 슬라이드 추가
        blank_layout = self.presentation.slide_layouts[6]
        new_slide = self.presentation.slides.add_slide(blank_layout)
        new_slide_element = new_slide.element

        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        # spTree 내용 복사
        source_sp_tree = source_slide_element.find('.//p:spTree', namespaces=ns)
        new_sp_tree = new_slide_element.find('.//p:spTree', namespaces=ns)

        if source_sp_tree is not None and new_sp_tree is not None:
            for child in list(new_sp_tree):
                new_sp_tree.remove(child)

            for child in source_sp_tree:
                new_child = deepcopy(child)
                new_sp_tree.append(new_child)

        # 배경 요소 복사
        for bg_element in new_slide_element.findall('p:bg', namespaces=ns):
            new_slide_element.remove(bg_element)
        for bg_ref_element in new_slide_element.findall('p:bgRef', namespaces=ns):
            new_slide_element.remove(bg_ref_element)

        source_bg = source_slide_element.find('p:bg', namespaces=ns)
        source_bg_ref = source_slide_element.find('p:bgRef', namespaces=ns)

        if source_bg is not None:
            new_bg = deepcopy(source_bg)
            new_slide_element.append(new_bg)
        elif source_bg_ref is not None:
            new_bg_ref = deepcopy(source_bg_ref)
            new_slide_element.append(new_bg_ref)

        return new_slide

    def _select_best_layout(self, content: SlideContent) -> int:
        """콘텐츠에 가장 적합한 레이아웃 선택 (통합 템플릿: +2 오프셋)"""
        if content.table_data and content.table_data_2:
            return 6  # two_tables (4 + 2)
        elif content.table_data:
            return 5  # large_table (3 + 2)
        elif content.image_path and content.image_path_2:
            return 14  # two_images (12 + 2)
        elif content.image_path:
            return 11  # image_right (본문 왼쪽 + 이미지 오른쪽, 일반적인 문서 형태)
        elif content.content:
            # 줄바꿈이 많으면 two_column
            line_count = content.content.count('\n')
            if line_count >= 3:
                return 3  # two_column (1 + 2)
            return 2  # basic (0 + 2)
        else:
            return 2  # basic (0 + 2)

    def _get_layout_index_from_name(self, layout_name: str) -> int:
        """레이아웃 이름에서 인덱스 가져오기 (통합 템플릿: 제목 0-1, 세션 2, 내용 3-15)"""
        layout_map = {
            # 제목 슬라이드 (0-1)
            'title': 0,
            
            # 내용 슬라이드 (2-14, 원래 인덱스 +2)
            'content_title_text': 2,  # 0 + 2
            'basic': 2,  # 0 + 2
            'two_column': 3,  # 1 + 2
            'three_column': 4,  # 2 + 2
            'large_table': 5,  # 3 + 2
            'two_tables': 6,  # 4 + 2
            'gantt_chart': 7,  # 5 + 2
            'cod_pipeline': 20,  
            'cod_pipeline_phase': 21,
            'equipment_procurement_case': 22,
            'text_small_table': 8,  # 6 + 2
            'text_large_table': 9,  # 7 + 2
            'permits': 9,  # 템플릿 10페이지(1-based) 사용
            'gantt_chart_template10': 9,  # 템플릿 10페이지(1-based) 강제 사용
            'two_tables_vertical': 10,  # 8 + 2
            'image_left': 11,  # 실제 템플릿 인덱스: 이미지 왼쪽 + 본문 오른쪽
            'image_right': 12,  # 실제 템플릿 인덱스: 본문 왼쪽 + 이미지 오른쪽
            'large_image': 13,  # 11 + 2
            'two_images': 14,  # 12 + 2
            'session_title': 15,
            'approval_request': 16,  # Approval request 전용 템플릿
            # 'executive_summary': 17,  # Executive Summary 템플릿
            'main_agreements_lease': 18,  # Main Agreements 템플릿 (템플릿 19페이지)
            'main_agreements': 17,  # Main Agreements 템플릿 (템플릿 19페이지)
            'title_subtitle_table': 19,  # 제목 + 부제목 + 테이블
            'bos_negotiation': 23,  # BOS negotiation 전용 템플릿(표 데이터 공란)
            
        }
        return layout_map.get(layout_name, 3)  # 기본값: basic (index 3)

    def _add_table_to_slide(self, slide, headers: List[str], data: List[List[str]]):
        """슬라이드에 표 추가"""
        from lxml import etree

        def _fit_rect_to_slide_bounds(left: int, top: int, width: int, height: int) -> tuple[int, int, int, int]:
            slide_width_emu = int(self.slide_width * 914400)
            slide_height_emu = int(self.slide_height * 914400)
            safe_margin = int(0.15 * 914400)

            if width <= 0:
                width = max(int(1.0 * 914400), slide_width_emu - (2 * safe_margin))
            if height <= 0:
                height = max(int(1.0 * 914400), slide_height_emu - (2 * safe_margin))

            left = max(safe_margin, left)
            top = max(safe_margin, top)

            max_width = max(int(1.0 * 914400), slide_width_emu - safe_margin - left)
            max_height = max(int(1.0 * 914400), slide_height_emu - safe_margin - top)
            width = min(width, max_width)
            height = min(height, max_height)

            return left, top, width, height

        # 일관된 플레이스홀더 이름 사용
        placeholders_to_try = [
            '표영역',
            '표영역1',
            '표영역2',
        ]

        rows = len(data) + 1  # 헤더 포함
        cols = len(headers)

        # deepcopy 후에는 slide.shapes가 업데이트되지 않으므로 XML 레벨에서 검색
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
              'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        slide_element = slide.element

        for placeholder in placeholders_to_try:
            # XML에서 플레이스홀더 텍스트 찾기
            for sp in slide_element.findall('.//p:sp', namespaces=ns):
                # 모든 <a:t> 요소 검색
                for text_elem in sp.findall('.//a:t', namespaces=ns):
                    if text_elem.text and placeholder in text_elem.text:
                        # spTree의 오프셋과 크기 정보 가져오기
                        # spPr 찾기 - 모든 가능한 위치 시도
                        sp_pr = sp.find('./p:spPr', namespaces=ns)
                        if sp_pr is None:
                            sp_pr = sp.find('./a:spPr', namespaces=ns)
                        if sp_pr is None:
                            # .//p:spPr 또는 .//a:spPr로 찾기
                            sp_pr = sp.find('.//p:spPr', namespaces=ns)
                            if sp_pr is None:
                                sp_pr = sp.find('.//a:spPr', namespaces=ns)

                        if sp_pr is not None:
                            # xfrm 찾기
                            xfrm = sp_pr.find('./a:xfrm', namespaces=ns)
                            if xfrm is None:
                                xfrm = sp_pr.find('.//a:xfrm', namespaces=ns)

                            if xfrm is not None:
                                # off와 ext 찾기
                                off = xfrm.find('./a:off', namespaces=ns)
                                ext = xfrm.find('./a:ext', namespaces=ns)

                                if off is not None and ext is not None:
                                    left = int(off.get('x', 0))
                                    top = int(off.get('y', 0))
                                    width = int(ext.get('cx', 0))
                                    height = int(ext.get('cy', 0))

                                    # shape 삭제
                                    sp.getparent().remove(sp)

                                    # 표 생성 (슬라이드 경계 내로 보정)
                                    left, top, width, height = _fit_rect_to_slide_bounds(left, top, width, height)
                                    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                                    self.style_table(table, headers, data)
                                    return  # 성공하면 종료

    def _add_table_to_slide_with_placeholder(self, slide, headers: List[str], data: List[List[str]], placeholder: str):
        """플레이스홀더 위치에 표 추가"""
        def _fit_rect_to_slide_bounds(left: int, top: int, width: int, height: int) -> tuple[int, int, int, int]:
            slide_width_emu = int(self.slide_width * 914400)
            slide_height_emu = int(self.slide_height * 914400)
            safe_margin = int(0.15 * 914400)

            if width <= 0:
                width = max(int(1.0 * 914400), slide_width_emu - (2 * safe_margin))
            if height <= 0:
                height = max(int(1.0 * 914400), slide_height_emu - (2 * safe_margin))

            left = max(safe_margin, left)
            top = max(safe_margin, top)

            max_width = max(int(1.0 * 914400), slide_width_emu - safe_margin - left)
            max_height = max(int(1.0 * 914400), slide_height_emu - safe_margin - top)
            width = min(width, max_width)
            height = min(height, max_height)

            return left, top, width, height

        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                if placeholder in shape.text_frame.text:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    # 플레이스홀더 삭제
                    sp = shape.element
                    sp.getparent().remove(sp)

                    # 표 생성 (슬라이드 경계 내로 보정)
                    rows = len(data) + 1
                    cols = len(headers)
                    left, top, width, height = _fit_rect_to_slide_bounds(left, top, width, height)
                    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                    self.style_table(table, headers, data)
                    return True
        return False

    def _add_table_to_slide_custom(self, slide, headers: List[str], data: List[List[str]], placeholder: str):
        return self._add_table_to_slide_with_placeholder(slide, headers, data, placeholder)

    def _add_table2_to_slide(self, slide, headers: List[str], data: List[List[str]]):
        """두 번째 표 추가 (two_tables, two_tables_vertical 레이아웃용)"""
        # 일관된 플레이스홀더 이름 사용 (두 번째 표는 표영역2)
        placeholders_to_try = ['표영역2']

        for placeholder in placeholders_to_try:
            if self._add_table_to_slide_with_placeholder(slide, headers, data, placeholder):
                return  # 성공하면 종료

        # 플레이스홀더를 찾지 못하면 경고 출력 (무시하고 계속)
        pass

    def _add_image_to_slide_v2(self, slide, image_path: str):
        """슬라이드에 이미지 추가 (플레이스홀더 교체, 비율 유지)

        사용자 요구사항:
        - 이미지 영역의 가로사이즈에 맞춰 비율 조정
        - 단, 원본 이미지가 플레이스홀더보다 작은 경우 확대하지 않음 (원본 크기 유지)
        """
        if not os.path.exists(image_path):
            return

        # 원본 이미지 크기 확인
        from PIL import Image as PILImage
        pil_img = PILImage.open(image_path)
        img_width, img_height = pil_img.size
        img_aspect = img_width / img_height

        # 원본 이미지 크기를 EMU로 변환
        img_width_emu = img_width * 914400 / 96  # 96 DPI 기준
        img_height_emu = img_height * 914400 / 96

        # 이미지 플레이스홀더 찾기 (일관된 이름)
        placeholders_to_try = [
            '이미지영역1',
            '이미지영역2',
            '이미지영역',
        ]

        # XML 직접 파싱으로 플레이스홀더 찾기 (python-pptx가 모든 shape를 노출하지 않는 문제 해결)
        from lxml import etree
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
              'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        slide_element = slide.element
        placeholder_found = False

        # 디버깅: 모든 텍스트 내용 출력
        all_texts = []
        for sp in slide_element.findall('.//p:sp', namespaces=ns):
            for text_elem in sp.findall('.//a:t', namespaces=ns):
                if text_elem.text:
                    all_texts.append(text_elem.text[:50] if len(text_elem.text) > 50 else text_elem.text)
        if all_texts:
            print(f"[DEBUG]   슬라이드의 모든 텍스트: {all_texts[:5]}")  # 처음 5개만 출력

        for placeholder in placeholders_to_try:
            # XML에서 플레이스홀더 텍스트 찾기 (모든 p:sp 요소 검색)
            for sp in slide_element.findall('.//p:sp', namespaces=ns):
                # 모든 <a:t> 요소 검색
                for text_elem in sp.findall('.//a:t', namespaces=ns):
                    if text_elem.text and placeholder in text_elem.text:
                        # spPr 찾기
                        sp_pr = sp.find('./p:spPr', namespaces=ns)
                        if sp_pr is None:
                            sp_pr = sp.find('./a:spPr', namespaces=ns)

                        if sp_pr is not None:
                            # xfrm에서 위치와 크기 정보 가져오기
                            xfrm = sp_pr.find('./a:xfrm', namespaces=ns)
                            if xfrm is not None:
                                off = xfrm.find('./a:off', namespaces=ns)
                                ext = xfrm.find('./a:ext', namespaces=ns)

                                if off is not None and ext is not None:
                                    left = int(off.get('x', 0))
                                    top = int(off.get('y', 0))
                                    placeholder_width = int(ext.get('cx', 0))
                                    placeholder_height = int(ext.get('cy', 0))
                                    placeholder_aspect = float(placeholder_width) / float(placeholder_height)
                                    placeholder_found = True
                                    break
                if placeholder_found:
                    break

            if placeholder_found:
                break

        if not placeholder_found:
            # 플레이스홀더를 찾지 못한 경우: 슬라이드 영역 사용
            print(f"[DEBUG]   플레이스홀더를 찾지 못함 -> 슬라이드 영역 사용")

            # 슬라이드 크기 가져오기 (첫 번째 shape에서 추론)
            slide_width = 12191695  # 16:9 기본값
            slide_height = 6858000

            # 첫 번째 shape(배경)의 크기 확인
            for s in slide.shapes:
                if s.shape_type == 1:  # Rectangle (배경)
                    slide_width = s.width
                    slide_height = s.height
                    break

            # 제목 영역(약 1.3인치)과 여백(0.5인치)을 제외한 영역 사용
            margin = int(0.5 * 914400)  # 0.5인치 여백
            title_height = int(1.3 * 914400)  # 1.3인치 제목 영역
            left = margin
            top = title_height
            placeholder_width = slide_width - 2 * margin
            placeholder_height = slide_height - top - margin

        # 플레이스홀더 크기 검사
        MIN_HEIGHT_EMU = 1.5 * 914400  # 최소 높이 1.5인치 (image_right 레이아웃의 세로 길이 고려)
        MIN_WIDTH_EMU = 3 * 914400   # 최소 너비 3인치 (image_left 레이아웃의 가로 길이 고려)

        # 플레이스홀더가 너무 작으면 슬라이드 사용 가능 영역으로 대체
        # (단, 이것은 템플릿이 실수로 수정되었을 때의 안전장치임)
        if placeholder_found and (placeholder_height < MIN_HEIGHT_EMU or placeholder_width < MIN_WIDTH_EMU):
            print(f"[DEBUG]   플레이스홀더가 너무 작음 -> 슬라이드 영역 사용")
            print(f"[DEBUG]     원본 플레이스홀더: {placeholder_width}x{placeholder_height} EMU")

            # 슬라이드 크기 가져오기 (첫 번째 shape에서 추론)
            slide_width = 12191695  # 16:9 기본값
            slide_height = 6858000

            # 첫 번째 shape(배경)의 크기 확인
            for s in slide.shapes:
                if s.shape_type == 1:  # Rectangle (배경)
                    slide_width = s.width
                    slide_height = s.height
                    break

            # 제목 영역(약 1.3인치)과 여백(0.5인치)을 제외한 영역 사용
            margin = int(0.5 * 914400)  # 0.5인치 여백
            title_height = int(1.3 * 914400)  # 1.3인치 제목 영역
            left = margin
            top = title_height
            placeholder_width = slide_width - 2 * margin
            placeholder_height = slide_height - top - margin
            print(f"[DEBUG]     조정된 영역: {placeholder_width}x{placeholder_height} EMU")

        placeholder_aspect = float(placeholder_width) / float(placeholder_height)

        # 디버깅 정보
        print(f"[DEBUG] 이미지: {os.path.basename(image_path)}")
        print(f"[DEBUG]   이미지 크기: {img_width}x{img_height} pixels ({img_width_emu:.0f}x{img_height_emu:.0f} EMU, aspect={img_aspect:.3f})")
        print(f"[DEBUG]   플레이스홀더: {placeholder_width}x{placeholder_height} EMU (aspect={placeholder_aspect:.3f})")

        # 이미지 크기 결정 (플레이스홀더 내에 완전히 들어가도록, 비율 유지)
        # 원본 이미지가 플레이스홀더보다 작은 경우: 원본 크기 사용 (확대하지 않음)
        # 원본 이미지가 플레이스홀더보다 큰 경우: 플레이스홀더에 맞춰 축소 (가로+세로 모두 고려)

        if img_width_emu <= placeholder_width and img_height_emu <= placeholder_height:
            # 원본이 플레이스홀더보다 작거나 같음: 원본 크기 사용
            width = int(img_width_emu)
            height = int(img_height_emu)
            print(f"[DEBUG]   원본 크기 사용 (확대하지 않음)")
        else:
            # 원본이 더 큼: 플레이스홀더에 맞춰 축소 (비율 유지)
            # 가로 기준 계산
            width_by_w = placeholder_width
            height_by_w = int(placeholder_width / img_aspect)

            # 세로 기준 계산
            height_by_h = placeholder_height
            width_by_h = int(placeholder_height * img_aspect)

            # 둘 중 플레이스홀더에 완전히 들어가는 쪽 선택
            if height_by_w <= placeholder_height:
                # 가로 기준으로 축소하면 세로도 안에 들어감
                width = width_by_w
                height = height_by_w
                print(f"[DEBUG]   가로 기준 축소 (width={width}, height={height})")
            else:
                # 세로 기준으로 축소해야 함
                width = width_by_h
                height = height_by_h
                print(f"[DEBUG]   세로 기준 축소 (width={width}, height={height})")

        # 세로 중앙 정렬 (이미지가 플레이스홀더보다 작은 경우)
        if height < placeholder_height:
            top = top + (placeholder_height - height) // 2

        print(f"[DEBUG]   조정된 크기: width={width}, height={height}")

        # 이미지 추가
        slide.shapes.add_picture(image_path, left, top, width, height)

        # 플레이스홀더가 있으면 제거 (XML 레벨에서)
        if placeholder_found:
            # 플레이스홀더 shape 제거
            for sp in slide_element.findall('.//p:sp', namespaces=ns):
                for text_elem in sp.findall('.//a:t', namespaces=ns):
                    if text_elem.text and any(p in text_elem.text for p in placeholders_to_try):
                        sp.getparent().remove(sp)
                        break
                else:
                    continue
                break

    def _add_image2_to_slide(self, slide, image_path: str):
        """두 번째 이미지 추가 (two_images 레이아웃용)"""
        # _add_image_to_slide_v2는 여러 플레이스홀더를 시도하므로
        # 첫 번째 이미지가 '이미지영역1'을 사용하면
        # 두 번째 호출은 자동으로 '이미지영역2'를 사용합니다
        self._add_image_to_slide_v2(slide, image_path)

    def generate_with_template(self, content_list: List[SlideContent], output_path: str) -> str:
        """
        템플릿을 기반으로 문서 생성

        방식 (통합 템플릿):
        1. 통합 템플릿을 출력 파일로 복사 (모든 리소스 포함)
        2. 첫 번째 슬라이드를 제목 슬라이드로 처리
        3. 불필요한 템플릿 슬라이드 제거 (슬라이드 2-14)
        4. 내용 슬라이드는 템플릿에서 복사하여 추가
        5. Placeholder만 교체

        Args:
            content_list: SlideContent 객체 리스트
            output_path: 출력 파일 경로

        Returns:
            생성된 파일 경로
        """
        import zipfile

        ratio_str = self.aspect_ratio.value

        # 통합 템플릿 경로
        unified_template_path = self.template_dir / f"templates_{ratio_str}.pptx"

        # 1. 통합 템플릿을 출력 파일로 복사
        shutil.copy(str(unified_template_path), str(output_path))

        # 주의: _ensure_table_styles()를 호출하면 python-pptx가 presentation.xml을 초기화함
        # 따라서 _remove_template_slides() 이후에 호출해야 함

        # 2. 첫 번째 슬라이드를 제목 슬라이드로 처리
        prs = Presentation(str(output_path))
        if content_list:
            # 사용자가 order를 지정했다면 그것을 기준으로 정렬
            # order가 문자열("01", "10")로 들어오는 경우가 있어 숫자로 변환해 정렬하지 않으면
            # "10"이 "2"보다 앞서는 문제가 발생함. 숫자로 캐스팅 가능할 때는 숫자 기준으로 정렬한다.
            def _normalized_order(val):
                if val is None:
                    return None
                if isinstance(val, (int, float)):
                    return float(val)
                if isinstance(val, str):
                    try:
                        return float(val)
                    except ValueError:
                        return None
                return None

            ordered = []
            for idx, item in enumerate(content_list):
                norm = _normalized_order(item.order)
                ordered.append((idx, item, norm))

            content_list = [
                item for _, item, _ in sorted(
                    ordered,
                    key=lambda x: (
                        0 if x[2] is not None else 1,          # order가 있으면 우선
                        x[2] if x[2] is not None else x[0]     # 숫자 order or 기존 순서
                    )
                )
            ]

            title_slide = prs.slides[0]
            title_content = content_list[0]

            text_map = {
                '{title}': title_content.title,
                '{subtitle}': title_content.subtitle or '',
                '제목을 입력하세요': title_content.title,
                '부제목을 입력하세요': title_content.subtitle or '',
                '날짜 또는 추가 정보': '',
            }
            self._replace_placeholders_in_slide(title_slide, text_map)

            prs.save(str(output_path))

        # 3. 불필요한 템플릿 슬라이드 제거 (슬라이드 인덱스 1-14)
        # 템플릿은 15개 슬라이드를 가짐: 0-1 (제목), 2-14 (내용 레이아웃)
        # 우리는 슬라이드 0만 제목으로 사용하고, 나머지는 직접 추가할 것이므로
        # 슬라이드 1-14 (두 번째 제목 + 모든 내용 레이아웃)를 제거해야 함
        template_slide_count = len(Presentation(str(unified_template_path)).slides)
        remove_start = 1
        remove_end = max(0, template_slide_count - 1)
        print("[DEBUG] 템플릿 슬라이드 제거 시작")
        print(f"[DEBUG]   제거 범위: start_idx={remove_start}, end_idx={remove_end} (총 {template_slide_count}장)")
        self._remove_template_slides(output_path, start_idx=remove_start, end_idx=remove_end)
        print("[DEBUG] 템플릿 슬라이드 제거 완료")

        # 3-1. LibreOffice 등으로 수정한 템플릿 복구 (tableStyles.xml 등 누락 파일 복구)
        # 슬라이드 제거 후에 호출해야 presentation.xml이 초기화되지 않음
        self._ensure_table_styles(output_path)

        # 4. 내용 슬라이드 처리
        # 통합 템플릿 프레젠테이션 로드 (source로 사용)
        template = Presentation(str(unified_template_path))

        for i, content in enumerate(content_list[1:], 1):
            # 레이아웃 인덱스 결정
            if content.layout:
                layout_idx = self._get_layout_index_from_name(content.layout)
            else:
                layout_idx = self._select_best_layout(content)

            if layout_idx < len(template.slides):
                # 임시 파일 사용 - 필요한 슬라이드만 남기기
                temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
                try:
                    os.close(temp_fd)

                    # 템플릿을 임시 파일로 복사
                    shutil.copy(str(unified_template_path), temp_path)

                    # 임시 프레젠테이션 로드
                    temp_prs = Presentation(temp_path)

                    # Placeholder 교체 - layout_idx에 해당하는 슬라이드 사용
                    slide = temp_prs.slides[layout_idx]

                    # Session Title은 전용 렌더러로 텍스트를 구성
                    if content.layout == 'session_title':
                        self._render_session_title_layout(slide, temp_prs, content)

                        # 임시 파일 저장
                        temp_prs.save(temp_path)

                        # 현재 출력 파일의 슬라이드 수 확인
                        with zipfile.ZipFile(output_path, 'r') as output_zip:
                            current_slide_count = len([f for f in output_zip.namelist() if f.startswith('ppt/slides/slide') and not f.endswith('.rels')])

                        # ZIP 레벨로 슬라이드 복사 (layout_idx에 해당하는 슬라이드)
                        self._copy_slide_zip(temp_path, str(output_path), layout_idx, current_slide_count)
                        continue

                    # Approval Request는 sample placeholder를 입력 데이터로 치환
                    if content.layout == 'approval_request':
                        self._render_approval_request_layout(slide, content)

                        # 임시 파일 저장
                        temp_prs.save(temp_path)

                        # 현재 출력 파일의 슬라이드 수 확인
                        with zipfile.ZipFile(output_path, 'r') as output_zip:
                            current_slide_count = len([f for f in output_zip.namelist() if f.startswith('ppt/slides/slide') and not f.endswith('.rels')])

                        # ZIP 레벨로 슬라이드 복사 (layout_idx에 해당하는 슬라이드)
                        self._copy_slide_zip(temp_path, str(output_path), layout_idx, current_slide_count)
                        continue

                    # Title + Subtitle + Table은 전용 렌더러로 직접 구성
                    if content.layout == 'title_subtitle_table':
                        self._render_title_subtitle_table_layout(slide, temp_prs, content)

                        # 임시 파일 저장
                        temp_prs.save(temp_path)

                        # 현재 출력 파일의 슬라이드 수 확인
                        with zipfile.ZipFile(output_path, 'r') as output_zip:
                            current_slide_count = len([f for f in output_zip.namelist() if f.startswith('ppt/slides/slide') and not f.endswith('.rels')])

                        # ZIP 레벨로 슬라이드 복사 (layout_idx에 해당하는 슬라이드)
                        self._copy_slide_zip(temp_path, str(output_path), layout_idx, current_slide_count)
                        continue

                    # Executive Summary는 전용 렌더러로 표를 구성
                    if content.layout == 'executive_summary':
                        self._render_executive_summary_layout(slide, temp_prs, content)

                        # 임시 파일 저장
                        temp_prs.save(temp_path)

                        # 현재 출력 파일의 슬라이드 수 확인
                        with zipfile.ZipFile(output_path, 'r') as output_zip:
                            current_slide_count = len([f for f in output_zip.namelist() if f.startswith('ppt/slides/slide') and not f.endswith('.rels')])

                        # ZIP 레벨로 슬라이드 복사 (layout_idx에 해당하는 슬라이드)
                        self._copy_slide_zip(temp_path, str(output_path), layout_idx, current_slide_count)
                        continue

                    # CAPEX phase 간트는 실제 막대형 차트로 렌더링
                    if content.layout == 'gantt_chart_template10':
                        self._render_gantt_plan_layout(slide, temp_prs, content)

                        # 임시 파일 저장
                        temp_prs.save(temp_path)

                        # 현재 출력 파일의 슬라이드 수 확인
                        with zipfile.ZipFile(output_path, 'r') as output_zip:
                            current_slide_count = len([f for f in output_zip.namelist() if f.startswith('ppt/slides/slide') and not f.endswith('.rels')])

                        # ZIP 레벨로 슬라이드 복사 (layout_idx에 해당하는 슬라이드)
                        self._copy_slide_zip(temp_path, str(output_path), layout_idx, current_slide_count)
                        continue

                    text_map = {
                        '{slide_title}': content.title,
                        '{title}': content.title,
                        '{subtitle}': content.subtitle or '',
                        '{content}': content.content or '',
                        '슬라이드 제목': content.title,
                        '제목을 입력하세요': content.title,
                        '부제목을 입력하세요': content.subtitle or '',
                    }

                    # 이미지 처리 (플레이스홀더 찾기 위해 텍스트 교체 전에 수행)
                    if content.image_path and os.path.exists(content.image_path):
                        self._add_image_to_slide_v2(slide, content.image_path)

                    if content.image_path_2 and os.path.exists(content.image_path_2):
                        self._add_image_to_slide_v2(slide, content.image_path_2)

                    # 표 처리
                    if content.table_data and content.table_headers:
                        self._add_table_to_slide(slide, content.table_headers, content.table_data)

                    if content.table_data_2 and content.table_headers_2:
                        self._add_table2_to_slide(slide, content.table_headers_2, content.table_data_2)

                    # 텍스트 처리 (이미지/표 추가 후 수행)
                    if content.content:
                        text_map.update({
                            '본문 내용': content.content,
                            '내용을 입력하세요': content.content,
                            '왼쪽 내용': content.content,
                            '오른쪽 내용': content.content,
                        })

                    # 템플릿의 기본 텍스트는 항상 제거 (content가 있든 없든)
                    # 주의: '내용'은 실제 컨텐츠에도 포함될 수 있으므로 제외하고
                    # _replace_placeholders_in_slide에서 별도 처리
                    text_map.update({
                        '• 여러 줄 작성 가능': '',
                        '• 줄간격 자동 조정': '',
                        '첫 번째': '',
                        '두 번째': '',
                        '세 번째': '',
                    })

                    # content가 없는 경우, 추가 placeholder들을 빈 문자열로 교체
                    if not content.content:
                        text_map.update({
                            '내용을 입력하세요': '',
                            '본문 내용': '',
                            '왼쪽 내용': '',
                            '오른쪽 내용': '',
                        })

                    self._replace_placeholders_in_slide(slide, text_map)

                    # 임시 파일 저장
                    temp_prs.save(temp_path)

                    # 현재 출력 파일의 슬라이드 수 확인
                    with zipfile.ZipFile(output_path, 'r') as output_zip:
                        current_slide_count = len([f for f in output_zip.namelist() if f.startswith('ppt/slides/slide') and not f.endswith('.rels')])

                    # ZIP 레벨로 슬라이드 복사 (layout_idx에 해당하는 슬라이드)
                    self._copy_slide_zip(temp_path, str(output_path), layout_idx, current_slide_count)

                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
            else:
                print(
                    f"[WARN] layout_idx {layout_idx} out of range "
                    f"(template slides: {len(template.slides)}). Using blank slide."
                )
                temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
                try:
                    os.close(temp_fd)
                    shutil.copy(str(unified_template_path), temp_path)
                    temp_prs = Presentation(temp_path)
                    blank_layout = temp_prs.slide_layouts[6]
                    blank_slide = temp_prs.slides.add_slide(blank_layout)
                    temp_prs.save(temp_path)

                    with zipfile.ZipFile(output_path, 'r') as output_zip:
                        current_slide_count = len([
                            f for f in output_zip.namelist()
                            if f.startswith('ppt/slides/slide') and not f.endswith('.rels')
                        ])

                    self._copy_slide_zip(temp_path, str(output_path), blank_slide, current_slide_count)
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)

        return str(output_path)

    def _remove_template_slides(self, output_path: str, start_idx: int, end_idx: int):
        """PPTX 파일에서 지정된 범위의 슬라이드를 제거 (ZIP 레벨 조작)"""
        import zipfile
        from lxml import etree

        # 슬라이드 파일 목록 (제거할 슬라이드)
        slides_to_remove = list(range(start_idx, end_idx + 1))
        print(f"[DEBUG] _remove_template_slides: 제거할 슬라이드 = {slides_to_remove}")

        # 임시 파일 생성
        temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
        try:
            os.close(temp_fd)

            # ZIP 파일을 직접 조작하여 슬라이드 제거
            with zipfile.ZipFile(output_path, 'r') as input_zip:
                with zipfile.ZipFile(temp_path, 'w', zipfile.ZIP_DEFLATED) as output_zip:
                    for item in input_zip.infolist():
                        # 슬라이드 XML 파일 제거
                        filename = item.filename
                        should_remove = False

                        # 슬라이드 파일 확인
                        if filename.startswith('ppt/slides/slide') and filename.endswith('.xml'):
                            # 슬라이드 번호 추출
                            try:
                                slide_num = int(filename.split('slide')[-1].replace('.xml', ''))
                                # 슬라이드 번호는 1-based, start_idx/end_idx는 0-based
                                if slide_num - 1 in slides_to_remove:
                                    should_remove = True
                                    print(f"[DEBUG]   제거: {filename}")
                            except ValueError:
                                pass

                        # 슬라이드 관계 파일 제거
                        if filename.startswith('ppt/slides/_rels/slide') and '.xml.rels' in filename:
                            try:
                                slide_num = int(filename.split('slide')[-1].replace('.xml.rels', ''))
                                if slide_num - 1 in slides_to_remove:
                                    should_remove = True
                                    print(f"[DEBUG]   제거: {filename}")
                            except ValueError:
                                pass

                        if not should_remove:
                            # 파일 내용 읽기 및 쓰기
                            data = input_zip.read(item.filename)
                            output_zip.writestr(item, data)

            # 원본 파일을 임시 파일로 교체
            shutil.move(temp_path, output_path)

            # presentation.xml 수정 (제거된 슬라이드 참조 업데이트)
            print("[DEBUG] presentation.xml 업데이트 시작")
            self._update_presentation_xml_after_removal(output_path, slides_to_remove)
            print("[DEBUG] presentation.xml 업데이트 완료")

        except Exception as e:
            # 오류 발생 시 임시 파일 정리
            print(f"[ERROR] _remove_template_slides: {e}")
            if os.path.exists(temp_path):
                os.unlink(temp_path)
            raise

    def _update_presentation_xml_after_removal(self, output_path: str, removed_indices: List[int]):
        """슬라이드 제거 후 presentation.xml의 슬라이드 ID와 참조 업데이트"""
        import zipfile
        from lxml import etree

        # 디버그: 제거할 인덱스 출력
        print(f"[DEBUG] _update_presentation_xml_after_removal: removed_indices = {removed_indices}")

        with zipfile.ZipFile(output_path, 'r') as zip_ref:
            # presentation.xml 읽기
            pres_xml = zip_ref.read('ppt/presentation.xml')

        # XML 파싱 및 수정
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
              'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
        root = etree.fromstring(pres_xml)

        # 모든 슬라이드 참조(p:sldId) 찾기 및 제거할 슬라이드 필터링
        sld_id_list = root.find('.//p:sldIdLst', namespaces=ns)
        if sld_id_list is not None:
            # 현재 슬라이드 ID 목록 가져오기
            sld_ids = sld_id_list.findall('p:sldId', namespaces=ns)
            print(f"[DEBUG]   현재 슬라이드 ID 수: {len(sld_ids)}")

            # 제거할 슬라이드 (인덱스 기준)
            # removed_indices는 0-based (slide2.xml = index 1, slide1.xml = index 0)
            # sld_ids 리스트도 0-based (sld_ids[0] = slide1.xml 참조)
            # 따라서 removed_indices의 값이 바로 sld_ids의 인덱스
            for idx in sorted(removed_indices, reverse=True):
                list_idx = idx  # 0-based 인덱스 직접 사용
                if 0 <= list_idx < len(sld_ids):
                    print(f"[DEBUG]   제거 sld_ids[{list_idx}]")
                    sld_id_list.remove(sld_ids[list_idx])

            print(f"[DEBUG]   남은 슬라이드 ID 수: {len(sld_id_list.findall('p:sldId', namespaces=ns))}")

        # 수정된 XML 다시 쓰기
        with zipfile.ZipFile(output_path, 'r') as zip_ref:
            # 모든 파일 읽기
            files = {name: zip_ref.read(name) for name in zip_ref.namelist()}

        # presentation.xml만 업데이트 (xml_declaration=False로 설정)
        files['ppt/presentation.xml'] = etree.tostring(root, encoding='UTF-8')

        # 새 ZIP 파일에 쓰기
        temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
        try:
            os.close(temp_fd)
            with zipfile.ZipFile(temp_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                for name, data in files.items():
                    zip_out.writestr(name, data)

            shutil.move(temp_path, output_path)
        except Exception:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
            raise

    def _ensure_table_styles(self, pptx_path: str):
        """tableStyles.xml이 없으면 생성 (LibreOffice 등으로 저장한 템플릿 대응)

        ZIP 파일을 직접 수정하는 대신, python-pptx로 파일을 열고 닫으면서
        자동으로 tableStyles.xml이 포함되도록 함
        """
        try:
            # 파일을 열고 즉시 저장하면 python-pptx가 누락된 파일들을 자동으로 복구
            prs = Presentation(str(pptx_path))
            prs.save(str(pptx_path))
            print("[DEBUG] 템플릿 복구 완료 (python-pptx)")
        except Exception as e:
            print(f"[WARNING] 템플릿 복구 실패 (무시함): {e}")
            # 실패해도 계속 진행 (style_table에서 직접 스타일 적용하므로)

    def _copy_slide_zip(self, template_path, output_path, source_slide_idx, dest_slide_idx):
        """슬라이드를 ZIP 레벨로 복사 (미디어 파일도 함께 복사)"""
        import zipfile
        from lxml import etree

        dest_slide_number = dest_slide_idx + 1

        source_slide_file = f'ppt/slides/slide{source_slide_idx + 1}.xml'
        source_rel_file = f'ppt/slides/_rels/slide{source_slide_idx + 1}.xml.rels'

        dest_slide_file = f'ppt/slides/slide{dest_slide_number}.xml'
        dest_rel_file = f'ppt/slides/_rels/slide{dest_slide_number}.xml.rels'

        # 슬라이드 XML과 relationship 파일 복사
        with zipfile.ZipFile(template_path, 'r') as template_zip:
            with zipfile.ZipFile(output_path, 'a') as out_zip:
                if source_slide_file in template_zip.namelist():
                    slide_data = template_zip.read(source_slide_file)
                    out_zip.writestr(dest_slide_file, slide_data)

                if source_rel_file in template_zip.namelist():
                    rel_data = template_zip.read(source_rel_file)
                    out_zip.writestr(dest_rel_file, rel_data)

        # 미디어 파일 복사
        self._copy_media_files_for_slide(template_path, output_path, source_slide_idx + 1)

        # Content_Types.xml에 슬라이드 추가
        self._add_content_type_for_slide(output_path, dest_slide_number)

        # 프레젠테이션 업데이트
        self._add_slide_to_presentation(output_path, dest_slide_idx)

    def _copy_all_media_from_template(self, template_path, output_path):
        """템플릿의 모든 미디어 파일을 출력 파일로 복사"""
        import zipfile
        from lxml import etree

        # 이미 복사된 미디어 파일 추적
        existing_media = set()
        existing_extensions = set()

        try:
            with zipfile.ZipFile(output_path, 'r') as output_zip:
                for name in output_zip.namelist():
                    if name.startswith('ppt/media/'):
                        existing_media.add(name)
                        ext = name.rsplit('.', 1)[-1].lower()
                        existing_extensions.add(ext)
        except Exception:
            pass

        new_extensions = set()

        # 템플릿의 모든 이미지 Content Type 읽기
        template_extensions = set()
        try:
            with zipfile.ZipFile(template_path, 'r') as template_zip:
                content_types_file = '[Content_Types].xml'
                if content_types_file in template_zip.namelist():
                    content_types_xml = template_zip.read(content_types_file)
                    types_root = etree.fromstring(content_types_xml)
                    for child in types_root:
                        if child.tag.endswith('}Default') or child.tag == 'Default':
                            ext = child.get('Extension')
                            if ext:
                                content_type = child.get('ContentType')
                                if content_type and 'image' in content_type:
                                    template_extensions.add(ext.lower())
        except Exception:
            pass

        # 템플릿에 있는 이미지 확장자 중 출력 파일에 없는 것들을 추가
        for ext in template_extensions:
            if ext not in existing_extensions:
                new_extensions.add(ext)

        # 모든 미디어 파일 복사
        with zipfile.ZipFile(template_path, 'r') as template_zip:
            with zipfile.ZipFile(output_path, 'a') as out_zip:
                for name in template_zip.namelist():
                    if name.startswith('ppt/media/') and name not in existing_media:
                        try:
                            media_data = template_zip.read(name)
                            out_zip.writestr(name, media_data)
                            existing_media.add(name)

                            ext = name.rsplit('.', 1)[-1].lower()
                            if ext not in existing_extensions:
                                new_extensions.add(ext)
                        except Exception as e:
                            print(f"Warning: Could not copy media file {name}: {e}")

        # 새로운 확장자가 있으면 Content_Types.xml 업데이트
        if new_extensions:
            self._add_content_types_for_extensions(output_path, new_extensions)

    def _copy_media_files_for_slide(self, template_path, output_path, slide_number):
        """슬라이드가 참조하는 모든 미디어 파일 복사"""
        import zipfile
        from lxml import etree

        rel_file = f'ppt/slides/_rels/slide{slide_number}.xml.rels'

        try:
            with zipfile.ZipFile(template_path, 'r') as template_zip:
                if rel_file not in template_zip.namelist():
                    return

                rel_xml = template_zip.read(rel_file)
        except Exception as e:
            print(f"Warning: Could not read {rel_file} from template: {e}")
            return

        root = etree.fromstring(rel_xml)

        existing_media = set()
        existing_extensions = set()
        try:
            with zipfile.ZipFile(output_path, 'r') as output_zip:
                for name in output_zip.namelist():
                    if name.startswith('ppt/media/'):
                        existing_media.add(name)
                        ext = name.rsplit('.', 1)[-1].lower()
                        existing_extensions.add(ext)
        except Exception:
            pass

        new_extensions = set()

        # 템플릿의 Content_Types.xml에서 이미 모든 이미지 확장자 정보를 가져오기
        template_extensions = set()
        try:
            with zipfile.ZipFile(template_path, 'r') as template_zip:
                content_types_file = '[Content_Types].xml'
                if content_types_file in template_zip.namelist():
                    content_types_xml = template_zip.read(content_types_file)
                    types_root = etree.fromstring(content_types_xml)
                    for child in types_root:
                        if child.tag.endswith('}Default') or child.tag == 'Default':
                            ext = child.get('Extension')
                            if ext:
                                content_type = child.get('ContentType')
                                if content_type and 'image' in content_type:
                                    template_extensions.add(ext.lower())
        except Exception:
            pass

        # 템플릿에 있는 이미지 확장자 중 출력 파일에 없는 것들을 추가
        for ext in template_extensions:
            if ext not in existing_extensions:
                new_extensions.add(ext)

        with zipfile.ZipFile(template_path, 'r') as template_zip:
            with zipfile.ZipFile(output_path, 'a') as out_zip:
                for rel in root:
                    rel_type = rel.get('Type')
                    if rel_type and 'image' in rel_type:
                        target = rel.get('Target')
                        if target:
                            if target.startswith('../'):
                                media_path = 'ppt/' + target[3:]
                            else:
                                media_path = 'ppt/' + target

                            if media_path in template_zip.namelist() and media_path not in existing_media:
                                try:
                                    media_data = template_zip.read(media_path)
                                    out_zip.writestr(media_path, media_data)
                                    existing_media.add(media_path)

                                    ext = media_path.rsplit('.', 1)[-1].lower()
                                    if ext not in existing_extensions:
                                        new_extensions.add(ext)
                                except Exception as e:
                                    print(f"Warning: Could not copy media file {media_path}: {e}")

        if new_extensions:
            self._add_content_types_for_extensions(output_path, new_extensions)

    def _copy_image_content_types(self, template_path, output_path):
        """템플릿의 모든 이미지 Content Type을 출력 파일에 복사"""
        import zipfile
        from lxml import etree

        content_types_file = '[Content_Types].xml'

        # 템플릿의 이미지 확장자들 읽기
        template_extensions = {}
        try:
            with zipfile.ZipFile(template_path, 'r') as template_zip:
                if content_types_file not in template_zip.namelist():
                    return

                content_types_xml = template_zip.read(content_types_file)
                types_root = etree.fromstring(content_types_xml)

                for child in types_root:
                    if child.tag.endswith('}Default') or child.tag == 'Default':
                        ext = child.get('Extension')
                        content_type = child.get('ContentType')
                        if ext and content_type and 'image' in content_type:
                            template_extensions[ext.lower()] = content_type
        except Exception as e:
            print(f"Warning: Could not read template content types: {e}")
            return

        # 출력 파일의 기존 확장자들 읽기
        existing_extensions = set()
        try:
            with zipfile.ZipFile(output_path, 'r') as output_zip:
                if content_types_file in output_zip.namelist():
                    content_types_xml = output_zip.read(content_types_file)
                    types_root = etree.fromstring(content_types_xml)

                    for child in types_root:
                        if child.tag.endswith('}Default') or child.tag == 'Default':
                            ext = child.get('Extension')
                            if ext:
                                existing_extensions.add(ext.lower())
        except Exception as e:
            print(f"Warning: Could not read output content types: {e}")
            return

        # 템플릿에 있고 출력 파일에 없는 이미지 확장자들 추가
        new_extensions = {k: v for k, v in template_extensions.items() if k not in existing_extensions}

        if new_extensions:
            # 출력 파일의 Content_Types.xml 업데이트
            with zipfile.ZipFile(output_path, 'r') as zip_ref:
                content_types_xml = zip_ref.read(content_types_file)

            root = etree.fromstring(content_types_xml)

            # 기존 Default들 찾기
            defaults = [child for child in root if child.tag.endswith('}Default') or child.tag == 'Default']

            # 새 Default 추가
            for ext, content_type in new_extensions.items():
                new_default = etree.Element('{http://schemas.openxmlformats.org/package/2006/content-types}Default')
                new_default.set('Extension', ext)
                new_default.set('ContentType', content_type)

                if defaults:
                    last_default = defaults[-1]
                    idx = list(root).index(last_default) + 1
                else:
                    idx = 0

                root.insert(idx, new_default)

            # ZIP 파일 수정
            temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
            try:
                os.close(temp_fd)

                with zipfile.ZipFile(output_path, 'r') as zin:
                    with zipfile.ZipFile(temp_path, 'w') as zout:
                        for item in zin.infolist():
                            data = zin.read(item.filename)
                            if item.filename == content_types_file:
                                zout.writestr(item, etree.tostring(root, encoding='UTF-8', xml_declaration=True, standalone=True))
                            else:
                                zout.writestr(item, data)

                shutil.move(temp_path, output_path)
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)

    def _add_content_types_for_extensions(self, output_path, extensions):
        """Content_Types.xml에 새로운 파일 확장자에 대한 Default 추가"""
        import zipfile
        from lxml import etree

        extension_to_mime = {
            'png': 'image/png',
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'gif': 'image/gif',
            'bmp': 'image/bmp',
            'tiff': 'image/tiff',
            'svg': 'image/svg+xml',
            'wmf': 'image/x-wmf',
            'emf': 'image/x-emf',
        }

        content_types_file = '[Content_Types].xml'

        with zipfile.ZipFile(output_path, 'r') as zip_ref:
            if content_types_file not in zip_ref.namelist():
                return

            content_types_xml = zip_ref.read(content_types_file)

        root = etree.fromstring(content_types_xml)

        existing_defaults = set()
        for child in root:
            if child.tag.endswith('}Default') or child.tag == 'Default':
                ext = child.get('Extension')
                if ext:
                    existing_defaults.add(ext.lower())

        new_defaults_added = False
        for ext in extensions:
            ext_lower = ext.lower()
            if ext_lower not in existing_defaults and ext_lower in extension_to_mime:
                defaults = [child for child in root if child.tag.endswith('}Default') or child.tag == 'Default']
                if defaults:
                    last_default = defaults[-1]
                    idx = list(root).index(last_default) + 1
                else:
                    idx = 0

                new_default = etree.Element('{http://schemas.openxmlformats.org/package/2006/content-types}Default')
                new_default.set('Extension', ext_lower)
                new_default.set('ContentType', extension_to_mime[ext_lower])
                root.insert(idx, new_default)
                existing_defaults.add(ext_lower)
                new_defaults_added = True

        if new_defaults_added:
            temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
            try:
                os.close(temp_fd)

                with zipfile.ZipFile(output_path, 'r') as zin:
                    with zipfile.ZipFile(temp_path, 'w') as zout:
                        for item in zin.infolist():
                            data = zin.read(item.filename)
                            if item.filename == content_types_file:
                                zout.writestr(item, etree.tostring(root, encoding='UTF-8', xml_declaration=True, standalone=True))
                            else:
                                zout.writestr(item, data)

                shutil.move(temp_path, output_path)
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)

    def _add_content_type_for_slide(self, output_path, slide_number):
        """Content_Types.xml에 슬라이드를 위한 Override 추가"""
        import zipfile
        from lxml import etree

        content_types_file = '[Content_Types].xml'
        slide_part_name = f'/ppt/slides/slide{slide_number}.xml'
        slide_content_type = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'

        with zipfile.ZipFile(output_path, 'r') as zip_ref:
            if content_types_file not in zip_ref.namelist():
                return
            content_types_xml = zip_ref.read(content_types_file)

        root = etree.fromstring(content_types_xml)

        existing = False
        for child in root:
            if child.tag.endswith('}Override') or child.tag == 'Override':
                part_name = child.get('PartName')
                if part_name == slide_part_name:
                    existing = True
                    break

        if not existing:
            overrides = [child for child in root if child.tag.endswith('}Override') or child.tag == 'Override']
            if overrides:
                last_override = overrides[-1]
                idx = list(root).index(last_override) + 1
            else:
                idx = 0

            new_override = etree.Element('{http://schemas.openxmlformats.org/package/2006/content-types}Override')
            new_override.set('PartName', slide_part_name)
            new_override.set('ContentType', slide_content_type)
            root.insert(idx, new_override)

            temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
            try:
                os.close(temp_fd)

                with zipfile.ZipFile(output_path, 'r') as zin:
                    with zipfile.ZipFile(temp_path, 'w') as zout:
                        for item in zin.infolist():
                            data = zin.read(item.filename)
                            if item.filename == content_types_file:
                                zout.writestr(item, etree.tostring(root, encoding='UTF-8', xml_declaration=True, standalone=True))
                            else:
                                zout.writestr(item, data)

                shutil.move(temp_path, output_path)
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)

    def _add_slide_to_presentation(self, output_path, slide_idx):
        """프레젠테이션에 새 슬라이드 추가 (presentation.xml과 presentation.xml.rels 동시 업데이트)"""
        import zipfile
        from lxml import etree

        # 1. 다음 사용 가능한 relationship ID 찾기
        pres_rel_file = 'ppt/_rels/presentation.xml.rels'
        with zipfile.ZipFile(output_path, 'r') as zip_ref:
            if pres_rel_file not in zip_ref.namelist():
                return
            rels_xml = zip_ref.read(pres_rel_file)

        rels_root = etree.fromstring(rels_xml)

        existing_rids = []
        for rel in rels_root:
            rid = rel.get('Id')
            if rid and rid.startswith('rId'):
                try:
                    num = int(rid[3:])
                    existing_rids.append(num)
                except ValueError:
                    pass

        if existing_rids:
            next_rid_num = max(existing_rids) + 1
        else:
            next_rid_num = 1

        new_rid = f'rId{next_rid_num}'

        # 2. presentation.xml에 새 슬라이드 ID 추가
        with zipfile.ZipFile(output_path, 'r') as zip_ref:
            pptx_xml = zip_ref.read('ppt/presentation.xml')

        ppt_root = etree.fromstring(pptx_xml)

        ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        }

        sld_id_lst = ppt_root.find('.//p:sldIdLst', ns)

        if sld_id_lst is not None:
            new_id = 256 + slide_idx + 1

            sld_id = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
            sld_id.set('id', str(new_id))
            sld_id.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', new_rid)

            sld_id_lst.append(sld_id)

        # 3. presentation.xml.rels에 새 relationship 추가
        new_target = f'slides/slide{slide_idx + 1}.xml'
        new_type = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'

        existing_rels = list(rels_root)
        if existing_rels:
            existing_tag = existing_rels[0].tag
            rel = etree.SubElement(rels_root, existing_tag)
        else:
            rel = etree.SubElement(rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')

        rel.set('Id', new_rid)
        rel.set('Target', new_target)
        rel.set('Type', new_type)

        # 4. ZIP 파일 수정
        temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
        try:
            os.close(temp_fd)

            with zipfile.ZipFile(output_path, 'r') as zin:
                with zipfile.ZipFile(temp_path, 'w') as zout:
                    for item in zin.infolist():
                        data = zin.read(item.filename)
                        if item.filename == 'ppt/presentation.xml':
                            zout.writestr(item, etree.tostring(ppt_root, encoding='UTF-8', xml_declaration=True, standalone=True))
                        elif item.filename == pres_rel_file:
                            zout.writestr(item, etree.tostring(rels_root, encoding='UTF-8', xml_declaration=True, standalone=True))
                        else:
                            zout.writestr(item, data)

            shutil.move(temp_path, output_path)
        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)

    def _replace_placeholders_in_slide(self, slide, placeholder_map: Dict[str, str]):
        """슬라이드 내의 모든 플레이스홀더 교체 (서식 유지)"""
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            text_frame = shape.text_frame

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text

                    for placeholder, value in placeholder_map.items():
                        if placeholder in text:
                            if text == placeholder:
                                self._replace_run_text_keep_format(run, value)
                            else:
                                new_text = text.replace(placeholder, value)
                                run.text = new_text

            # 빈 paragraph나 bullet만 있는 paragraph 삭제 (placeholder 교체 후 남은 빈 줄 제거)
            # 뒤에서부터 삭제하여 인덱스 문제 방지
            for i in range(len(text_frame.paragraphs) - 1, -1, -1):
                para = text_frame.paragraphs[i]
                para_text = para.text.strip()

                # paragraph가 비어있거나 bullet symbol만 있는 경우 삭제
                should_delete = False
                if not para_text:
                    should_delete = True
                elif para_text in ['•', '• ', '●', '● ', '-', '- ']:
                    # bullet symbol만 있는 경우
                    should_delete = True
                elif para_text == '내용':
                    # '내용' 단독으로 있는 경우만 삭제 (템플릿 placeholder)
                    # 실제 컨텐츠에 포함된 '내용'은 삭제되지 않도록 함
                    should_delete = True

                if should_delete:
                    try:
                        p = text_frame.paragraphs[i]
                        element = p._element
                        element.getparent().remove(element)
                    except Exception:
                        pass  # 삭제 실패 시 무시

    def _replace_run_text_keep_format(self, run, new_text: str):
        """run의 텍스트를 교체하며 서식 유지"""
        # python-pptx는 run.text를 설정해도 기존 서식을 유지함
        run.text = new_text

    def get_slide_count(self) -> int:
        """현재 슬라이드 수 반환"""
        return len(self.presentation.slides)

    def get_preview_info(self) -> Dict[str, Any]:
        """문서 미리보기 정보 반환"""
        return {
            'aspect_ratio': self.aspect_ratio.value,
            'slide_width': self.slide_width,
            'slide_height': self.slide_height,
            'slide_count': self.get_slide_count(),
            'template_dir': str(self.template_dir)
        }


# 편의 함수들
def create_pptx_document(aspect_ratio: str = '4_3', template_type: str = 'default') -> PPTXGenerator:
    """PPTX 문서 생성자 생성 (편의 함수)"""
    return PPTXGenerator(aspect_ratio=aspect_ratio, template_type=template_type)


def ensure_templates_exist():
    """템플릿이 존재하는지 확인하고, 없으면 생성"""
    airun_path_str = getVarVal('AIRUN_PATH', str(Path(__file__).parent))
    airun_path = Path(airun_path_str)
    template_dir = airun_path / "templates" / "pptx" / "default"

    if not template_dir.exists():
        print("[INFO] 템플릿 디렉토리를 생성합니다...")
        template_dir.mkdir(parents=True, exist_ok=True)

    ratio_templates = [
        f"title_templates_4_3.pptx",
        f"content_templates_4_3.pptx",
        f"title_templates_16_9.pptx",
        f"content_templates_16_9.pptx"
    ]

    missing_templates = []
    for template in ratio_templates:
        template_path = template_dir / template
        if not template_path.exists():
            missing_templates.append(template)

    if missing_templates:
        print(f"[INFO] 누락된 템플릿을 생성합니다: {missing_templates}")
        from scripts.create_pptx_templates import create_templates
        create_templates()
        print("[INFO] 템플릿 생성 완료")


if __name__ == "__main__":
    ensure_templates_exist()

    print("=" * 80)
    print("AIRUN PPTX Generator - 예제 문서 생성")
    print("=" * 80)

    for ratio in ['4_3', '16_9']:
        print(f"\n[{ratio} 비율]")

        pptx = create_pptx_document(aspect_ratio=ratio)

        pptx.add_title_slide(
            title="AIRUN PPTX Generator",
            subtitle="깔끔한 템플릿 기반 문서 생성"
        )

        pptx.add_section_slide(
            title="1. 시스템 개요",
            subtitle="주요 기능 소개"
        )

        content1 = SlideContent(
            title="주요 기능",
            content="• 템플릿 자동 관리\n• AI 레이아웃 자동 선택\n• 다양한 콘텐츠 지원"
        )
        pptx.add_content_slide(content1)

        image_path = "/home/hamonikr/문서/하모나이즈-인포그래픽.png"
        if os.path.exists(image_path):
            pptx.add_image_slide(
                title="시각적 표현",
                image_path=image_path
            )

        pptx.add_table_slide(
            title="기능 비교",
            headers=["기능", "지원 여부", "비고"],
            data=[
                ["템플릿 관리", "✓", "자동 생성"],
                ["AI 레이아웃", "✓", "자동 선택"],
                ["이미지 지원", "✓", "비율 유지"],
                ["표 지원", "✓", "스타일 자동"]
            ]
        )

        output_path = f"/home/hamonikr/문서/generated-pptx/airun_pptx_example_{ratio}.pptx"
        pptx.save(output_path)
