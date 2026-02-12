#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AIRUN PPTX Generator 전체 테스트
모든 레이아웃의 표, 이미지, 본문 기능 검증
"""

import sys
import os
import json
from pathlib import Path
from typing import Any
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# 현재 디렉토리를 sys.path 최우선으로 설정 (pv_solar 로컬 solar_pptx 사용)
current_dir = Path(__file__).resolve().parent
sys.path.insert(0, str(current_dir))
# 프로젝트 루트가 sys.path에 이미 있으면 제거해 로컬 모듈 충돌을 방지
project_root = str(current_dir.parent)
sys.path = [p for p in sys.path if p != project_root]
# 이전에 로드된 모듈이 있으면 제거 후 재로딩
sys.modules.pop('airun_pptx', None)
sys.modules.pop('solar_pptx', None)

from solar_pptx import PPTXGenerator
try:
    from capex_content_builder import (
        build_content_list_from_api,
        _build_capex_approval_values,
        _build_technical_solution_text,
    )
except ModuleNotFoundError:
    from pv_solar.capex_content_builder import (
        build_content_list_from_api,
        _build_capex_approval_values,
        _build_technical_solution_text,
    )


def _apply_approval_request_replacements(pptx_path: Path, replacements: dict[str, str], slide_index: int = 3) -> None:
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if slide_index >= len(prs.slides):
        return

    slide = prs.slides[slide_index]
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            continue
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                for target, value in replacements.items():
                    if target in text:
                        run.text = text.replace(target, value)

    prs.save(str(pptx_path))


def prune_slides(pptx_path: Path, keep: int) -> None:
    """생성 후 슬라이드 개수를 기대치에 맞춰 정리한다."""

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if len(prs.slides) <= keep:
        return

    sld_id_lst = prs.slides._sldIdLst
    while len(sld_id_lst) > keep:
        sld_id_lst.remove(sld_id_lst[-1])

    prs.save(str(pptx_path))
    print(f"⚙️  불필요 슬라이드 제거: 총 {len(prs.slides)}장으로 정리")


def apply_main_agreements_slide(pptx_path: Path,
                                slide_idx: int,
                                title_text: str,
                                table_headers: list[str],
                                table_rows: list[list[str]]) -> None:
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if slide_idx >= len(prs.slides):
        return

    slide = prs.slides[slide_idx]
    if not table_rows:
        prs.save(str(pptx_path))
        return

    # 템플릿에 placeholder/컬럼 헤더가 없는 경우 조용히 스킵
    template_texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            template_texts.append(str(getattr(shape, "text", "") or "").strip())
    joined_template_text = "\n".join(template_texts)
    has_placeholder_token = any(token in joined_template_text for token in ("{{agreement_", "{agreement_", "<<agreement_"))
    has_column_headers = ("Current Status" in joined_template_text) and ("Next Steps" in joined_template_text)
    if not has_placeholder_token and not has_column_headers:
        prs.save(str(pptx_path))
        return

    default_text = "입력"

    def _normalize_value(value: Any, *, fill_default: bool = False) -> str:
        if value is None:
            return default_text if fill_default else ""
        text = str(value).strip()
        if not text:
            return default_text if fill_default else ""
        return text

    def build_replacements(rows: list[list[str]]) -> dict[str, str]:
        keys = ["agreement", "completion", "current_status", "next_steps"]
        replacements: dict[str, str] = {}
        for idx, row in enumerate(rows, start=1):
            for key, value in zip(keys, row):
                val = _normalize_value(
                    value,
                    fill_default=(key in {"current_status", "next_steps"}),
                )
                tokens = [
                    f"{{{{{key}_{idx}}}}}",
                    f"{{{key}_{idx}}}",
                    f"<<{key}_{idx}>>",
                ]
                for token in tokens:
                    replacements[token] = val
        return replacements

    replacements = build_replacements(table_rows)


    def replace_in_text_frame(text_frame, mapping: dict[str, str]) -> bool:
        changed = False
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                if not text:
                    continue
                new_text = text
                for target, value in mapping.items():
                    if target in new_text:
                        new_text = new_text.replace(target, value)
                if new_text != text:
                    run.text = new_text
                    changed = True
        return changed

    replaced_any = False
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            continue
        if replace_in_text_frame(text_frame, replacements):
            replaced_any = True

    def _find_header_shape(text: str) -> Any | None:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if str(getattr(shape, "text", "") or "").strip() == text:
                return shape
        return None

    def _collect_column_shapes(header_shape: Any | None,
                               min_width: int = 2000000,
                               min_height: int = 700000,
                               max_height: int = 1400000) -> list[Any]:
        if not header_shape:
            return []
        header_center = float(header_shape.left) + float(header_shape.width) / 2.0
        header_bottom = int(header_shape.top) + int(header_shape.height)
        candidates: list[tuple[float, Any]] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if int(shape.width) < min_width:
                continue
            if int(shape.height) < min_height or int(shape.height) > max_height:
                continue
            if int(shape.top) <= header_bottom:
                continue
            center = float(shape.left) + float(shape.width) / 2.0
            candidates.append((abs(center - header_center), shape))
        if not candidates:
            return []
        best_left = int(min(candidates, key=lambda x: x[0])[1].left)
        return [shape for _, shape in candidates if int(shape.left) == best_left]

    def _fill_column(shapes: list[Any], values: list[str], fill_missing: str | None = None) -> bool:
        if not shapes:
            return False
        ordered = sorted(shapes, key=lambda s: int(s.top))
        extended = list(values)
        if fill_missing is not None and len(extended) < len(ordered):
            extended.extend([fill_missing] * (len(ordered) - len(extended)))
        filled = False
        for shape, value in zip(ordered, extended):
            if value is None:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            shape.text_frame.text = value
            filled = True
        return filled

    current_header = _find_header_shape("Current Status")
    next_header = _find_header_shape("Next Steps")
    current_shapes = _collect_column_shapes(current_header)
    next_shapes = _collect_column_shapes(next_header)

    current_values = [
        _normalize_value(row[2], fill_default=True) if len(row) > 2 else default_text
        for row in table_rows
    ]
    next_values = [
        _normalize_value(row[3], fill_default=True) if len(row) > 3 else default_text
        for row in table_rows
    ]

    filled_current = _fill_column(current_shapes, current_values, fill_missing=default_text)
    filled_next = _fill_column(next_shapes, next_values, fill_missing=default_text)

    if not replaced_any and not (filled_current or filled_next):
        print("[INFO] Main Agreements template placeholders were not matched; skipped text replacement.")

    # 요청사항: 제목을 제외한 표/본문 폰트는 9pt로 고정
    title_norm = (title_text or "").strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            continue
        shape_text = str(getattr(shape, "text", "") or "").strip().lower()
        if shape_text == title_norm:
            continue
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(9)

    prs.save(str(pptx_path))


def apply_cod_pipeline_slide(pptx_path: Path,
                             slide_idx: int,
                             title_text: str,
                             technical_solution_text: str,
                             table_title_1: str,
                             subtitle_2: str) -> None:
    """COD pipeline 슬라이드의 템플릿 잔여 placeholder 텍스트를 정리한다."""

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if slide_idx >= len(prs.slides):
        return

    slide = prs.slides[slide_idx]

    technical_shape = None
    subtitle1_shape = None
    subtitle2_shape = None
    title_shape = None
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        tf = getattr(shape, "text_frame", None)
        if tf is None:
            continue
        text = str(getattr(shape, "text", "") or "").strip()
        if not text:
            continue

        # 제목은 content 기준으로 재확인
        if text == title_text or "BOLT#2 is Korean PV Portfolio" in text:
            title_shape = shape
            tf.text = title_text
            continue

        # 서브타이틀 1/2 채우기
        compact = text.replace("\xa0", "")
        if "서비타이틀1" in compact or "서비스타이틀1" in compact or "서브타이틀1" in compact:
            subtitle1_shape = shape
            tf.text = table_title_1
            continue
        if "서브타이틀2" in compact:
            subtitle2_shape = shape
            tf.text = subtitle_2
            continue

        # 서브타이틀2 하단 본문 placeholder를 실제 본문으로 치환
        if "텍스트 영역" in text or "왼쪽" in text:
            technical_shape = shape
            tf.clear()

    is_construction_notes = subtitle_2.strip().lower() == "construction notes"

    if technical_shape is not None:
        tf = getattr(technical_shape, "text_frame", None)
        if tf is None:
            prs.save(str(pptx_path))
            return
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        lines = [ln.strip() for ln in (technical_solution_text or "").splitlines() if ln.strip()]

        for idx, line in enumerate(lines):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = 1.15 if is_construction_notes else 1.0
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(2 if is_construction_notes else 1)

            run = paragraph.add_run()
            if is_construction_notes and line.lower().startswith(("engineering", "procurement", "construction")):
                run.text = line
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0, 79, 153)
            elif line == "Equipment Configuration Summary":
                run.text = line
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0, 0, 0)
            elif line.startswith("-"):
                run.text = f"▪  {line.lstrip('- ').strip()}"
                run.font.bold = False
                run.font.size = Pt(9 if is_construction_notes else 10)
                run.font.color.rgb = RGBColor(56, 56, 56)
            else:
                run.text = line
                run.font.bold = True
                run.font.size = Pt(10 if is_construction_notes else 12)
                run.font.color.rgb = RGBColor(0, 0, 0)

    # 레이아웃 위치/크기 보정 (레퍼런스 이미지 비율 반영)
    if title_shape is not None:
        title_shape.left = Inches(0.2)
        title_shape.top = Inches(0.2)
        title_shape.width = Inches(12.9)
        title_shape.height = Inches(0.5)
        tf = getattr(title_shape, "text_frame", None)
        if tf is not None and tf.paragraphs:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = RGBColor(68, 114, 196)
    if subtitle1_shape is not None:
        subtitle1_shape.left = Inches(0.2)
        subtitle1_shape.top = Inches(0.72)
        subtitle1_shape.width = Inches(3.4)
        subtitle1_shape.height = Inches(0.28)
        tf = getattr(subtitle1_shape, "text_frame", None)
        if tf is not None and tf.paragraphs and tf.paragraphs[0].runs:
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(14)
            run.font.bold = False
            run.font.color.rgb = RGBColor(96, 68, 42)
    if subtitle2_shape is not None:
        subtitle2_shape.left = Inches(7.2)
        subtitle2_shape.top = Inches(0.72)
        subtitle2_shape.width = Inches(3.6)
        subtitle2_shape.height = Inches(0.28)
        tf = getattr(subtitle2_shape, "text_frame", None)
        if tf is not None and tf.paragraphs and tf.paragraphs[0].runs:
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(14)
            run.font.bold = False
            run.font.color.rgb = RGBColor(96, 68, 42)
    if technical_shape is not None:
        technical_shape.left = Inches(7.2)
        technical_shape.top = Inches(1.05)
        technical_shape.width = Inches(5.9)
        technical_shape.height = Inches(6.2)

    def _style_cod_pipeline_table(table) -> None:
        import lxml.etree as etree

        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        border_color = '2A2A2A'
        header_bg_color = '0F70B7'
        body_bg_color = 'D9D9D9'
        cod_col_bg_color = 'FFFFFF'
        subtotal_bg_color = 'D9D9D9'
        total_bg_color = 'FFF200'

        def set_cell_border(cell_element, width: int = 9000):
            tc_pr = cell_element.find('./a:tcPr', namespaces=ns)
            if tc_pr is None:
                tc_pr = etree.SubElement(cell_element, '{%s}tcPr' % ns['a'])
            tc_borders = tc_pr.find('./a:tcBorders', namespaces=ns)
            if tc_borders is None:
                tc_borders = etree.SubElement(tc_pr, '{%s}tcBorders' % ns['a'])
            for side in ['top', 'left', 'bottom', 'right']:
                border = tc_borders.find('./a:%s' % side, namespaces=ns)
                if border is None:
                    border = etree.SubElement(tc_borders, '{%s}%s' % (ns['a'], side))
                border.set('val', 'single')
                border.set('w', str(width))
                border.set('color', border_color)

        def set_cell_fill(cell_element, fill_color: str):
            tc_pr = cell_element.find('./a:tcPr', namespaces=ns)
            if tc_pr is None:
                tc_pr = etree.SubElement(cell_element, '{%s}tcPr' % ns['a'])
            for existing in tc_pr.findall('./a:solidFill', namespaces=ns):
                tc_pr.remove(existing)
            solid_fill = etree.SubElement(tc_pr, '{%s}solidFill' % ns['a'])
            srgb_clr = etree.SubElement(solid_fill, '{%s}srgbClr' % ns['a'])
            srgb_clr.set('val', fill_color)

        total_width = int(Inches(6.8))
        col_count = len(table.columns)
        if col_count == 7:
            ratios = [0.10, 0.09, 0.06, 0.30, 0.17, 0.13, 0.15]
        elif col_count == 8:
            ratios = [0.08, 0.08, 0.06, 0.28, 0.14, 0.12, 0.12, 0.12]
        else:
            ratios = [1 / max(1, col_count)] * col_count
        for cidx, ratio in enumerate(ratios[:col_count]):
            table.columns[cidx].width = int(total_width * ratio)

        total_height = sum(int(row.height) for row in table.rows)
        if total_height > 0 and len(table.rows) > 1:
            header_h = min(int(0.52 * 914400), int(total_height * 0.18))
            body_h = int((total_height - header_h) / (len(table.rows) - 1))
            table.rows[0].height = header_h
            for ridx in range(1, len(table.rows)):
                table.rows[ridx].height = body_h

        for cidx in range(len(table.columns)):
            cell = table.cell(0, cidx)
            tf_h = cell.text_frame
            tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
            if tf_h.paragraphs:
                p = tf_h.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.line_spacing = 1.0
                if p.runs:
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(9)
                        run.font.color.rgb = RGBColor(255, 255, 255)
            set_cell_fill(cell._tc, header_bg_color)
            set_cell_border(cell._tc)

        label_col_idx = 3 if col_count >= 4 else max(0, col_count - 1)
        subtotal_rows: list[int] = []
        total_row_idx = None
        for ridx in range(1, len(table.rows)):
            label = (table.cell(ridx, label_col_idx).text or '').strip().lower().replace(' ', '')
            if label == 'subtotal':
                subtotal_rows.append(ridx)
            if label == 'total':
                total_row_idx = ridx

            no_value = (table.cell(ridx, 2).text or '').strip()
            is_numbered_project_row = no_value.isdigit()

            for cidx in range(len(table.columns)):
                cell = table.cell(ridx, cidx)
                tf = cell.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    center_cols = {0, 1, 2}
                    if col_count == 7:
                        center_cols |= {5, 6}
                    elif col_count == 8:
                        center_cols |= {4, 5, 6, 7}
                    if cidx in center_cols:
                        p.alignment = PP_ALIGN.CENTER
                    else:
                        p.alignment = PP_ALIGN.LEFT
                    if p.runs:
                        for run in p.runs:
                            run.font.size = Pt(9)
                            run.font.color.rgb = RGBColor(74, 57, 43)
                if is_numbered_project_row and 2 <= cidx <= (col_count - 1):
                    set_cell_fill(cell._tc, 'FFFFFF')
                elif (col_count == 7 and cidx == 6) or (col_count == 8 and cidx in {6, 7}):
                    set_cell_fill(cell._tc, cod_col_bg_color)
                else:
                    set_cell_fill(cell._tc, body_bg_color)
                set_cell_border(cell._tc)

        phase1_start = 1
        phase1_end = subtotal_rows[0] if subtotal_rows else (total_row_idx - 1 if total_row_idx else len(table.rows) - 1)
        if total_row_idx is not None:
            spv_end = total_row_idx
        else:
            spv_end = subtotal_rows[-1] if subtotal_rows else phase1_end

        if spv_end >= phase1_start and col_count >= 1:
            spv_cell = table.cell(phase1_start, 0)
            spv_cell.merge(table.cell(spv_end, 0))
            spv_cell.text = spv_cell.text or 'BOLT#2'
            spv_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            if spv_cell.text_frame.paragraphs:
                p = spv_cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    p.runs[0].font.bold = True
                    p.runs[0].font.size = Pt(10)
            set_cell_fill(spv_cell._tc, '3DDC84')
            set_cell_border(spv_cell._tc)

        if phase1_end >= phase1_start and col_count >= 2:

            phase1_cell = table.cell(phase1_start, 1)
            phase1_cell.merge(table.cell(phase1_end, 1))
            phase1_cell.text = '1'
            phase1_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            if phase1_cell.text_frame.paragraphs:
                p = phase1_cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    p.runs[0].font.bold = True
                    p.runs[0].font.size = Pt(10)
            set_cell_fill(phase1_cell._tc, body_bg_color)
            set_cell_border(phase1_cell._tc)

            cod_col_idx = 6 if col_count == 7 else None
            if cod_col_idx is not None and cod_col_idx < col_count:
                cod1_cell = table.cell(phase1_start, cod_col_idx)
                cod1_cell.merge(table.cell(phase1_end, cod_col_idx))
                cod1_cell.text = cod1_cell.text or table.cell(phase1_start, cod_col_idx).text or ''
                cod1_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if cod1_cell.text_frame.paragraphs:
                    p = cod1_cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    if p.runs:
                        p.runs[0].font.bold = True
                        p.runs[0].font.size = Pt(9)

        if len(subtotal_rows) >= 2 and col_count >= 2:
            phase2_start = subtotal_rows[0] + 1
            phase2_end = subtotal_rows[1]
            phase2_cell = table.cell(phase2_start, 1)
            phase2_cell.merge(table.cell(phase2_end, 1))
            phase2_cell.text = '2'
            phase2_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            if phase2_cell.text_frame.paragraphs:
                p = phase2_cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    p.runs[0].font.bold = True
                    p.runs[0].font.size = Pt(10)

            cod_col_idx = 6 if col_count == 7 else None
            if cod_col_idx is not None and cod_col_idx < col_count:
                cod2_cell = table.cell(phase2_start, cod_col_idx)
                cod2_cell.merge(table.cell(phase2_end, cod_col_idx))
                cod2_cell.text = cod2_cell.text or table.cell(phase2_start, cod_col_idx).text or ''
                cod2_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                if cod2_cell.text_frame.paragraphs:
                    p = cod2_cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    if p.runs:
                        p.runs[0].font.bold = True
                        p.runs[0].font.size = Pt(9)

        for ridx in subtotal_rows:
            subtotal_start = 2 if col_count >= 4 else 0
            subtotal_end = 3 if col_count >= 4 else subtotal_start
            subtotal_cell = table.cell(ridx, subtotal_start)
            subtotal_cell.merge(table.cell(ridx, subtotal_end))
            subtotal_cell.text = 'Sub Total'
            subtotal_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            if subtotal_cell.text_frame.paragraphs:
                p = subtotal_cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    p.runs[0].font.bold = True
                    p.runs[0].font.size = Pt(10)
            for cidx in range(0, col_count):
                cell = table.cell(ridx, cidx)
                set_cell_fill(cell._tc, subtotal_bg_color)
                set_cell_border(cell._tc)
                if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].font.bold = True
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)

        if total_row_idx is not None and col_count >= 2:
            total_start = 1
            total_end = 4 if col_count >= 5 else (col_count - 1)
            total_label_cell = table.cell(total_row_idx, total_start)
            total_label_cell.merge(table.cell(total_row_idx, total_end))
            total_label_cell.text = 'TOTAL'
            total_label_cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            if total_label_cell.text_frame.paragraphs:
                p = total_label_cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    p.runs[0].font.bold = True
                    p.runs[0].font.size = Pt(11)
            for cidx in range(0, col_count):
                cell = table.cell(total_row_idx, cidx)
                set_cell_fill(cell._tc, total_bg_color)
                set_cell_border(cell._tc)
                if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].font.bold = True
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

        # 요청사항: 표 내 폰트 사이즈를 전부 9pt로 통일
        for row in table.rows:
            for cell in row.cells:
                tf = cell.text_frame
                for paragraph in tf.paragraphs:
                    if paragraph.runs:
                        for run in paragraph.runs:
                            run.font.size = Pt(9)

    # COD Pipeline 표는 전 셀 폰트를 9pt로 통일
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        shape.left = Inches(0.2)
        shape.top = Inches(1.05)
        shape.width = Inches(6.8)
        shape.height = Inches(6.0)
        table = getattr(shape, "table", None)
        if table is None:
            continue
        _style_cod_pipeline_table(table)

    prs.save(str(pptx_path))


def apply_equipment_procurement_title(pptx_path: Path, slide_idx: int, title_text: str, body_text: str = "") -> None:
    """템플릿 23 기반 슬라이드의 타이틀을 강제로 적용한다."""
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if slide_idx >= len(prs.slides):
        return

    slide = prs.slides[slide_idx]

    title_shape = None
    # 1) 기존 텍스트 패턴으로 우선 탐색
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = str(getattr(shape, "text", "") or "").strip()
        if not text:
            continue
        if "Equipment procurement process" in text or "Illustrative case" in text:
            title_shape = shape
            break

    # 2) 못 찾으면 상단 텍스트 박스 중 가장 큰 것을 타이틀로 간주
    if title_shape is None:
        candidates = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if getattr(shape, "top", 0) <= Inches(1.2):
                area = int(getattr(shape, "width", 0)) * int(getattr(shape, "height", 0))
                candidates.append((area, shape))
        if candidates:
            title_shape = sorted(candidates, key=lambda x: x[0], reverse=True)[0][1]

    if title_shape is not None:
        tf = getattr(title_shape, "text_frame", None)
        if tf is None:
            prs.save(str(pptx_path))
            return
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.color.rgb = RGBColor(68, 114, 196)

        # 위치도 타이틀 슬라이드들과 유사하게 정렬
        title_shape.left = Inches(0.2)
        title_shape.top = Inches(0.2)
        title_shape.width = Inches(12.9)
        title_shape.height = Inches(0.5)
    else:
        # 타이틀 shape가 없으면 새로 생성
        title_shape = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(12.9), Inches(0.5))
        tf = title_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.color.rgb = RGBColor(68, 114, 196)

    # 본문 불릿도 주입 (하드코딩 문구 대신 content 사용)
    if body_text:
        body_shape = None
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = str(getattr(shape, "text", "") or "").strip()
            if not text:
                continue
            if text.startswith("•") or text.startswith("-") or "CLUSTER SET" in text:
                body_shape = shape
                break

        if body_shape is None:
            body_shape = slide.shapes.add_textbox(Inches(0.24), Inches(0.95), Inches(12.8), Inches(1.55))

        tf_body = getattr(body_shape, "text_frame", None)
        if tf_body is not None:
            tf_body.clear()
            tf_body.word_wrap = True
            tf_body.vertical_anchor = MSO_ANCHOR.TOP
            lines = [ln.strip() for ln in body_text.splitlines() if ln.strip()]
            for idx, line in enumerate(lines):
                p = tf_body.paragraphs[0] if idx == 0 else tf_body.add_paragraph()
                p.text = f"•  {line.lstrip('- ').strip()}"
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = 1.0
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                if p.runs:
                    run = p.runs[0]
                    run.font.size = Pt(8.8)
                    run.font.bold = False
                    run.font.color.rgb = RGBColor(0, 0, 0)

            body_shape.left = Inches(0.24)
            body_shape.top = Inches(0.95)
            body_shape.width = Inches(12.8)
            body_shape.height = Inches(1.55)

    prs.save(str(pptx_path))


def _build_project_detail_values(data_dir: Path) -> dict[str, Any]:
    general_path = data_dir / "general.json"
    opex_path = data_dir / "opex_year1.json"

    def _as_dict(value: Any) -> dict[str, Any]:
        return value if isinstance(value, dict) else {}

    try:
        general = json.loads(general_path.read_text(encoding="utf-8"))
    except Exception:
        general = {}
    try:
        opex = json.loads(opex_path.read_text(encoding="utf-8"))
    except Exception:
        opex = {}

    vendor_path = data_dir / "equipment_vendors.json"
    try:
        vendor_overrides = json.loads(vendor_path.read_text(encoding="utf-8"))
    except Exception:
        vendor_overrides = {}

    projects = general.get("projects") if isinstance(general, dict) else []
    if not isinstance(projects, list):
        projects = []
    cluster_projects = [p for p in projects[:4] if isinstance(p, dict)]

    regions: list[str] = []
    cod_dates: list[datetime] = []
    dc_wp = 0.0
    names: set[str] = set()
    for p in cluster_projects:
        names.add(str(p.get("name") or "").strip())
        gi = _as_dict(p.get("general_inputs"))
        region = str(gi.get("Region") or p.get("region") or "").strip()
        if region and region not in regions:
            regions.append(region)
        cap_mw = gi.get("Total Capacity DC")
        try:
            dc_wp += float(str(cap_mw).replace(",", "")) * 1000
        except Exception:
            pass
        cod_raw = gi.get("COD date")
        if cod_raw:
            txt = str(cod_raw).strip().replace(" ", "T")
            try:
                cod_dates.append(datetime.fromisoformat(txt))
            except Exception:
                try:
                    cod_dates.append(datetime.strptime(txt.split("T")[0], "%Y-%m-%d"))
                except Exception:
                    pass

    # 요청사항: LOCATION 값은 고정 문구 사용
    location_text = "see the LOCATION"
    if cod_dates:
        min_cod = min(cod_dates)
        quarter = ((min_cod.month - 1) // 3) + 1
        cod_text = f"Q{quarter} {str(min_cod.year)[-2:]}"
    else:
        cod_text = "TBD"

    # quote/budget from opex costs for first 4 projects
    quote = 0.0
    contingency = 0.0
    grid_connection_total = 0.0
    opex_projects = opex.get("projects") if isinstance(opex, dict) else []
    if not isinstance(opex_projects, list):
        opex_projects = []
    for p in opex_projects:
        if not isinstance(p, dict):
            continue
        if str(p.get("name") or "").strip() not in names:
            continue
        oy = _as_dict(p.get("opex_year1"))
        def _cost(key: str) -> float:
            block = _as_dict(oy.get(key))
            if block:
                try:
                    return float(str(block.get("Cost") or "0").replace(",", ""))
                except Exception:
                    return 0.0
            return 0.0
        grid_cost = _cost("Grid Connection Cost")
        quote += _cost("Modules") + _cost("BOS") + grid_cost
        grid_connection_total += grid_cost
        contingency += _cost("Contingency")

    budget = quote + contingency
    delta = budget - quote
    total_wp = dc_wp * 1000.0
    cubicle_won_per_wp = (grid_connection_total / total_wp) if total_wp > 0 else 0.0
    cubicle_amount_for_1000kwp = cubicle_won_per_wp * 1_000_000

    cubicle_companies_raw = (
        ((vendor_overrides.get("cluster_set_1_4") or {}).get("cubicle") or {}).get("companies")
        if isinstance(vendor_overrides, dict)
        else None
    )
    cubicle_companies: list[str] = []
    if isinstance(cubicle_companies_raw, list):
        for item in cubicle_companies_raw:
            text = str(item or "").strip()
            if text:
                cubicle_companies.append(text)

    return {
        "location": location_text,
        "cod": cod_text,
        "dc_wp": f"{dc_wp:,.0f}" if dc_wp > 0 else "",
        "total_quote": f"{quote:,.0f}" if quote > 0 else "",
        "equip_budget": f"{budget:,.0f}" if budget > 0 else "",
        "delta": f"{delta:,.0f}" if delta > 0 else "",
        "cubicle_amount": f"{cubicle_amount_for_1000kwp:,.0f}" if cubicle_amount_for_1000kwp > 0 else "",
        "cubicle_won_per_wp": f"{cubicle_won_per_wp:,.1f}" if cubicle_won_per_wp > 0 else "",
        "cubicle_companies": cubicle_companies,
    }


def apply_equipment_project_detail_values(pptx_path: Path, slide_idx: int, data_dir: Path) -> None:
    """슬라이드 14 좌측 PROJECT DETAIL 영역 값(수치)을 data/*.json 기반으로 주입."""
    values = _build_project_detail_values(data_dir)

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    if slide_idx >= len(prs.slides):
        return
    slide = prs.slides[slide_idx]

    label_shapes: dict[str, Any] = {}
    budget_shapes: list[Any] = []
    amount_header_shapes: list[Any] = []
    won_header_shapes: list[Any] = []
    company_header_shapes: list[Any] = []
    cubicle_header_shape = None
    location_shape = None
    cod_shape = None
    cluster_shape = None

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = str(getattr(shape, "text", "") or "").strip()
        norm = text.lower().replace(" ", "")
        if norm == "location":
            label_shapes["location"] = shape
        elif norm == "cod":
            label_shapes["cod"] = shape
        elif norm == "dcwp":
            label_shapes["dc_wp"] = shape
        elif "totalquote" in norm:
            label_shapes["total_quote"] = shape
        elif "equip.budget" in norm:
            label_shapes["equip_budget"] = shape
        elif norm in {"δ+", "d+", "a+", "△+", "Δ+".lower()} or "δ+" in norm or "Δ+" in text:
            label_shapes["delta"] = shape
        elif "see the location" in text.lower():
            location_shape = shape
        elif text.lower().startswith("q") and len(text) <= 6:
            cod_shape = shape
        elif "cluster set" in text.lower():
            cluster_shape = shape
        elif text.strip() == "Budget":
            budget_shapes.append(shape)
        elif text.strip() == "Amount(\\)":
            amount_header_shapes.append(shape)
        elif text.strip() == "Won(\\)/wp":
            won_header_shapes.append(shape)
        elif text.strip() == "Company":
            company_header_shapes.append(shape)

        if "CUBICLE [1200kVA]" in text and "PV MONITORING" in text:
            cubicle_header_shape = shape

    def _set_shape_text(shape: Any, value: str, size: float = 9.0):
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            return
        shape.text = value
        if not tf.paragraphs:
            return
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(size)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0, 0, 0)

    if location_shape is not None:
        _set_shape_text(location_shape, values["location"])
    if cod_shape is not None:
        _set_shape_text(cod_shape, values["cod"])
    if cluster_shape is not None:
        _set_shape_text(cluster_shape, "CLUSTER SET 1 ~ 4")

    # LOCATION/COD는 기본 값 shape가 없을 때만 라벨 우측 값 칸 생성
    for key in ["location", "cod"]:
        if key == "location" and location_shape is not None:
            continue
        if key == "cod" and cod_shape is not None:
            continue
        label = label_shapes.get(key)
        if label is None:
            continue
        left = label.left + label.width + Inches(0.02)
        top = label.top
        width = Inches(1.9)
        height = label.height
        val_shape = slide.shapes.add_textbox(left, top, width, height)
        _set_shape_text(val_shape, values[key], 9.0)

    # Numeric rows often have empty value boxes; place overlay text right to each label
    for key in ["dc_wp", "total_quote", "equip_budget", "delta"]:
        label = label_shapes.get(key)
        if label is None:
            continue
        left = label.left + label.width + Inches(0.02)
        top = label.top
        width = Inches(1.35)
        height = label.height
        val_shape = slide.shapes.add_textbox(left, top, width, height)
        _set_shape_text(val_shape, values[key], 9.0)

    # CUBICLE [1200kVA]+'DER-AVM'+PV MONITORING budget row values (data-driven)
    if cubicle_header_shape is not None:
        cubicle_left = cubicle_header_shape.left
        cubicle_right = cubicle_header_shape.left + cubicle_header_shape.width

        def _center_x(shape: Any) -> int:
            return int(shape.left + (shape.width / 2))

        cubicle_budget = next(
            (
                shp
                for shp in sorted(budget_shapes, key=lambda s: s.top)
                if cubicle_left <= _center_x(shp) <= cubicle_right
            ),
            None,
        )
        cubicle_amount_header = next(
            (
                shp
                for shp in amount_header_shapes
                if cubicle_left <= _center_x(shp) <= cubicle_right
            ),
            None,
        )
        cubicle_won_header = next(
            (
                shp
                for shp in won_header_shapes
                if cubicle_left <= _center_x(shp) <= cubicle_right
            ),
            None,
        )
        cubicle_company_header = next(
            (
                shp
                for shp in company_header_shapes
                if cubicle_left <= _center_x(shp) <= cubicle_right
            ),
            None,
        )

        if cubicle_budget is not None:
            amount_text = values.get("cubicle_amount", "")
            won_text = values.get("cubicle_won_per_wp", "")

            if cubicle_amount_header is not None and amount_text:
                amount_shape = slide.shapes.add_textbox(
                    cubicle_amount_header.left,
                    cubicle_budget.top,
                    cubicle_amount_header.width,
                    cubicle_budget.height,
                )
                _set_shape_text(amount_shape, amount_text, 9.0)

            if cubicle_won_header is not None and won_text:
                won_shape = slide.shapes.add_textbox(
                    cubicle_won_header.left,
                    cubicle_budget.top,
                    cubicle_won_header.width,
                    cubicle_budget.height,
                )
                _set_shape_text(won_shape, won_text, 9.0)

        company_values = values.get("cubicle_companies")
        if cubicle_company_header is not None and isinstance(company_values, list) and company_values:
            min_top = cubicle_company_header.top + int(Inches(0.10))
            max_top = cubicle_company_header.top + int(Inches(0.80))
            skip_texts = {"Budget", "Amount(\\)", "Won(\\)/wp", "Company"}
            cubicle_company_cells = [
                shp
                for shp in slide.shapes
                if getattr(shp, "has_text_frame", False)
                and cubicle_left <= _center_x(shp) <= cubicle_right
                and min_top <= shp.top <= max_top
                and str(getattr(shp, "text", "") or "").strip() not in skip_texts
            ]
            cubicle_company_cells.sort(key=lambda s: s.top)

            for idx, name in enumerate(company_values):
                if idx >= len(cubicle_company_cells):
                    break
                _set_shape_text(cubicle_company_cells[idx], str(name), 8.0)

    # 요청사항: 슬라이드 14 표 영역 폰트 크기 6pt
    table_top_threshold = Inches(2.55)
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "top", 0) < table_top_threshold:
            continue
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            continue
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(6)
                run.font.bold = False
                run.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(str(pptx_path))


def _resolve_capex_data_dir() -> Path:
    env_data_dir = os.getenv("CAPEX_DATA_DIR", "").strip()
    if env_data_dir:
        return Path(env_data_dir).resolve()
    return (Path(__file__).resolve().parent.parent / "data").resolve()


def _resolve_capex_output_path(default_output_dir: Path) -> Path:
    env_output_path = os.getenv("CAPEX_OUTPUT_PATH", "").strip()
    if env_output_path:
        return Path(env_output_path).resolve()
    env_output_dir = os.getenv("CAPEX_OUTPUT_DIR", "").strip()
    output_dir = Path(env_output_dir).resolve() if env_output_dir else default_output_dir.resolve()
    return output_dir / "result.pptx"


def test_airun_pptx():
    """AIRUN PPTX Generator 전체 기능 테스트"""
    print("=" * 80)
    print("AIRUN PPTX Generator 전체 테스트")
    print("=" * 80)

    # 출력 디렉토리 생성 (workspace 내 고정)
    data_dir = _resolve_capex_data_dir()
    output_dir = Path(__file__).resolve().parent.parent / "resuot_output"
    output_path = _resolve_capex_output_path(output_dir)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # v2 방식: content_list에 모든 슬라이드 내용을 담기 (API 연동)
    content_list = build_content_list_from_api(data_dir=data_dir)
    
    # ========== 문서 생성 ==========

    # PPTXGenerator 생성
    gen = PPTXGenerator(aspect_ratio='16_9', template_type='custom')
    gen.generate_with_template(content_list, str(output_path))

    approval_values = _build_capex_approval_values(data_dir)
    _apply_approval_request_replacements(
        output_path,
        {
            "KRW 5,086M": approval_values["capex_krw_text"],
            "EUR 3,459,727": approval_values["capex_eur_text"],
            "4.203MWp": approval_values["capacity_text"],
        },
    )

    # 세션 타이틀 슬라이드 재구성 (content 적용 보장)
    def _normalized_order(val: Any):
        if val is None:
            return None
        if isinstance(val, (int, float)):
            return float(val)
        if isinstance(val, str):
            try:
                return float(val)
            except ValueError:
                return None
        return None

    ordered = []
    for idx, item in enumerate(content_list):
        norm = _normalized_order(item.order)
        ordered.append((idx, item, norm))

    ordered_content = [
        item for _, item, _ in sorted(
            ordered,
            key=lambda x: (
                0 if x[2] is not None else 1,
                x[2] if x[2] is not None else x[0]
            )
        )
    ]

    for idx, c in enumerate(ordered_content):
        # 템플릿 스타일 유지를 위해 대부분 레이아웃은 generate_with_template 결과를 그대로 사용.
        # 일부 레이아웃은 placeholder 구조가 특수하여 별도 후처리를 수행한다.
        if c.layout == "cod_pipeline":
            legacy_table_title = (
                (c.content or "").strip()
                if c.content and "pipeline" in c.content.lower() and "\n" not in c.content
                else ""
            )
            table_title_1 = c.subtitle or legacy_table_title or "COD 2026 Pipeline"
            technical_solution_text = (
                c.content
                if c.content and c.content != table_title_1
                else _build_technical_solution_text(data_dir)
            )
            apply_cod_pipeline_slide(
                output_path,
                idx,
                c.title,
                technical_solution_text,
                table_title_1,
                "Technical Solution",
            )

        if c.layout == "cod_pipeline_phase":
            table_title_1 = c.subtitle or "Construction baseline summary"
            apply_cod_pipeline_slide(
                output_path,
                idx,
                c.title,
                c.content or "",
                table_title_1,
                "Construction Notes",
            )

        if c.layout == "main_agreements":
            apply_main_agreements_slide(
                output_path,
                idx,
                c.title,
                c.table_headers or [],
                c.table_data or [],
            )

        if c.layout == "equipment_procurement_case":
            apply_equipment_procurement_title(
                output_path,
                idx,
                c.title,
                c.content or "",
            )
            apply_equipment_project_detail_values(
                output_path,
                idx,
                data_dir,
            )

    # 템플릿 기본 슬라이드가 남지 않도록 생성된 슬라이드를 content_list 길이에 맞춰 정리
    prune_slides(output_path, keep=len(content_list))

    # title 위치/스타일은 각 레이아웃 템플릿 기준으로 재적용한다.
    # (1페이지 표지, session_title 레이아웃은 제외)
    # apply_layout_template_title_positions(output_path, ordered_content)
    # enforce_left_alignment_for_titles(output_path, ordered_content)

    print("=" * 80)
    print(f"\n✅ 문서 생성 완료: {output_path}")
    print(f"   저장 폴더: {output_path.parent}")
    print(f"   총 슬라이드: {len(content_list)}")

    # 로고 확인
    from pptx import Presentation
    prs = Presentation(str(output_path))
    logo_count = 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                left_inch = float(shape.left) / 914400
                if left_inch > 10:  # 오른쪽 상단 로고
                    logo_count += 1

    print(f"   로고 포함: {logo_count}개 슬라이드")

    # 레이아웃 요약
    print("\n" + "=" * 80)
    print("테스트된 레이아웃 요약")
    print("=" * 80)
    print("슬라이드 1: 표지 (title)")
    print("슬라이드 2: Executive Summary (two_tables)")
    print("슬라이드 3: Approval request")
    print("슬라이드 4: 세션 표지 (session_title)")
    print("슬라이드 5: Executive Summary 상세 (executive_summary)")
    print("=" * 80)

    return output_path


if __name__ == "__main__":
    test_airun_pptx()
    def _as_dict(value: Any) -> dict[str, Any]:
        return value if isinstance(value, dict) else {}
